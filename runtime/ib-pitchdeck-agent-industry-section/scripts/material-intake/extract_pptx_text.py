#!/usr/bin/env python3
"""Extract readable text from PPT/PPTX files.

Use python-pptx when available, otherwise fallback to a minimal XML parser from
the pptx archive.
"""

from __future__ import annotations

# Runtime scripts can be run directly. Shared helpers remain in runtime
# `scripts/`; production tools live under role scripts; validators live under QC.
import sys as _ib_sys
from pathlib import Path as _IbPath
_IB_ROLE_SCRIPT_DIR = _IbPath(__file__).resolve().parent
_IB_RUNTIME_ROOT = next(
    _p for _p in _IbPath(__file__).resolve().parents
    if (_p / 'configs').is_dir() and (_p / 'scripts').is_dir()
)
_IB_SHARED_SCRIPT_DIR = _IB_RUNTIME_ROOT / "scripts"
_IB_ROLE_SCRIPT_DIRS = sorted(_p for _p in (_IB_RUNTIME_ROOT / 'scripts').iterdir() if _p.is_dir())
_IB_QC_VALIDATOR_DIRS = sorted((_IB_RUNTIME_ROOT / 'scripts' / 'qc' / 'validators').glob('*'))
_IB_IMPORT_PATHS = [str(_IB_ROLE_SCRIPT_DIR)]
for _ib_dir in [*_IB_ROLE_SCRIPT_DIRS, *_IB_QC_VALIDATOR_DIRS]:
    _ib_text = str(_ib_dir)
    if _ib_text not in _IB_IMPORT_PATHS:
        _IB_IMPORT_PATHS.append(_ib_text)
_IB_IMPORT_PATHS.append(str(_IB_SHARED_SCRIPT_DIR))
for _ib_path in list(_IB_IMPORT_PATHS):
    if _ib_path in _ib_sys.path:
        _ib_sys.path.remove(_ib_path)
for _ib_path in reversed(_IB_IMPORT_PATHS):
    _ib_sys.path.insert(0, _ib_path)

import argparse
import re
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path


TEXT_PATH_RE = re.compile(r"\[slide\s+(\d+)\]")


def _walk_shapes(shape, prefix: str, lines: list[str]) -> None:
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            _walk_shapes(child, prefix, lines)
        return

    if hasattr(shape, "text") and str(shape.text).strip():
        lines.append(f"{prefix} {str(shape.text).strip()}")
        return

    if hasattr(shape, "table"):
        table = shape.table
        for row in table.rows:
            row_values = [cell.text_frame.text.strip() if cell.text_frame else "" for cell in row.cells]
            filtered = [value for value in row_values if value]
            if filtered:
                lines.append(f"{prefix} {' | '.join(filtered)}")


def _extract_from_xml(xml_bytes: bytes) -> str:
    root = ET.fromstring(xml_bytes)
    text_nodes = [node.text.strip() for node in root.iter() if node.text and node.text.strip()]
    if text_nodes:
        return " | ".join(text_nodes)
    # Common PPTX text nodes are under a:t, so include deeper extraction.
    ns = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }
    a_t_nodes = root.findall(".//a:t", ns)
    t_vals = [node.text.strip() for node in a_t_nodes if node.text and node.text.strip()]
    return " | ".join(t_vals)


def _extract_with_python_pptx(source: Path, output_path: str | None = None) -> tuple[str, list[str], int]:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:
        return "", [f"python-pptx unavailable ({exc})"], 1

    try:
        presentation = Presentation(str(source))
    except Exception as exc:
        return "", [f"cannot open PPTX: {exc}"], 1

    lines: list[str] = []
    for slide_idx, slide in enumerate(presentation.slides, start=1):
        prefix = f"[slide {slide_idx}]"
        for shape in slide.shapes:
            _walk_shapes(shape, prefix, lines)

    text_value = "\n".join([line.strip() for line in lines if line.strip()])
    if not text_value:
        return "", ["no readable text extracted from this pptx"], 1
    if output_path:
        Path(output_path).write_text(text_value, encoding="utf-8")
    return text_value, [], 0


def _extract_with_zip_archive(source: Path, output_path: str | None = None) -> tuple[str, list[str], int]:
    lines: list[str] = []
    try:
        with zipfile.ZipFile(source, "r") as archive:
            for name in sorted(
                item
                for item in archive.namelist()
                if item.startswith("ppt/slides/slide") and item.endswith(".xml")
            ):
                m = TEXT_PATH_RE.search(name.split("/")[-1])
                slide_no = m.group(1) if m else ""
                prefix = f"[slide {slide_no}]" if slide_no else "[slide]"
                text = _extract_from_xml(archive.read(name))
                if text:
                    lines.append(f"{prefix} {text}")
    except Exception as exc:
        return "", [f"failed to read pptx package: {exc}"], 1

    text_value = "\n".join(line for line in lines if line.strip())
    if not text_value:
        return "", ["no readable text extracted from this pptx"], 1
    if output_path:
        Path(output_path).write_text(text_value, encoding="utf-8")
    return text_value, [], 0


def extract_pptx_text(pptx_path: str, output_path: str | None = None, material_id: str = "") -> tuple[str, list[str], int]:
    source = Path(pptx_path)
    if not source.exists():
        return "", [f"source path not found: {source}"], 1

    if source.suffix.lower() == ".ppt":
        return "", [".ppt binary text extraction not supported without external tools"], 1

    text_value, warnings, status = _extract_with_python_pptx(source, output_path)
    if status == 0:
        return text_value, warnings, status

    # Fallback to direct XML parse only for .pptx files.
    return _extract_with_zip_archive(source, output_path)


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--pptx-path", required=True)
    parser.add_argument("--output-text", required=True)
    parser.add_argument("--material-id", default="")
    args = parser.parse_args()
    text_value, warnings, status = extract_pptx_text(args.pptx_path, args.output_text, args.material_id)
    if status != 0:
        payload = {
            "schema_version": "pptx_extraction_v1",
            "status": "failed",
            "material_id": args.material_id,
            "source_path": args.pptx_path,
            "extracted_text_path": args.output_text,
            "text": "",
            "warnings": warnings,
        }
        print(payload)
        return 1
    payload = {
        "schema_version": "pptx_extraction_v1",
        "status": "complete",
        "material_id": args.material_id,
        "source_path": args.pptx_path,
        "extracted_text_path": args.output_text,
        "text_length": len(text_value),
        "text": text_value,
        "warnings": warnings,
    }
    print(payload)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
