#!/usr/bin/env python3
"""Run material intake end-to-end and emit manifest/extract artifacts."""

from __future__ import annotations

# Runtime scripts can be run directly. Shared helpers remain in runtime
# `scripts/`; production tools live under role scripts; validators live under QC.
import sys as _ib_sys
from pathlib import Path as _IbPath
_IB_ROLE_SCRIPT_DIR = _IbPath(__file__).resolve().parent
_IB_RUNTIME_ROOT = next(
    _p for _p in _IbPath(__file__).resolve().parents
    if (_p / 'configs').is_dir() and (_p / 'scripts').is_dir()
)
_IB_SHARED_SCRIPT_DIR = _IB_RUNTIME_ROOT / "scripts" / "_lib"
_IB_ROLE_SCRIPT_DIRS = sorted(_p for _p in (_IB_RUNTIME_ROOT / 'scripts').iterdir() if _p.is_dir())
_IB_QC_VALIDATOR_DIRS = sorted((_IB_RUNTIME_ROOT / 'scripts' / 'qc' / 'validators').glob('*'))
_IB_IMPORT_PATHS = [str(_IB_ROLE_SCRIPT_DIR)]
for _ib_dir in [*_IB_ROLE_SCRIPT_DIRS, *_IB_QC_VALIDATOR_DIRS]:
    _ib_text = str(_ib_dir)
    if _ib_text not in _IB_IMPORT_PATHS:
        _IB_IMPORT_PATHS.append(_ib_text)
_IB_IMPORT_PATHS.append(str(_IB_SHARED_SCRIPT_DIR))
for _ib_path in list(_IB_IMPORT_PATHS):
    if _ib_path in _ib_sys.path:
        _ib_sys.path.remove(_ib_path)
for _ib_path in reversed(_IB_IMPORT_PATHS):
    _ib_sys.path.insert(0, _ib_path)

import argparse
import csv
import hashlib
import html
import json
import re
import urllib.error
import zipfile
import zlib
import xml.etree.ElementTree as ET
from datetime import datetime, timezone
from html.parser import HTMLParser
from pathlib import Path
from typing import Any
from urllib.request import Request, urlopen

from material_intake_common import (
    as_bool,
    clean_text_block,
    classify_access,
    is_url,
    normalize_source_type_hint,
    file_fingerprint,
    read_text_file,
    text,
    infer_material_kind,
)
from material_intake_common import CANONICAL_SOURCE_TYPES, normalize_source_type


CJK_RE = re.compile(r"[\u4e00-\u9fff]")
PDF_TEXT_RE = re.compile(rb"\((.*?)\)")
PDF_STREAM_TEXT_RE = re.compile(rb"BT(.*?)ET", re.S)
PDF_STREAM_FILTER_RE = re.compile(rb"/Filter\s*/([A-Za-z0-9]+)")
PPT_TEXT_PATH_RE = re.compile(r"\[slide\s+(\d+)\]")


def _decode_pdf_text_bytes(raw: bytes) -> str:
    text_chunks: list[str] = []
    for segment in PDF_STREAM_TEXT_RE.finditer(raw):
        block = segment.group(0)
        data = block
        if b"/Filter" in data:
            filter_match = PDF_STREAM_FILTER_RE.search(data)
            if filter_match and filter_match.group(1).lower() == b"flatedecode":
                start = data.find(b"stream") + 6
                end = data.find(b"endstream")
                if start > 6 and end > start:
                    compressed = data[start:end].strip(b"\r\n")
                    try:
                        data = zlib.decompress(compressed)
                    except Exception:
                        pass
        for raw_text in PDF_TEXT_RE.finditer(data):
            raw_fragment = raw_text.group(1)
            try:
                text_chunks.append(raw_fragment.decode("latin1", errors="ignore"))
            except Exception:
                text_chunks.append(str(raw_fragment))
    if text_chunks:
        return clean_text_block("\n".join(text_chunks))
    return clean_text_block(raw.decode("latin1", errors="ignore"))


def extract_pdf_text(pdf_path: str, output_path: str | None = None) -> tuple[str, list[str], int]:
    warnings: list[str] = []
    source = Path(pdf_path)
    if not source.exists():
        return "", [f"source path not found: {source}"], 1

    try:
        from PyPDF2 import PdfReader  # type: ignore

        reader = PdfReader(str(source))
        text_chunks = []
        for page in reader.pages:
            try:
                value = page.extract_text() or ""
            except Exception:
                continue
            if value:
                text_chunks.append(value)
        text_value = clean_text_block("\n".join(text_chunks))
        if not text_value:
            raise RuntimeError("PyPDF2 returned empty text")
        if output_path:
            Path(output_path).write_text(text_value, encoding="utf-8")
        return text_value, warnings, 0
    except Exception:
        pass

    text_value = _decode_pdf_text_bytes(source.read_bytes())
    if not text_value.strip():
        warnings.append("PDF text extraction fallback found no readable text stream; manual review required.")
        return "", warnings, 1
    if output_path:
        Path(output_path).write_text(text_value, encoding="utf-8")
    return text_value, warnings, 0


def _walk_shapes(shape: Any, prefix: str, lines: list[str]) -> None:
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            _walk_shapes(child, prefix, lines)
        return
    if hasattr(shape, "text") and str(shape.text).strip():
        lines.append(f"{prefix} {str(shape.text).strip()}")
        return
    if hasattr(shape, "table"):
        table = shape.table
        for row in table.rows:
            row_values = [cell.text_frame.text.strip() if cell.text_frame else "" for cell in row.cells]
            filtered = [value for value in row_values if value]
            if filtered:
                lines.append(f"{prefix} {' | '.join(filtered)}")


def _extract_pptx_xml(xml_bytes: bytes) -> str:
    root = ET.fromstring(xml_bytes)
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    a_t_nodes = root.findall(".//a:t", ns)
    t_vals = [node.text.strip() for node in a_t_nodes if node.text and node.text.strip()]
    if t_vals:
        return " | ".join(t_vals)
    text_nodes = [node.text.strip() for node in root.iter() if node.text and node.text.strip()]
    return " | ".join(text_nodes)


def _extract_pptx_with_python_pptx(source: Path, output_path: str | None = None) -> tuple[str, list[str], int]:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:
        return "", [f"python-pptx unavailable ({exc})"], 1

    try:
        presentation = Presentation(str(source))
    except Exception as exc:
        return "", [f"cannot open PPTX: {exc}"], 1

    lines: list[str] = []
    for slide_idx, slide in enumerate(presentation.slides, start=1):
        prefix = f"[slide {slide_idx}]"
        for shape in slide.shapes:
            _walk_shapes(shape, prefix, lines)

    text_value = "\n".join([line.strip() for line in lines if line.strip()])
    if not text_value:
        return "", ["no readable text extracted from this pptx"], 1
    if output_path:
        Path(output_path).write_text(text_value, encoding="utf-8")
    return text_value, [], 0


def _extract_pptx_with_zip(source: Path, output_path: str | None = None) -> tuple[str, list[str], int]:
    lines: list[str] = []
    try:
        with zipfile.ZipFile(source, "r") as archive:
            for name in sorted(
                item
                for item in archive.namelist()
                if item.startswith("ppt/slides/slide") and item.endswith(".xml")
            ):
                match = PPT_TEXT_PATH_RE.search(name.split("/")[-1])
                slide_no = match.group(1) if match else ""
                prefix = f"[slide {slide_no}]" if slide_no else "[slide]"
                text_value = _extract_pptx_xml(archive.read(name))
                if text_value:
                    lines.append(f"{prefix} {text_value}")
    except Exception as exc:
        return "", [f"failed to read pptx package: {exc}"], 1

    text_value = "\n".join(line for line in lines if line.strip())
    if not text_value:
        return "", ["no readable text extracted from this pptx"], 1
    if output_path:
        Path(output_path).write_text(text_value, encoding="utf-8")
    return text_value, [], 0


def extract_pptx_text(pptx_path: str, output_path: str | None = None, material_id: str = "") -> tuple[str, list[str], int]:
    del material_id
    source = Path(pptx_path)
    if not source.exists():
        return "", [f"source path not found: {source}"], 1
    if source.suffix.lower() == ".ppt":
        return "", [".ppt binary text extraction not supported without external tools"], 1

    text_value, warnings, status = _extract_pptx_with_python_pptx(source, output_path)
    if status == 0:
        return text_value, warnings, status
    return _extract_pptx_with_zip(source, output_path)


def _read_csv_text(path: Path) -> str:
    lines: list[str] = []
    with path.open("r", encoding="utf-8", errors="replace") as handle:
        reader = csv.reader(handle)
        for row in reader:
            if row:
                line = " | ".join([cell.strip() for cell in row if str(cell).strip()])
                if line:
                    lines.append(line)
    return "\n".join(lines)


def _read_xlsx_text(path: Path) -> str:
    try:
        with zipfile.ZipFile(path, "r") as archive:
            names = archive.namelist()
            shared_strings = []
            if "xl/sharedStrings.xml" in names:
                shared_xml = ET.fromstring(archive.read("xl/sharedStrings.xml"))
                for item in shared_xml.findall(".//{http://schemas.openxmlformats.org/spreadsheetml/2006/main}t"):
                    shared_strings.append("".join(item.itertext()).strip())

            ns = {"ns": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
            sheets = [name for name in names if name.startswith("xl/worksheets/") and name.endswith(".xml")]
            rows: list[str] = []
            for sheet_name in sorted(sheets):
                sheet = ET.fromstring(archive.read(sheet_name))
                for row in sheet.findall(".//ns:sheetData/ns:row", ns):
                    values: list[str] = []
                    for cell in row.findall("ns:c", ns):
                        type_attr = cell.attrib.get("t", "")
                        if type_attr == "inlineStr":
                            inline_text = " ".join(
                                "".join(node.itertext()).strip()
                                for node in cell.findall(".//ns:is//ns:t", ns)
                                if "".join(node.itertext()).strip()
                            )
                            if inline_text:
                                values.append(inline_text.strip())
                            continue

                        value_node = cell.find("ns:v", ns)
                        if value_node is None:
                            continue
                        value = value_node.text or ""
                        if type_attr == "s":
                            try:
                                value = shared_strings[int(value)]
                            except Exception:
                                pass
                        values.append(str(value).strip())
                    line = " | ".join([item for item in values if item])
                    if line:
                        rows.append(line)
            return "\n".join(rows)
    except Exception:
        return ""


def extract_excel_text(excel_path: str, output_path: str | None = None, material_id: str = "") -> tuple[str, list[str], int]:
    del material_id
    source = Path(excel_path)
    if not source.exists():
        return "", [f"source path not found: {source}"], 1
    suffix = source.suffix.lower()
    if suffix == ".csv":
        text_value = _read_csv_text(source)
    elif suffix in {".xls", ".xlsx"}:
        if suffix == ".xls":
            return "", ["xls binary format is unsupported without external dependencies"], 1
        text_value = _read_xlsx_text(source)
    else:
        text_value = read_text_file(source)
        if not text_value.strip():
            return "", [f"unsupported file extension for structured extract: {suffix}"], 1

    if not text_value.strip():
        return "", ["no readable table rows extracted from file"], 1

    if output_path:
        Path(output_path).write_text(text_value, encoding="utf-8")
    return text_value, [], 0


class _TextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self._chunks: list[str] = []

    def handle_data(self, data: str) -> None:
        if data:
            self._chunks.append(data)

    def text(self) -> str:
        return clean_text_block(html.unescape(" ".join(self._chunks)))


def _strip_scripts(html_text: str) -> str:
    html_text = re.sub(r"<script[\s\S]*?</script>", " ", html_text, flags=re.IGNORECASE)
    html_text = re.sub(r"<style[\s\S]*?</style>", " ", html_text, flags=re.IGNORECASE)
    return html_text


def extract_web_url(url: str, output_path: str | None = None, material_id: str = "") -> tuple[str, list[str], int]:
    del material_id
    request = Request(url, headers={"User-Agent": "Mozilla/5.0 (compatible; ib-pitchdeck-agent/1.0)"})
    try:
        with urlopen(request, timeout=20) as response:
            raw = response.read()
            mime = response.headers.get_content_charset() or "utf-8"
    except urllib.error.URLError as exc:
        return "", [f"url fetch failed: {exc}"], 1
    except Exception as exc:
        return "", [f"unexpected fetch error: {exc}"], 1

    try:
        html_text = raw.decode(mime, errors="replace")
    except Exception:
        html_text = raw.decode("utf-8", errors="replace")
    if "<html" in html_text.lower():
        parser = _TextExtractor()
        parser.feed(_strip_scripts(html_text))
        parser.close()
        text_content = parser.text()
    else:
        text_content = html_text

    if not text_content.strip():
        return "", ["url content empty or non-text"], 1

    cleaned = clean_text_block(text_content)
    if output_path:
        Path(output_path).write_text(cleaned, encoding="utf-8")
    return cleaned, [], 0


def _write_json(path: Path, payload: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def _new_material_id(index: int) -> str:
    return f"MAT-{index:03d}"


def _material_access(material: dict[str, Any], source_path: str) -> str:
    return text(material.get("source_access") or classify_access(material.get("source_type", ""), source_path))


def _source_hash(source_path: str) -> str:
    if is_url(source_path):
        return hashlib.sha256(text(source_path).encode("utf-8")).hexdigest()
    return file_fingerprint(Path(source_path))


def _build_material_entry(
    material_id: str,
    source: str,
    source_type: str,
    source_access: str,
    title: str,
    brief_excerpt: str | None = None,
) -> dict[str, Any]:
    source_type = normalize_source_type(source_type)
    material = {
        "material_id": material_id,
        "source_type": normalize_source_type(source_type),
        "source_access": source_access,
        "file_path_or_url": source,
        "locator": "" if source == "inline_user_text" else source,
        "material_title": text(title),
        "material_kind": infer_material_kind(source, source_type),
        "extraction_status": "pending",
        "extraction_limitations": "not_processed",
        "can_be_used_as_evidence": False,
    }
    if brief_excerpt:
        material["brief_excerpt"] = clean_text_block(brief_excerpt)
    return material


def _default_extraction_status(source_text: str, source_file: str) -> tuple[str, list[str], bool]:
    if not source_text.strip():
        return "failed", ["no readable text extracted"], False
    if "placeholder" in source_text.lower() and len(source_text) < 200:
        return "partial", ["auto extraction returned short/placeholder-like content; confirm against source"], False
    if not source_file:
        return "partial", ["source text written to transient location only"], False
    return "complete", [], True


def _emit_entry(
    *,
    material_id: str,
    source: str,
    source_type: str,
    source_access: str,
    title: str,
    brief_excerpt: str | None = None,
    source_hash: str = "",
) -> tuple[dict[str, Any], dict[str, Any]]:
    material = _build_material_entry(
        material_id=material_id,
        source=source,
        source_type=source_type,
        source_access=source_access,
        title=title,
        brief_excerpt=brief_excerpt,
    )
    source_classification_entry = {
        "material_id": material_id,
        "source_type": material["source_type"],
        "source_access": material["source_access"],
        "file_path_or_url": source,
        "source_hash": source_hash,
        "source_date": datetime.now(timezone.utc).isoformat(),
    }
    return material, source_classification_entry


def _extract_one(material: dict[str, Any], material_texts_dir: Path) -> dict[str, Any]:
    material_id = material["material_id"]
    source = material["file_path_or_url"]
    source_type = material["source_type"]
    text_path = material_texts_dir / f"{material_id}.txt"
    text_path.parent.mkdir(parents=True, exist_ok=True)
    limitations: list[str] = []
    extracted_text = ""
    code = 1

    if source_type == "ppt_template":
        material["extracted_text_path"] = ""
        material["extraction_status"] = "template_registered"
        material["raw_text_available"] = False
        material["raw_text_extraction_status"] = "template_registered"
        material["extraction_limitations"] = (
            "PPT template registered for Template Layer analysis/rendering; "
            "not treated as project or industry evidence."
        )
        material["can_be_used_as_evidence"] = False
        material["extracted_text_preview"] = ""
        return material

    if source_type == "manual_url_ingestion":
        extracted_text, extract_limitations, code = extract_web_url(source, str(text_path), material_id=material_id)
        limitations.extend(extract_limitations)
    elif source.lower().endswith(".pdf"):
        extracted_text, extract_limitations, code = extract_pdf_text(source, str(text_path))
        limitations.extend(extract_limitations)
    elif source.lower().endswith((".ppt", ".pptx")):
        extracted_text, extract_limitations, code = extract_pptx_text(source, str(text_path))
        limitations.extend(extract_limitations)
    elif source.lower().endswith((".xls", ".xlsx", ".csv")):
        extracted_text, extract_limitations, code = extract_excel_text(source, str(text_path))
        limitations.extend(extract_limitations)
    else:
        if source == "inline_user_text":
            extracted_text = material.get("brief_excerpt") or material.get("material_title") or ""
            code = 0
            try:
                text_path.write_text(extracted_text, encoding="utf-8")
            except Exception:
                limitations.append("inline brief could not be written to extracted text path")
                code = 1
        else:
            code = 1
            limitations.append(f"unsupported source_type: {source_type}")
            limitations.append("unsupported extension for extraction")
    extraction_status, quality_limitations, can_use = _default_extraction_status(extracted_text, str(text_path))
    if quality_limitations:
        limitations.extend(quality_limitations)
    if code == 0 and extraction_status == "complete":
        text_path.write_text(extracted_text, encoding="utf-8")
        can_use = bool(text_path.exists() and extracted_text.strip())
    elif code != 0 and extracted_text:
        text_path.write_text(extracted_text, encoding="utf-8")
    material["extracted_text_path"] = str(text_path)
    material["extraction_status"] = extraction_status if code == 0 else "failed"
    material["raw_text_available"] = bool(can_use) and material["extraction_status"] == "complete"
    material["raw_text_extraction_status"] = material["extraction_status"]
    material["extraction_limitations"] = "; ".join(limitations) if limitations else "none"
    # Raw text capture is not evidence review. A role LLM must extract facts,
    # metrics, quotes, unknowns, and use limits before downstream evidence use.
    material["can_be_used_as_evidence"] = False
    material["extracted_text_preview"] = clean_text_block(extracted_text)[:320]
    return material


def ingest_materials(
    *,
    brief_text: str | None,
    files: list[str],
    template_files: list[str] | None = None,
    urls: list[str],
    default_file_source_type: str,
    default_url_source_type: str,
    output_material_manifest: Path,
    output_material_extracts: Path,
    output_source_classification: Path | None = None,
    dry_run: bool = False,
) -> tuple[dict[str, Any], dict[str, Any], dict[str, Any]]:

    materials: list[dict[str, Any]] = []
    extract_entries: list[dict[str, Any]] = []
    source_classification: list[dict[str, Any]] = []
    idx = 1

    # User brief is not considered a file/url source.
    if brief_text:
        material, classification_entry = _emit_entry(
            material_id=_new_material_id(idx),
            source="inline_user_text",
            source_type=normalize_source_type_hint("inline_user_text", "project_specific_material"),
            source_access="user_provided",
            title="User inline brief",
            brief_excerpt=brief_text,
        )
        materials.append(material)
        source_classification.append(classification_entry)
        idx += 1

    for path in files:
        source_type = normalize_source_type_hint(path, default_file_source_type)
        source_access = _material_access({"source_type": source_type}, path)
        material, classification_entry = _emit_entry(
            material_id=_new_material_id(idx),
            source=path,
            source_type=source_type,
            source_access=source_access,
            title=Path(path).name,
            source_hash=_source_hash(path),
        )
        materials.append(material)
        source_classification.append(classification_entry)
        idx += 1

    template_files = template_files or []

    for path in template_files:
        material, classification_entry = _emit_entry(
            material_id=_new_material_id(idx),
            source=path,
            source_type="ppt_template",
            source_access="user_provided",
            title=Path(path).name,
            source_hash=_source_hash(path),
        )
        material["material_kind"] = "ppt_template"
        material["extraction_status"] = "template_registered"
        material["extraction_limitations"] = "registered for template selection; not evidence"
        materials.append(material)
        source_classification.append(classification_entry)
        idx += 1

    for url in urls:
        source_type = normalize_source_type_hint(url, default_url_source_type)
        source_access = _material_access({"source_type": source_type}, url)
        material, classification_entry = _emit_entry(
            material_id=_new_material_id(idx),
            source=url,
            source_type=source_type,
            source_access=source_access,
            title=url,
            source_hash=_source_hash(url),
        )
        materials.append(material)
        source_classification.append(classification_entry)
        idx += 1

    material_texts_dir = output_material_extracts.parent / "material_texts"
    for material in materials:
        extracted = _extract_one(material, material_texts_dir)
        extract_entries.append(
            {
                "material_id": material["material_id"],
                "source_type": material["source_type"],
                "source_access": material["source_access"],
                "file_path_or_url": material["file_path_or_url"],
                "extracted_text_path": extracted["extracted_text_path"],
                "raw_text_path": extracted["extracted_text_path"],
                "raw_text_available": extracted["raw_text_available"],
                "raw_text_extraction_status": extracted["raw_text_extraction_status"],
                "content_capture_status": "captured" if extracted["raw_text_available"] else "capture_failed",
                "extraction_status": extracted["extraction_status"],
                "extraction_limitations": extracted["extraction_limitations"],
                "llm_extraction_status": (
                    "not_relevant_for_knowledge"
                    if extracted["source_type"] == "ppt_template"
                    else
                    "pending_llm_extraction"
                    if extracted["raw_text_available"]
                    else "blocked_no_readable_text"
                ),
                "can_be_used_as_evidence": False,
                "extracted_facts": [],
                "extracted_metrics": [],
                "quoted_excerpts": [],
                "unknowns_or_conflicts": [],
                "claim_use_limitations": (
                    "Raw content captured only; do not use as evidence until a role LLM extracts "
                    "source-faithful facts, metrics, quoted excerpts, unknowns/conflicts, and use limits."
                    if extracted["raw_text_available"]
                    else (
                        "Template material; use only in Template Layer for style/layout/rendering."
                        if extracted["source_type"] == "ppt_template"
                        else extracted["extraction_limitations"]
                    )
                ),
                "evidence_snapshot": extracted["extracted_text_preview"],
            }
        )

    manifest = {
        "schema_version": "material_manifest_v1",
        "created_at": datetime.now(timezone.utc).isoformat(),
        "policy_context": "pre_mandate_client_pitch",
        "materials": materials,
        "source_type_policy": {
            "project_specific_material": "user_provided user input and uploaded files; candidate until extraction and review",
            "user_curated_industry_report": "curated report candidate; subject to formal evidence extraction/review",
            "manual_url_ingestion": "public source candidate; requires source locator and archive for formal evidence",
            "ppt_template": "user-provided presentation template; use for Template Layer selection/analysis, not evidence",
        },
    }

    extracts = {
        "schema_version": "material_extracts_v1",
        "artifact_semantics": (
            "Content capture plus LLM extraction workspace. Raw text availability does not mean "
            "evidence usability. Project facts should be transcribed into input_card.json; "
            "industry facts from provided reports should be extracted by a role LLM before entering Knowledge."
        ),
        "materials_source": str(output_material_manifest),
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "extracts": extract_entries,
    }

    source_classification_payload = {
        "schema_version": "source_classification_v1",
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "materials": source_classification,
    }

    if not dry_run:
        _write_json(output_material_manifest, manifest)
        _write_json(output_material_extracts, extracts)
        if output_source_classification:
            _write_json(output_source_classification, source_classification_payload)
    return manifest, extracts, source_classification_payload


def _slug(value: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9_\-\u4e00-\u9fff]+", "_", value.strip())
    cleaned = re.sub(r"_+", "_", cleaned).strip("_")
    return cleaned or "case"


def _read_brief(args: argparse.Namespace) -> str:
    if args.brief_text:
        return str(args.brief_text).strip()
    if args.brief_file:
        return Path(args.brief_file).read_text(encoding="utf-8").strip()
    raise SystemExit("Either --brief-text or --brief-file is required.")


def _default_run_dir(case_name: str, work_root: Path) -> Path:
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    return work_root / "runs" / _slug(case_name) / f"attempt_{stamp}"


def _detect_language(brief: str) -> tuple[str, str]:
    if CJK_RE.search(brief):
        return "zh-CN", "Chinese"
    return "en", "English"


def _mark_inline_brief_transcribed(extracts: dict[str, Any]) -> dict[str, Any]:
    for item in extracts.get("extracts", []):
        if not isinstance(item, dict):
            continue
        if item.get("file_path_or_url") == "inline_user_text":
            item["llm_extraction_status"] = "project_brief_transcribed_to_input_card"
            item["can_be_used_as_evidence"] = False
            item["claim_use_limitations"] = (
                "Project-specific brief transcribed into input_card.json. "
                "Do not use as external industry evidence."
            )
            item.setdefault("unknowns_or_conflicts", [])
    return extracts


def _build_input_card(
    *,
    brief: str,
    request_language: str,
    output_language: str,
    target_company: str,
    transaction_type: str,
    industry: str,
    subsector: str,
    geography: str,
) -> dict[str, Any]:
    user_paths = ["deal_context", "target_business_summary", "source_materials"]
    normalized_paths = [
        "engagement_context.stage",
        "engagement_context.audience",
        "engagement_context.objective",
        "engagement_context.tone",
        "language",
    ]
    for field_name, path_name in [
        (target_company, "target_company"),
        (transaction_type, "transaction_type"),
        (industry, "industry"),
        (subsector, "subsector"),
        (geography, "geography"),
    ]:
        if field_name:
            user_paths.append(path_name)
    return {
        "_provenance": {
            "_description": (
                "Generated by ingest_materials.py start-brief in transcription mode. "
                "Only explicit CLI fields and the exact user brief are populated."
            ),
            "request_language": request_language,
            "user_provided_paths": user_paths,
            "normalized_metadata_paths": normalized_paths,
        },
        "target_company": target_company,
        "transaction_type": transaction_type,
        "industry": industry,
        "subsector": subsector,
        "geography": geography,
        "language": output_language,
        "deal_context": brief,
        "engagement_context": {
            "stage": "pre_mandate_transaction_pitch",
            "audience": "potential_client_or_management",
            "objective": "demonstrate sector understanding first, then transaction-relevant judgment and selective target context where supported",
            "tone": "insightful_but_not_over_selling",
        },
        "target_business_summary": brief,
        "user_provided_target_facts": [],
        "known_risks_or_limits": [],
        "source_materials": [
            {
                "source_name": "User inline brief",
                "source_type": "project_specific_material",
                "source_access": "user_provided",
                "source_access_path": "inline_user_text",
                "source_date": "",
                "geography": geography,
                "fact_type": "factual",
                "confidence": "high",
                "scope": "target-level",
                "notes": "Exact user-provided brief saved in materials/user_brief.md and artifacts/material_texts/MAT-001.txt.",
            }
        ],
        "management_hypotheses": [],
        "peer_set": [],
        "must_cover_topics": [],
        "must_avoid_topics": [],
        "research_direction": {
            "_description": "Optional user-provided controls only. Planner-generated research choices belong in artifacts/formal_search_plan.json.",
            "priority_websites": [],
            "preferred_source_domains": [],
            "preferred_source_packs": [],
            "priority_topics": [],
            "peer_set": [],
            "avoid_topics": [],
            "avoid_sources": [],
        },
    }


def start_brief(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Start a run from a plain user brief without hand-building intake files.")
    parser.add_argument("--case-name", required=True)
    parser.add_argument("--brief-text")
    parser.add_argument("--brief-file")
    parser.add_argument("--run-dir", help="Existing or desired run directory. If omitted, runs/<case>/attempt_<timestamp> is created under --work-root.")
    parser.add_argument("--work-root", default=".", help="Work root used when --run-dir is omitted.")
    parser.add_argument("--file", action="append", default=[], help="Additional user material path.")
    parser.add_argument("--url", action="append", default=[], help="Additional user-provided URL.")
    parser.add_argument("--template-file", action="append", default=[], help="User PPT/POTX template path.")
    parser.add_argument("--target-company", default="")
    parser.add_argument("--transaction-type", default="")
    parser.add_argument("--industry", default="")
    parser.add_argument("--subsector", default="")
    parser.add_argument("--geography", default="")
    args = parser.parse_args(argv)

    brief = _read_brief(args)
    work_root = Path(args.work_root).resolve()
    run_dir = Path(args.run_dir).resolve() if args.run_dir else _default_run_dir(args.case_name, work_root)
    artifacts = run_dir / "artifacts"
    materials = run_dir / "materials"
    artifacts.mkdir(parents=True, exist_ok=True)
    materials.mkdir(parents=True, exist_ok=True)
    (materials / "user_brief.md").write_text(brief + "\n", encoding="utf-8")

    manifest, extracts, source_classification = ingest_materials(
        brief_text=brief,
        files=list(args.file),
        template_files=list(args.template_file),
        urls=list(args.url),
        default_file_source_type="project_specific_material",
        default_url_source_type="manual_url_ingestion",
        output_material_manifest=artifacts / "material_manifest.json",
        output_material_extracts=artifacts / "material_extracts.json",
        output_source_classification=artifacts / "source_classification.json",
        dry_run=False,
    )
    extracts = _mark_inline_brief_transcribed(extracts)
    _write_json(artifacts / "material_extracts.json", extracts)
    _write_json(artifacts / "source_classification.json", source_classification)

    request_language, output_language = _detect_language(brief)
    input_card = _build_input_card(
        brief=brief,
        request_language=request_language,
        output_language=output_language,
        target_company=args.target_company.strip(),
        transaction_type=args.transaction_type.strip(),
        industry=args.industry.strip(),
        subsector=args.subsector.strip(),
        geography=args.geography.strip(),
    )
    _write_json(run_dir / "input_card.json", input_card)

    result = {
        "schema_version": "case_start_result_v1",
        "run_dir": str(run_dir),
        "material_manifest": str(artifacts / "material_manifest.json"),
        "material_extracts": str(artifacts / "material_extracts.json"),
        "source_classification": str(artifacts / "source_classification.json"),
        "input_card": str(run_dir / "input_card.json"),
        "template_registered": bool(args.template_file),
        "next_step": "Material role should review input_card transcription, then QC can run material/input validators.",
        "note": "This deterministic starter does not infer industry conclusions or external evidence.",
        "material_count": len(manifest.get("materials", [])),
    }
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0


def ingest_cli(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--brief-text", help="User brief text")
    parser.add_argument("--file", action="append", default=[], help="Upload file path")
    parser.add_argument("--template-file", action="append", default=[], help="User-provided PPT/POTX template path. Preferred over bundled template for Template/Output.")
    parser.add_argument("--url", action="append", default=[], help="Source URL")
    parser.add_argument("--default-file-source-type", default="project_specific_material")
    parser.add_argument("--default-url-source-type", default="manual_url_ingestion")
    parser.add_argument("--output-manifest", required=True)
    parser.add_argument("--output-extracts", required=True)
    parser.add_argument("--output-source-classification", default="")
    parser.add_argument("--dry-run", action="store_true")
    args = parser.parse_args(argv)

    output_manifest = Path(args.output_manifest)
    output_extracts = Path(args.output_extracts)
    output_source_classification = Path(args.output_source_classification) if args.output_source_classification else None

    if any(not normalize_source_type(st) in CANONICAL_SOURCE_TYPES for st in [args.default_file_source_type, args.default_url_source_type]):
        parser.error("default file/url source type must normalize to a canonical source_type")

    manifest, extracts, source_classification = ingest_materials(
        brief_text=args.brief_text,
        files=args.file,
        template_files=args.template_file,
        urls=args.url,
        default_file_source_type=args.default_file_source_type,
        default_url_source_type=args.default_url_source_type,
        output_material_manifest=output_manifest,
        output_material_extracts=output_extracts,
        output_source_classification=output_source_classification,
        dry_run=as_bool(args.dry_run),
    )

    result = {
        "is_valid": bool(manifest.get("materials")) if not args.dry_run else True,
        "material_manifest": str(output_manifest),
        "material_extracts": str(output_extracts),
        "source_classification": str(output_source_classification) if output_source_classification else "",
        "material_count": len(manifest.get("materials") or []),
        "extracted_count": len(extracts.get("extracts") or []),
        "validation": {
            "failed": [
                item for item in [
                    "material source_type invalid" if not item.get("source_type") in CANONICAL_SOURCE_TYPES else None
                    for item in manifest.get("materials", [])
                ]
                if item
            ]
        },
    }
    print(json.dumps(result, ensure_ascii=False, indent=2))
    if result["validation"]["failed"]:
        return 1
    return 0


def main(argv: list[str] | None = None) -> int:
    argv = list(_ib_sys.argv[1:] if argv is None else argv)
    if argv and argv[0] == "start-brief":
        return start_brief(argv[1:])
    return ingest_cli(argv)


if __name__ == "__main__":
    raise SystemExit(main())
