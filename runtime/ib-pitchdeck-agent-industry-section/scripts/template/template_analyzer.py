#!/usr/bin/env python3
"""Analyze a PPT template and emit a deterministic template profile."""

from __future__ import annotations

# Runtime scripts can be run directly. Shared helpers remain in runtime
# `scripts/`; production tools live under role scripts; validators live under QC.
import sys as _ib_sys
from pathlib import Path as _IbPath
_IB_ROLE_SCRIPT_DIR = _IbPath(__file__).resolve().parent
_IB_RUNTIME_ROOT = next(
    _p for _p in _IbPath(__file__).resolve().parents
    if (_p / 'configs').is_dir() and (_p / 'scripts').is_dir()
)
_IB_SHARED_SCRIPT_DIR = _IB_RUNTIME_ROOT / "scripts" / "_lib"
_IB_ROLE_SCRIPT_DIRS = sorted(_p for _p in (_IB_RUNTIME_ROOT / 'scripts').iterdir() if _p.is_dir())
_IB_QC_VALIDATOR_DIRS = sorted((_IB_RUNTIME_ROOT / 'scripts' / 'qc' / 'validators').glob('*'))
_IB_IMPORT_PATHS = [str(_IB_ROLE_SCRIPT_DIR)]
for _ib_dir in [*_IB_ROLE_SCRIPT_DIRS, *_IB_QC_VALIDATOR_DIRS]:
    _ib_text = str(_ib_dir)
    if _ib_text not in _IB_IMPORT_PATHS:
        _IB_IMPORT_PATHS.append(_ib_text)
_IB_IMPORT_PATHS.append(str(_IB_SHARED_SCRIPT_DIR))
for _ib_path in list(_IB_IMPORT_PATHS):
    if _ib_path in _ib_sys.path:
        _ib_sys.path.remove(_ib_path)
for _ib_path in reversed(_IB_IMPORT_PATHS):
    _ib_sys.path.insert(0, _ib_path)

import argparse
import json
import math
import re
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from deck_blueprint_utils import FIXED_PAGE_ROLES, active_body_fields
from json_utils import load_json_file

try:
    from pptx import Presentation
except Exception:
    Presentation = None


ROOT = _IB_RUNTIME_ROOT
TOKEN_PATTERN = re.compile(r"\{\{[^{}]+\}\}")
SPACE_UNIT = 0.3
PUNCTUATION_UNIT = 0.35
ASCII_UNIT = 0.55
FULLWIDTH_PUNCTUATION_UNIT = 0.55
DEFAULT_CHAR_UNIT = 1.0


def display_units(text: str) -> float:
    """Approximate rendered line width in CJK-character units."""
    units = 0.0
    for ch in re.sub(r"\[\[/?(?:b|hl)\]\]", "", text or ""):
        code = ord(ch)
        if ch in "\n\r":
            continue
        if ch.isspace():
            units += SPACE_UNIT
        elif ch in ",.;:!?()[]{}<>/\\|-_+=~'\"":
            units += PUNCTUATION_UNIT
        elif code < 128:
            units += ASCII_UNIT
        elif 0xFF61 <= code <= 0xFF9F:
            units += FULLWIDTH_PUNCTUATION_UNIT
        else:
            units += DEFAULT_CHAR_UNIT
    return units


def estimate_lines(text: str, max_line_units: float) -> int:
    if not text or max_line_units <= 0:
        return 0
    return sum(
        max(1, math.ceil(display_units(segment) / max_line_units))
        for segment in re.split(r"\r?\n", text)
    )


def layout_rules_for(slide_no: int, page_type: str, layout_budget: dict[str, Any] | None) -> dict[str, Any]:
    if not layout_budget:
        return {}
    slide_key = f"{slide_no}:{page_type}"
    slide_rules = layout_budget.get("slide_budgets", {}).get(slide_key)
    if isinstance(slide_rules, dict):
        return slide_rules
    page_rules = layout_budget.get("page_type_budgets", {}).get(page_type)
    return page_rules if isinstance(page_rules, dict) else {}


def _profile_path(path: Path | str) -> str:
    p = Path(path)
    try:
        return str(p.resolve().relative_to(ROOT))
    except Exception:
        return str(path)


DEFAULT_STYLE = {
    "colors": {
        "brand_primary": "#0D57AA",
        "accent_red": "#C03C28",
        "grid_gray": "#D9D9D9",
        "text_gray": "#555555",
        "panel_fill": "#FAFBFC",
        "matrix_axis": "#FFFFFF",
    },
    "typography": {
        "body": "Microsoft YaHei",
        "table_header": "Microsoft YaHei",
        "table_body": "Microsoft YaHei",
        "legend_pt": 8.0,
        "table_header_pt": 10.0,
        "table_body_pt": 10.0,
    },
}


def _load_json(path: Path | str) -> dict[str, Any]:
    data = load_json_file(Path(path))
    if not isinstance(data, dict):
        raise ValueError(f"{path} must be a JSON object")
    return data


def _layout_config_paths(path: Path | str) -> dict[str, Path]:
    config_path = Path(path)
    if not config_path.is_absolute():
        candidate = Path.cwd() / config_path
        config_path = candidate if candidate.exists() else ROOT / config_path
    config = _load_json(config_path)
    if config.get("schema_version") != "layout_config_v1":
        raise ValueError(f"{config_path} must use schema_version layout_config_v1")
    files = config.get("files")
    if not isinstance(files, dict):
        raise ValueError(f"{config_path} must define object field 'files'")
    resolved: dict[str, Path] = {}
    for key, raw in files.items():
        candidate = Path(str(raw))
        resolved[key] = candidate if candidate.is_absolute() else ROOT / candidate
    return resolved


def _ppt_mapping_by_slide(ppt_mapping: dict[str, Any]) -> dict[int, dict[str, Any]]:
    return {
        int(item["slide_no"]): item
        for item in ppt_mapping.get("slides", [])
        if isinstance(item, dict) and isinstance(item.get("slide_no"), int)
    }


def _layout_binding_by_slide(ppt_mapping: dict[str, Any]) -> dict[int, dict[str, Any]]:
    raw = ppt_mapping.get("layout_binding_by_slide", {})
    if not isinstance(raw, dict):
        return {}
    return {int(key): value for key, value in raw.items() if str(key).isdigit() and isinstance(value, dict)}


def _registry_supports(renderer: str, page_type: str, deck_contract: dict[str, Any]) -> dict[str, bool]:
    required = set(deck_contract.get("required_objects") or [])
    preferred = set(deck_contract.get("preferred_objects") or [])
    objects = required | preferred
    return {
        "chart": "chart" in objects or "chart" in renderer or "chart" in page_type,
        "table": "table" in objects or "table" in renderer or "table" in page_type,
        "matrix": "matrix" in renderer or "matrix" in page_type,
        "cards": "card" in renderer or "card" in page_type or page_type in {"moat_page", "trend_page"},
    }


def _variant_field_roles(
    slide_no: int,
    page_type: str,
    ppt_mapping_by_slide: dict[int, dict[str, Any]],
    layout_binding: dict[int, dict[str, Any]],
) -> dict[str, str]:
    binding = layout_binding.get(slide_no, {})
    if page_type in (binding.get("variants") or {}):
        return dict((binding["variants"][page_type].get("field_roles") or {}))
    if binding.get("selected_page_type") == page_type:
        return dict(binding.get("field_roles") or {})

    mapping_entry = ppt_mapping_by_slide.get(slide_no, {})
    if "tokens" in mapping_entry:
        return {
            token.get("field_name", ""): token.get("role", token.get("field_name", ""))
            for token in mapping_entry.get("tokens", [])
            if isinstance(token, dict) and token.get("field_name")
        }
    variant = (mapping_entry.get("controlled_variants") or {}).get(page_type, {})
    return {
        token.get("field_name", ""): token.get("role", token.get("field_name", ""))
        for token in variant.get("tokens", [])
        if isinstance(token, dict) and token.get("field_name")
    }


def build_registry(
    *,
    template: Path,
    slide_registry_path: Path,
    page_type_rules_path: Path,
    ppt_mapping_path: Path,
    layout_budget_path: Path,
    text_fit_rules_path: Path,
) -> dict[str, Any]:
    slide_registry = _load_json(slide_registry_path)
    page_type_rules = _load_json(page_type_rules_path)
    ppt_mapping = _load_json(ppt_mapping_path)
    layout_budget = _load_json(layout_budget_path)
    text_fit_rules = _load_json(text_fit_rules_path)

    allowed_by_slide = {
        int(item["slide_no"]): set(item.get("page_types") or [])
        for item in page_type_rules.get("slides", [])
        if isinstance(item, dict) and isinstance(item.get("slide_no"), int)
    }
    ppt_mapping_by_slide = _ppt_mapping_by_slide(ppt_mapping)
    layout_binding = _layout_binding_by_slide(ppt_mapping)
    budget_by_key = layout_budget.get("slide_budgets", {}) if isinstance(layout_budget, dict) else {}
    text_fit_by_key = text_fit_rules.get("fields", {}) if isinstance(text_fit_rules, dict) else {}

    slides = []
    for item in slide_registry.get("slides", []):
        slide_no = int(item.get("slide_no"))
        variants = []
        for page_type, variant in (item.get("variants") or {}).items():
            renderer_contract = variant.get("renderer_contract") or {}
            token_contract = variant.get("token_contract") or {}
            field_roles = _variant_field_roles(slide_no, page_type, ppt_mapping_by_slide, layout_binding)
            key = f"{slide_no}:{page_type}"
            variants.append(
                {
                    "page_type": page_type,
                    "renderer": str(variant.get("renderer") or ""),
                    "formal_allowed": page_type in allowed_by_slide.get(slide_no, set()),
                    "render_layout_key": str(variant.get("render_layout_key") or page_type),
                    "physical_slide": str(variant.get("physical_slide") or ""),
                    "supports": _registry_supports(str(variant.get("renderer") or ""), page_type, renderer_contract),
                    "required_body_fields": list(token_contract.get("required_body_fields") or []),
                    "field_roles": field_roles,
                    "capacity_notes": {
                        "layout_budget": budget_by_key.get(key, {}),
                        "text_fit_fields": {
                            field_key: value
                            for field_key, value in text_fit_by_key.items()
                            if str(field_key).startswith(key + ":")
                        },
                    },
                    "deprecation_status": "active" if page_type in allowed_by_slide.get(slide_no, set()) else "deprecated",
                }
            )
        slides.append(
            {
                "slide_no": slide_no,
                "fixed_page_role": FIXED_PAGE_ROLES.get(slide_no, ""),
                "slide_key": str(item.get("slide_key") or ""),
                "selection_mode": str(item.get("selection_mode") or ""),
                "default_variant": str(item.get("default_variant") or ""),
                "variants": variants,
            }
        )

    return {
        "schema_version": "template_registry_v1",
        "template_file": str(template),
        "slides": sorted(slides, key=lambda row: row["slide_no"]),
    }


def _template_from_selection(path: Path | str) -> Path | None:
    selection_path = Path(path)
    if not selection_path.exists():
        return None
    data = _load_json(selection_path)
    selected = str(data.get("selected_template_path") or "").strip()
    return Path(selected) if selected else None


def _as_str(value: Any) -> str:
    text = str(value or "").strip()
    return text


def _as_color(value: Any, fallback: str) -> str:
    if value is None:
        return fallback
    if isinstance(value, str):
        raw = value.strip().lstrip("#")
        if len(raw) == 6 and all(ch in "0123456789ABCDEFabcdef" for ch in raw):
            return f"#{raw.upper()}"
    if isinstance(value, int):
        if 0 <= value <= 0xFFFFFF:
            return f"#{value:06X}"
    if isinstance(value, (list, tuple)) and len(value) >= 3:
        try:
            return f"#{int(value[0]):02X}{int(value[1]):02X}{int(value[2]):02X}"
        except Exception:
            return fallback
    if isinstance(value, dict):
        for key in ("hex", "color", "rgb", "value", "hex_value"):
            nested = value.get(key)
            converted = _as_color(nested, fallback)
            if converted != fallback:
                return converted
    return fallback


def _as_font_size(value: Any, fallback: float) -> float:
    try:
        value = float(value)
    except (TypeError, ValueError):
        return fallback
    if value <= 0:
        return fallback
    return value


def _most_common(values: list[Any], fallback: str) -> str:
    if not values:
        return fallback
    counts = Counter([str(v).strip() for v in values if str(v).strip()])
    if not counts:
        return fallback
    return counts.most_common(1)[0][0]


def _emu_to_inches(value: Any) -> float:
    try:
        return round(float(value) / 914400.0, 3)
    except Exception:
        return 0.0


def _shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    try:
        return str(shape.text or "").strip()
    except Exception:
        return ""


def _placeholder_type(shape: Any) -> str:
    if not getattr(shape, "is_placeholder", False):
        return ""
    try:
        return str(shape.placeholder_format.type)
    except Exception:
        return "unknown"


def _shape_role(shape: Any, slide_height: float) -> str:
    text = _shape_text(shape).lower()
    y = _emu_to_inches(getattr(shape, "top", 0))
    if "source" in text or "来源" in text or "资料来源" in text:
        return "source_footer"
    if slide_height and y >= slide_height * 0.82 and len(text) <= 180:
        return "footer_or_source_area"
    if getattr(shape, "has_chart", False):
        return "chart_slot"
    if getattr(shape, "has_table", False):
        return "table_slot"
    if getattr(shape, "shape_type", None) is not None and "PICTURE" in str(shape.shape_type):
        return "image_slot"
    if getattr(shape, "has_text_frame", False):
        return "text_slot"
    return "shape"


def _slot_capacity(shape: Any) -> dict[str, Any]:
    width = _emu_to_inches(getattr(shape, "width", 0))
    height = _emu_to_inches(getattr(shape, "height", 0))
    # Conservative text density approximation used only for fit planning.
    return {
        "width_in": width,
        "height_in": height,
        "area_in2": round(width * height, 3),
        "estimated_text_units": max(12, int(width * height * 18)) if width and height else 0,
    }


def _collect_template_inventory(template_path: Path) -> tuple[dict[str, Any], list[str]]:
    warnings: list[str] = []
    if Presentation is None:
        return {
            "slide_count": 0,
            "slides": [],
            "analysis_fallback": True,
        }, ["python-pptx unavailable; template inventory fell back to deterministic config only"]

    try:
        prs = Presentation(str(template_path))
    except Exception as exc:
        return {
            "slide_count": 0,
            "slides": [],
            "analysis_fallback": True,
            "analysis_error": str(exc),
        }, [f"could not open PPTX template for inventory: {exc}"]

    page_width = _emu_to_inches(prs.slide_width)
    page_height = _emu_to_inches(prs.slide_height)
    slides: list[dict[str, Any]] = []
    aggregate = Counter()
    for idx, slide in enumerate(prs.slides, start=1):
        shapes: list[dict[str, Any]] = []
        roles = Counter()
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            role = _shape_role(shape, page_height)
            roles[role] += 1
            aggregate[role] += 1
            text = _shape_text(shape)
            shapes.append(
                {
                    "shape_index": shape_idx,
                    "name": str(getattr(shape, "name", "") or ""),
                    "role": role,
                    "placeholder_type": _placeholder_type(shape),
                    "has_text": bool(text),
                    "text_sample": text[:160],
                    "token_count": len(TOKEN_PATTERN.findall(text)),
                    "geometry": {
                        "x_in": _emu_to_inches(getattr(shape, "left", 0)),
                        "y_in": _emu_to_inches(getattr(shape, "top", 0)),
                        **_slot_capacity(shape),
                    },
                }
            )
        slides.append(
            {
                "slide_no": idx,
                "shape_count": len(shapes),
                "roles": dict(sorted(roles.items())),
                "slots": shapes,
                "supports": {
                    "chart": roles.get("chart_slot", 0) > 0,
                    "table": roles.get("table_slot", 0) > 0,
                    "image": roles.get("image_slot", 0) > 0,
                    "source_footer": roles.get("source_footer", 0) > 0 or roles.get("footer_or_source_area", 0) > 0,
                    "text": roles.get("text_slot", 0) > 0,
                },
                "information_density": "high" if len(shapes) >= 18 else "medium" if len(shapes) >= 9 else "low",
            }
        )

    return {
        "slide_count": len(slides),
        "page_size": {"width_in": page_width, "height_in": page_height},
        "slides": slides,
        "aggregate_roles": dict(sorted(aggregate.items())),
        "analysis_fallback": False,
    }, warnings


def _normalize_render_layouts(raw: dict[str, Any]) -> dict[str, Any]:
    if isinstance(raw, dict) and isinstance(raw.get("slides"), dict):
        return raw
    if isinstance(raw, dict):
        return {"slides": raw}
    return {"slides": {}}


def _collect_template_style(template_path: Path) -> tuple[dict[str, Any], bool]:
    if Presentation is None:
        return {
            "colors": DEFAULT_STYLE["colors"],
            "typography": DEFAULT_STYLE["typography"],
            "analysis_fallback": True,
        }, True

    fill_colors: list[str] = []
    line_colors: list[str] = []
    text_colors: list[str] = []
    fonts_body: list[str] = []
    fonts_headers: list[str] = []
    fonts_table: list[str] = []
    body_sizes: list[float] = []
    header_sizes: list[float] = []
    table_sizes: list[float] = []
    fallback = False

    try:
        prs = Presentation(str(template_path))
    except Exception as exc:
        return {
            "colors": DEFAULT_STYLE["colors"],
            "typography": DEFAULT_STYLE["typography"],
            "analysis_fallback": True,
            "analysis_error": str(exc),
        }, True

    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                fill = getattr(shape, "fill", None)
                if fill is not None and getattr(fill, "fore_color", None) is not None and getattr(fill.fore_color, "type", None) is not None:
                    fill_colors.append(_as_color(fill.fore_color.rgb, DEFAULT_STYLE["colors"]["panel_fill"]))
            except Exception:
                pass
            try:
                line = getattr(shape, "line", None)
                if line is not None and getattr(line, "color", None) is not None and getattr(line.color, "type", None) is not None:
                    line_colors.append(_as_color(line.color.rgb, DEFAULT_STYLE["colors"]["grid_gray"]))
            except Exception:
                pass

            text_frame = getattr(shape, "text_frame", None)
            if not text_frame:
                continue
            for paragraph in text_frame.paragraphs:
                is_token = bool(TOKEN_PATTERN.search(paragraph.text or ""))
                for run in paragraph.runs:
                    run_font = run.font
                    if run_font is None:
                        continue
                    if run_font.name:
                        fonts_body.append(_as_str(run_font.name))
                        if is_token:
                            fonts_table.append(_as_str(run_font.name))
                    if run_font.size:
                        pt = _as_font_size(run_font.size.pt, DEFAULT_STYLE["typography"]["table_body_pt"])
                        body_sizes.append(pt)
                        if is_token and run_font.bold:
                            header_sizes.append(pt)
                            if run_font.name:
                                fonts_headers.append(_as_str(run_font.name))
                        elif is_token:
                            table_sizes.append(pt)
                    try:
                        if run_font.color is not None and getattr(run_font.color, "type", None) is not None:
                            text_colors.append(_as_color(run_font.color.rgb, DEFAULT_STYLE["colors"]["text_gray"]))
                    except Exception:
                        pass

                if is_token and paragraph.font is not None and paragraph.font.name:
                    fonts_headers.append(_as_str(paragraph.font.name))
                if paragraph.font is not None and paragraph.font.size is not None:
                    pt = _as_font_size(paragraph.font.size.pt, DEFAULT_STYLE["typography"]["table_body_pt"])
                    body_sizes.append(pt)

                try:
                    if paragraph.font is not None and paragraph.font.color is not None and getattr(paragraph.font.color, "type", None) is not None:
                        text_colors.append(_as_color(paragraph.font.color.rgb, DEFAULT_STYLE["colors"]["text_gray"]))
                except Exception:
                    pass

                if not is_token and paragraph.font is not None and paragraph.font.name:
                    fonts_body.append(_as_str(paragraph.font.name))

    if not (fill_colors or line_colors or text_colors or fonts_body or fonts_headers):
        fallback = True

    body_font = _most_common(fonts_body, DEFAULT_STYLE["typography"]["body"])
    header_font = _most_common(fonts_headers, body_font)
    table_font = _most_common(fonts_table, header_font or body_font)

    return {
        "colors": {
            "brand_primary": _most_common(fill_colors, DEFAULT_STYLE["colors"]["brand_primary"]),
            "accent_red": _most_common(
                [c for c in text_colors + line_colors if c != _most_common(fill_colors, DEFAULT_STYLE["colors"]["brand_primary"])],
                DEFAULT_STYLE["colors"]["accent_red"],
            ),
            "grid_gray": _most_common(line_colors, DEFAULT_STYLE["colors"]["grid_gray"]),
            "text_gray": _most_common(text_colors, DEFAULT_STYLE["colors"]["text_gray"]),
            "panel_fill": DEFAULT_STYLE["colors"]["panel_fill"],
            "matrix_axis": DEFAULT_STYLE["colors"]["matrix_axis"],
        },
        "typography": {
            "body": body_font,
            "table_header": header_font,
            "table_body": table_font,
            "legend_pt": _as_font_size(sum(body_sizes) / len(body_sizes), DEFAULT_STYLE["typography"]["legend_pt"]) if body_sizes else DEFAULT_STYLE["typography"]["legend_pt"],
            "table_header_pt": _as_font_size(sum(header_sizes) / len(header_sizes), DEFAULT_STYLE["typography"]["table_header_pt"]) if header_sizes else DEFAULT_STYLE["typography"]["table_header_pt"],
            "table_body_pt": _as_font_size(
                sum(table_sizes) / len(table_sizes),
                _as_font_size(sum(body_sizes) / len(body_sizes), DEFAULT_STYLE["typography"]["table_body_pt"]) if body_sizes else DEFAULT_STYLE["typography"]["table_body_pt"],
            ) if table_sizes else _as_font_size(
                sum(body_sizes) / len(body_sizes), DEFAULT_STYLE["typography"]["table_body_pt"]
            ) if body_sizes else DEFAULT_STYLE["typography"]["table_body_pt"],
        },
        "analysis_fallback": fallback,
    }, fallback


def _extract_source_footer_fields(template_registry: dict[str, Any]) -> list[str]:
    fields: list[str] = []
    for slide in template_registry.get("slides") or []:
        for variant in slide.get("variants") or []:
            for key in (variant.get("field_roles") or {}):
                if str(key).lower() == "source_footer":
                    fields.append("source_footer")
                    break
    return sorted(set(fields))


def _build_variants(template_registry: dict[str, Any], layout_budget: dict[str, Any], source_footer_fields: list[str]) -> list[dict[str, Any]]:
    budget_by_key = layout_budget.get("slide_budgets", {})
    if not isinstance(budget_by_key, dict):
        budget_by_key = {}

    variants: list[dict[str, Any]] = []
    for slide in template_registry.get("slides") or []:
        if not isinstance(slide, dict):
            continue
        slide_no = int(slide.get("slide_no") or 0)
        for variant in slide.get("variants") or []:
            if not isinstance(variant, dict):
                continue
            page_type = str(variant.get("page_type") or "").strip()
            if not page_type:
                continue
            supports = variant.get("supports") if isinstance(variant.get("supports"), dict) else {}
            field_roles = variant.get("field_roles") if isinstance(variant.get("field_roles"), dict) else {}
            budget = budget_by_key.get(f"{slide_no}:{page_type}", {})
            body_fields = [str(item) for item in (variant.get("required_body_fields") or []) if str(item)]
            variants.append(
                {
                    "slide_no": slide_no,
                    "page_type": page_type,
                    "render_layout": str(variant.get("render_layout_key") or page_type),
                    "supports": {
                        "chart": bool(supports.get("chart")),
                        "table": bool(supports.get("table")),
                        "matrix": bool(supports.get("matrix")),
                        "cards": bool(supports.get("cards")),
                    },
                    "required_body_fields": body_fields,
                    "field_roles": {str(k): str(v) for k, v in field_roles.items() if str(k)},
                    "capacity_rules": {
                        "body_fields_max_units": dict((budget or {}).get("body_fields_max_units", {})),
                        "layout_rules": {
                            "table": dict((budget or {}).get("table", {})),
                            "matrix": dict((budget or {}).get("matrix", {})),
                        },
                    },
                    "source_footer_required": bool("source_footer" in field_roles or "source_footer" in source_footer_fields),
                }
            )
    return variants


def _merge_supports(current: dict[str, Any], incoming: dict[str, Any]) -> dict[str, Any]:
    keys = sorted(set(current) | set(incoming) | {"chart", "table", "matrix", "cards", "image", "source_footer", "text"})
    return {key: bool(current.get(key)) or bool(incoming.get(key)) for key in keys}


def _build_page_type_capability(variants: list[dict[str, Any]], inventory: dict[str, Any]) -> dict[str, Any]:
    capabilities: dict[str, dict[str, Any]] = {}
    inventory_by_slide: dict[int, dict[str, Any]] = {}
    for item in inventory.get("slides", []):
        if not isinstance(item, dict):
            continue
        try:
            inventory_by_slide[int(item.get("slide_no"))] = item
        except Exception:
            continue

    for variant in variants:
        page_type = _as_str(variant.get("page_type"))
        if not page_type:
            continue
        slide_no = int(variant.get("slide_no") or 0)
        inv = inventory_by_slide.get(slide_no, {})
        inv_supports = inv.get("supports") if isinstance(inv.get("supports"), dict) else {}
        supports = _merge_supports(variant.get("supports") if isinstance(variant.get("supports"), dict) else {}, inv_supports)
        row = capabilities.setdefault(
            page_type,
            {
                "page_type": page_type,
                "slides": [],
                "supports": {},
                "required_body_fields": [],
                "render_layouts": [],
                "density_classes": [],
            },
        )
        row["slides"].append(slide_no)
        row["supports"] = _merge_supports(row.get("supports", {}), supports)
        for field in variant.get("required_body_fields", []):
            if field not in row["required_body_fields"]:
                row["required_body_fields"].append(field)
        render_layout = _as_str(variant.get("render_layout"))
        if render_layout and render_layout not in row["render_layouts"]:
            row["render_layouts"].append(render_layout)
        density = _as_str(inv.get("information_density")) or "unknown"
        if density not in row["density_classes"]:
            row["density_classes"].append(density)

    return {
        "schema_version": "page_type_capability_v1",
        "page_types": dict(sorted(capabilities.items())),
    }


def _build_source_area(inventory: dict[str, Any], source_policy: dict[str, Any]) -> dict[str, Any]:
    source_slots: list[dict[str, Any]] = []
    for slide in inventory.get("slides", []):
        if not isinstance(slide, dict):
            continue
        for slot in slide.get("slots", []):
            if not isinstance(slot, dict):
                continue
            role = _as_str(slot.get("role"))
            if role not in {"source_footer", "footer_or_source_area"}:
                continue
            source_slots.append(
                {
                    "slide_no": slide.get("slide_no"),
                    "shape_index": slot.get("shape_index"),
                    "role": role,
                    "geometry": slot.get("geometry", {}),
                    "text_sample": slot.get("text_sample", ""),
                }
            )
    return {
        "schema_version": "source_area_v1",
        "required_source_footer": bool(source_policy.get("required_source_footer")),
        "source_footer_fields": source_policy.get("source_footer_fields", ["source_footer"]),
        "source_slots": source_slots,
        "has_detected_source_area": bool(source_slots),
    }


def _build_density_budget(layout_budget: dict[str, Any], variants: list[dict[str, Any]]) -> dict[str, Any]:
    slide_budgets = layout_budget.get("slide_budgets") if isinstance(layout_budget.get("slide_budgets"), dict) else {}
    rows: list[dict[str, Any]] = []
    for variant in variants:
        slide_no = int(variant.get("slide_no") or 0)
        page_type = _as_str(variant.get("page_type"))
        key = f"{slide_no}:{page_type}"
        budget = slide_budgets.get(key, {}) if isinstance(slide_budgets.get(key), dict) else {}
        body_fields = budget.get("body_fields_max_units") if isinstance(budget.get("body_fields_max_units"), dict) else {}
        rows.append(
            {
                "slide_no": slide_no,
                "page_type": page_type,
                "body_fields_max_units": body_fields,
                "table": budget.get("table", {}),
                "matrix": budget.get("matrix", {}),
                "max_body_units_total": sum(float(value or 0) for value in body_fields.values()) if body_fields else 0,
            }
        )
    return {
        "schema_version": "density_budget_v1",
        "global": layout_budget.get("global", {}),
        "slides": rows,
    }


def _load_paths(args) -> tuple[dict[str, Path], Path]:
    layout_paths = _layout_config_paths(args.layout_config)
    if args.slide_registry is not None:
        layout_paths["slide_registry"] = Path(args.slide_registry)
    if args.page_type_rules is not None:
        layout_paths["page_type_rules"] = Path(args.page_type_rules)
    if args.ppt_mapping is not None:
        layout_paths["ppt_mapping"] = Path(args.ppt_mapping)
    layout_paths["render_layouts"] = Path(args.render_layouts) if args.render_layouts else layout_paths["render_layouts"]
    layout_paths["text_fit_rules"] = Path(args.text_fit_rules) if args.text_fit_rules else layout_paths["text_fit_rules"]
    layout_paths["layout_budget"] = Path(args.layout_budget) if args.layout_budget else layout_paths["layout_budget"]
    template_path = Path(args.template)
    return layout_paths, template_path


def _build_profile(layout_paths: dict[str, Path], template_path: Path, output_path: Path, skip_pptextract: bool) -> dict[str, Any]:
    template_registry = build_registry(
        template=template_path,
        slide_registry_path=layout_paths["slide_registry"],
        page_type_rules_path=layout_paths["page_type_rules"],
        ppt_mapping_path=layout_paths["ppt_mapping"],
        layout_budget_path=layout_paths["layout_budget"],
        text_fit_rules_path=layout_paths["text_fit_rules"],
    )
    style_payload, style_fallback = _collect_template_style(template_path) if not skip_pptextract else (
        {"colors": DEFAULT_STYLE["colors"], "typography": DEFAULT_STYLE["typography"], "analysis_fallback": True},
        True,
    )

    if style_payload.get("analysis_fallback"):
        style_fallback = True
    inventory, inventory_warnings = _collect_template_inventory(template_path) if not skip_pptextract else (
        {"slide_count": 0, "slides": [], "analysis_fallback": True},
        ["template PPTX inventory skipped by --skip-pptextract"],
    )

    layout_budget = _load_json(layout_paths["layout_budget"])
    text_fit_rules = _load_json(layout_paths["text_fit_rules"])

    source_footer_fields = _extract_source_footer_fields(template_registry)
    render_layouts = _normalize_render_layouts(_load_json(layout_paths["render_layouts"]))
    variant_payload = _build_variants(template_registry, layout_budget, source_footer_fields)
    source_policy = {
        "source_footer_fields": source_footer_fields or ["source_footer"],
        "required_source_footer": True if source_footer_fields else bool(
            any((item.get("supports") or {}).get("source_footer") for item in inventory.get("slides", []) if isinstance(item, dict))
        ),
    }

    return {
        "schema_version": "template_profile_v1",
        "template_file": _profile_path(template_path),
        "analysis_source": "template_analyzer.py",
        "analysis_timestamp": datetime.now(timezone.utc).isoformat(),
        "analysis_fallback": bool(style_fallback),
        "analysis_errors": inventory_warnings,
        "render_layouts_source": _profile_path(layout_paths["render_layouts"]),
        "text_fit_rules_source": _profile_path(layout_paths["text_fit_rules"]),
        "layout_budget_source": _profile_path(layout_paths["layout_budget"]),
        "ppt_mapping_source": _profile_path(layout_paths["ppt_mapping"]),
        "visual_style": {
            "colors": style_payload["colors"],
            "typography": style_payload["typography"],
        },
        "text_geometry": {
            "analysis_fallback": bool(style_fallback),
            "font_samples": {
                "body": style_payload["typography"]["body"],
                "table_header": style_payload["typography"]["table_header"],
                "table_body": style_payload["typography"]["table_body"],
            },
        },
        "template_inventory": inventory,
        "render_layouts": render_layouts,
        "page_type_capability": _build_page_type_capability(variant_payload, inventory),
        "source_area": _build_source_area(inventory, source_policy),
        "density_budget": _build_density_budget(layout_budget, variant_payload),
        "dynamic_slots": {
            "slide_count": inventory.get("slide_count", 0),
            "slides": [
                {
                    "slide_no": item.get("slide_no"),
                    "supports": item.get("supports", {}),
                    "information_density": item.get("information_density", "unknown"),
                    "slot_count": item.get("shape_count", 0),
                }
                for item in inventory.get("slides", [])
                if isinstance(item, dict)
            ],
        },
        "layout": {
            "render_layouts": render_layouts,
            "text_fit_rules": {
                "fields": text_fit_rules.get("fields", {}),
                "renderer_field_aliases": text_fit_rules.get("renderer_field_aliases", {}),
            },
            "layout_budget": layout_budget,
        },
        "slide_variants": variant_payload,
        "source_policy": source_policy,
        "output_path": str(output_path),
        "template_registry_file": _profile_path(layout_paths["slide_registry"]),
    }


def _validate_profile(profile: dict[str, Any]) -> list[str]:
    errors = []
    if profile.get("schema_version") != "template_profile_v1":
        errors.append("schema_version must be template_profile_v1")
    if "visual_style" not in profile:
        errors.append("visual_style is required")
    if "layout" not in profile:
        errors.append("layout is required")
    for key in ("render_layouts", "page_type_capability", "source_area", "density_budget"):
        if not isinstance(profile.get(key), dict):
            errors.append(f"{key} is required")
    return errors


def _fit_as_list(value: Any) -> list[Any]:
    return value if isinstance(value, list) else []


def _fit_slide_variants(profile: dict[str, Any]) -> dict[tuple[int, str], dict[str, Any]]:
    variants = {}
    for item in profile.get("slide_variants", []):
        if not isinstance(item, dict):
            continue
        try:
            slide_no = int(item.get("slide_no"))
            page_type = str(item.get("page_type") or "")
        except Exception:
            continue
        variants[(slide_no, page_type)] = item
    return variants


def _fit_is_blank(value: Any) -> bool:
    return not bool(str(value or "").strip())


def _fit_has_payload(value: Any) -> bool:
    if value is None:
        return False
    if isinstance(value, dict):
        return bool(value)
    if isinstance(value, list):
        return bool(value)
    return bool(str(value).strip())


def _fit_render_layouts_by_slide(profile: dict[str, Any]) -> dict[str, dict[str, dict[str, Any]]]:
    raw = profile.get("layout", {}).get("render_layouts", {})
    if not isinstance(raw, dict):
        return {}
    slides = raw.get("slides", raw)
    if not isinstance(slides, dict):
        return {}
    return slides


def _fit_text_rules(profile: dict[str, Any]) -> tuple[dict[str, Any], dict[str, str]]:
    text_fit = profile.get("layout", {}).get("text_fit_rules", {})
    fields = text_fit.get("fields", {})
    aliases = text_fit.get("renderer_field_aliases", {})
    if not isinstance(fields, dict):
        fields = {}
    if not isinstance(aliases, dict):
        aliases = {}
    return fields, aliases


def _fit_capacity_conflict(slide_no: int, field_path: str, message: str, recommendation: str) -> dict[str, Any]:
    return {
        "conflict_type": "template_capacity_conflict",
        "slide_no": slide_no,
        "field_path": field_path,
        "message": message,
        "repair_owner": "generation",
        "repair_action": recommendation,
        "downstream_blocked": True,
    }


def _fit_check_source_footer(slide: dict[str, Any], variant: dict[str, Any], warnings: list[str], blocking: list[str]) -> None:
    if not bool(variant.get("source_footer_required")):
        return
    if _fit_is_blank(slide.get("source_note")):
        blocking.append(f"slide {slide.get('slide_no')}: source footer required by template variant but source_note is empty")
        return
    note = str(slide.get("source_note") or "")
    if note.strip().startswith("来源：") and len(note.strip()) <= 3:
        warnings.append(f"slide {slide.get('slide_no')}: source footer is present but appears empty")


def _fit_check_text_capacity(slide: dict[str, Any], layout_budget: dict[str, Any], warnings: list[str], blocking: list[str]) -> None:
    slide_no = int(slide.get("slide_no") or 0)
    page_type = str(slide.get("selected_page_type") or "")
    rules = layout_rules_for(slide_no, page_type, layout_budget)
    if not isinstance(rules, dict):
        return
    body_copy = slide.get("body_copy") or {}
    if not isinstance(body_copy, dict):
        return
    global_rules = layout_budget.get("global", {})
    global_body = global_rules.get("body_copy", {})
    default_limit = float(global_body.get("max_bullet_units_default", 88))
    field_limits = rules.get("body_fields_max_units", {})
    table_limit = float(rules.get("table", {}).get("max_cell_units", global_rules.get("table", {}).get("max_cell_units", 22)))
    for field_name, value in body_copy.items():
        if not isinstance(value, str):
            continue
        field_limit = float(field_limits.get(field_name, default_limit))
        actual = display_units(value)
        if actual > field_limit:
            blocking.append(
                f"slide {slide_no}: '{field_name}' is {actual:.1f} layout units, "
                f"template budget is {field_limit:.1f}; edit Generation output or reduce copy"
            )
        if field_name.startswith("table_") and "|" in value:
            for idx, cell in enumerate(value.split("|"), start=1):
                cell_units = display_units(cell)
                if cell_units > table_limit:
                    warnings.append(f"slide {slide_no}: table field '{field_name}' cell {idx} is {cell_units:.1f} units; >{table_limit:.1f} may wrap")


def _fit_check_text_lines(slide: dict[str, Any], text_fit_fields: dict[str, Any], aliases: dict[str, str], warnings: list[str], blocking: list[str]) -> None:
    slide_no = int(slide.get("slide_no") or 0)
    page_type = str(slide.get("selected_page_type") or "")
    for field, value in {"headline": slide.get("headline", ""), "main_message": slide.get("main_message", "")}.items():
        if not isinstance(value, str) or not value.strip():
            continue
        alias = aliases.get(field, field)
        rule = text_fit_fields.get(f"{slide_no}:{page_type}:{alias}")
        if not isinstance(rule, dict):
            continue
        max_line_units = float(rule.get("max_line_units") or 0)
        if max_line_units <= 0:
            continue
        target_lines = int(rule.get("target_lines") or 0)
        max_lines = int(rule.get("max_lines") or 0)
        actual_lines = estimate_lines(value, max_line_units)
        placeholder = rule.get("placeholder", "")
        if target_lines and actual_lines > target_lines:
            warnings.append(
                f"slide {slide_no}: '{field}' exceeds target density for {placeholder}; "
                f"estimated {actual_lines} line(s), target is {target_lines}"
            )
        if max_lines and actual_lines > max_lines:
            blocking.append(
                f"slide {slide_no}: '{field}' exceeds template max lines for {placeholder}; "
                f"estimated {actual_lines} line(s), max is {max_lines}"
            )


def _fit_check_payload_support(slide: dict[str, Any], variant: dict[str, Any], warnings: list[str], blocking: list[str]) -> None:
    supports = variant.get("supports") if isinstance(variant.get("supports"), dict) else {}
    slide_no = slide.get("slide_no")
    page_type = str(slide.get("selected_page_type") or "")
    if _fit_has_payload(slide.get("chart_data")) and not bool(supports.get("chart", False)):
        blocking.append(f"slide {slide_no}: chart data exists for '{page_type}' but template variant reports no chart support")
    if _fit_has_payload(slide.get("compare_table_data")) and not bool(supports.get("table", False)):
        blocking.append(f"slide {slide_no}: compare_table_data exists for '{page_type}' but template variant reports no table support")
    if _fit_has_payload(slide.get("matrix_data")) and not bool(supports.get("matrix", False)):
        warnings.append(f"slide {slide_no}: matrix payload exists for '{page_type}' but matrix support flag is false")
    required_fields = variant.get("required_body_fields") if isinstance(variant.get("required_body_fields"), list) else []
    required_fields = active_body_fields(required_fields, page_type, slide)
    body_copy = slide.get("body_copy") or {}
    if isinstance(body_copy, dict):
        missing = [name for name in required_fields if _fit_is_blank(body_copy.get(name))]
        if missing:
            blocking.append(f"slide {slide_no}: required body fields missing: {', '.join(missing)}")


def _fit_check_render_layout_presence(slide: dict[str, Any], layout_slides: dict[str, dict[str, dict[str, Any]]], blocking: list[str]) -> None:
    slide_no = int(slide.get("slide_no") or 0)
    page_type = str(slide.get("selected_page_type") or "")
    variant_missing = page_type not in layout_slides.get(str(slide_no), {})
    needs_render_layout = (
        _fit_has_payload(slide.get("chart_data"))
        or _fit_has_payload(slide.get("compare_table_data"))
        or _fit_has_payload(slide.get("matrix_data"))
        or _fit_has_payload(slide.get("exhibit"))
    )
    if variant_missing and needs_render_layout:
        blocking.append(f"slide {slide_no}: render layout for '{page_type}' not found in template profile; PPT rendering will produce broken output")
    elif variant_missing:
        blocking.append(f"slide {slide_no}: render layout for '{page_type}' not found in template profile; formal delivery requires explicit layout support")


def _fit_check_visual_style(profile: dict[str, Any], warnings: list[str]) -> None:
    colors = profile.get("visual_style", {}).get("colors", {})
    typography = profile.get("visual_style", {}).get("typography", {})
    if not isinstance(colors, dict) or not isinstance(typography, dict):
        warnings.append("template_profile missing visual_style; fit checks cannot validate color/font constraints")
        return
    for key, fallback in {"brand_primary": "#0D57AA", "accent_red": "#C03C28", "grid_gray": "#D9D9D9", "text_gray": "#555555"}.items():
        value = str(colors.get(key) or "")
        if value == fallback or not value:
            warnings.append(f"template_profile uses fallback color for {key}: {value}")
    for key in ("body", "table_header", "table_body"):
        if not str(typography.get(key) or "").strip():
            warnings.append(f"template_profile missing typography field {key}")


def _run_fit(renderer_spec: dict[str, Any], profile: dict[str, Any]) -> tuple[bool, list[str], list[str]]:
    warnings: list[str] = []
    blocking: list[str] = []
    slides = renderer_spec.get("slides")
    if not isinstance(slides, list):
        raise ValueError("renderer_spec.slides must be a list")
    layout_budget = profile.get("layout", {}).get("layout_budget", {})
    if not isinstance(layout_budget, dict):
        layout_budget = {}
    text_fit_fields, aliases = _fit_text_rules(profile)
    variants = _fit_slide_variants(profile)
    layout_slides = _fit_render_layouts_by_slide(profile)
    _fit_check_visual_style(profile, warnings)
    for slide in slides:
        if not isinstance(slide, dict):
            continue
        slide_no = int(slide.get("slide_no") or 0)
        page_type = str(slide.get("selected_page_type") or "")
        variant = variants.get((slide_no, page_type))
        if not isinstance(variant, dict):
            blocking.append(f"slide {slide_no}: template profile missing variant '{page_type}'")
            continue
        _fit_check_source_footer(slide, variant, warnings, blocking)
        _fit_check_text_capacity(slide, layout_budget, warnings, blocking)
        _fit_check_text_lines(slide, text_fit_fields, aliases, warnings, blocking)
        _fit_check_payload_support(slide, variant, warnings, blocking)
        _fit_check_render_layout_presence(slide, layout_slides, blocking)
    return not blocking, warnings, blocking


def _fit_inventory_by_slide(profile: dict[str, Any]) -> dict[int, dict[str, Any]]:
    inventory = profile.get("template_inventory") if isinstance(profile.get("template_inventory"), dict) else {}
    return {
        int(item.get("slide_no")): item
        for item in inventory.get("slides", [])
        if isinstance(item, dict) and str(item.get("slide_no", "")).isdigit()
    }


def _fit_slot_assignments(slide: dict[str, Any], variant: dict[str, Any], inventory_slide: dict[str, Any]) -> list[dict[str, Any]]:
    assignments: list[dict[str, Any]] = []
    slide_no = int(slide.get("slide_no") or 0)
    page_type = str(slide.get("selected_page_type") or "")
    body_copy = slide.get("body_copy") if isinstance(slide.get("body_copy"), dict) else {}
    field_roles = variant.get("field_roles") if isinstance(variant.get("field_roles"), dict) else {}
    for field_name in sorted(body_copy):
        assignments.append({"slide_no": slide_no, "page_type": page_type, "content_field": f"body_copy.{field_name}", "template_role": field_roles.get(field_name, field_name), "slot_type": "text", "placement": "body"})
    if _fit_has_payload(slide.get("chart_data")):
        assignments.append({"slide_no": slide_no, "page_type": page_type, "content_field": "chart_data", "template_role": "chart", "slot_type": "chart", "placement": "visual"})
    if _fit_has_payload(slide.get("compare_table_data")):
        assignments.append({"slide_no": slide_no, "page_type": page_type, "content_field": "compare_table_data", "template_role": "table", "slot_type": "table", "placement": "visual"})
    if _fit_has_payload(slide.get("source_note")):
        assignments.append({"slide_no": slide_no, "page_type": page_type, "content_field": "source_note", "template_role": "source_footer", "slot_type": "source_footer", "placement": "footer"})
    if inventory_slide:
        for assignment in assignments:
            assignment["template_slide_inventory"] = {"slot_count": inventory_slide.get("shape_count", 0), "information_density": inventory_slide.get("information_density", "unknown")}
    return assignments


def _build_fit_plan(renderer_spec: dict[str, Any], profile: dict[str, Any]) -> dict[str, Any]:
    layout_budget = profile.get("layout", {}).get("layout_budget", {})
    if not isinstance(layout_budget, dict):
        layout_budget = {}
    variants = _fit_slide_variants(profile)
    inventory = _fit_inventory_by_slide(profile)
    fields, aliases = _fit_text_rules(profile)
    assignments: list[dict[str, Any]] = []
    recommendations: list[dict[str, Any]] = []
    conflicts: list[dict[str, Any]] = []
    for slide in _fit_as_list(renderer_spec.get("slides")):
        if not isinstance(slide, dict):
            continue
        slide_no = int(slide.get("slide_no") or 0)
        page_type = str(slide.get("selected_page_type") or "")
        variant = variants.get((slide_no, page_type), {})
        if not variant:
            conflicts.append(_fit_capacity_conflict(slide_no, "selected_page_type", f"template profile has no variant for page type '{page_type}'", "Choose a registered page type or regenerate template_profile from the selected template."))
            continue
        assignments.extend(_fit_slot_assignments(slide, variant, inventory.get(slide_no, {})))
        body_copy = slide.get("body_copy") if isinstance(slide.get("body_copy"), dict) else {}
        rules = layout_rules_for(slide_no, page_type, layout_budget)
        field_limits = rules.get("body_fields_max_units", {}) if isinstance(rules, dict) else {}
        default_limit = float(layout_budget.get("global", {}).get("body_copy", {}).get("max_bullet_units_default", 88))
        for field_name, value in body_copy.items():
            if not isinstance(value, str):
                continue
            limit = float(field_limits.get(field_name, default_limit))
            actual = display_units(value)
            if actual > limit:
                over_by = round(actual - limit, 1)
                conflicts.append(_fit_capacity_conflict(slide_no, f"body_copy.{field_name}", f"body_copy.{field_name} exceeds template capacity by {over_by:.1f} layout units", "Return to Generation and compress/restructure this field; do not silently truncate."))
                recommendations.append({"slide_no": slide_no, "field_path": f"body_copy.{field_name}", "recommendation_type": "copy_compression", "current_units": actual, "max_units": limit, "message": f"Compress or split copy before rendering; over budget by {over_by:.1f} units."})
        supports = variant.get("supports") if isinstance(variant.get("supports"), dict) else {}
        if _fit_has_payload(slide.get("chart_data")) and not supports.get("chart"):
            conflicts.append(_fit_capacity_conflict(slide_no, "chart_data", "renderer_spec contains chart_data but selected template variant has no chart slot", "Choose a chart-capable page type or revise visual plan in Generation."))
        if _fit_has_payload(slide.get("compare_table_data")) and not supports.get("table"):
            conflicts.append(_fit_capacity_conflict(slide_no, "compare_table_data", "renderer_spec contains compare_table_data but selected template variant has no table slot", "Choose a table-capable page type or revise visual plan in Generation."))
        for field, value in {"headline": slide.get("headline", ""), "main_message": slide.get("main_message", "")}.items():
            alias = aliases.get(field, field)
            rule = fields.get(f"{slide_no}:{page_type}:{alias}")
            if not isinstance(rule, dict) or not isinstance(value, str) or not value.strip():
                continue
            max_line_units = float(rule.get("max_line_units") or 0)
            max_lines = int(rule.get("max_lines") or 0)
            if max_line_units and max_lines and estimate_lines(value, max_line_units) > max_lines:
                actual_lines = estimate_lines(value, max_line_units)
                conflicts.append(_fit_capacity_conflict(slide_no, field, f"{field} estimates to {actual_lines} lines but template allows {max_lines}", "Return to Generation and compress the headline/message before rendering."))
    return {
        "schema_version": "template_fit_plan_v1",
        "analysis_source": "template_analyzer.py fit",
        "template_profile": str(profile.get("output_path") or profile.get("template_file") or ""),
        "page_assignments": assignments,
        "copy_compression_recommendations": recommendations,
        "capacity_conflicts": conflicts,
        "template_capacity_conflict": bool(conflicts),
        "fit_decision": "template_capacity_conflict" if conflicts else "template_ready",
    }


def fit_cli(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Run Template Fit checks for renderer_spec against template profile.")
    parser.add_argument("--renderer-spec", required=True)
    parser.add_argument("--template-profile", default=str(ROOT / "configs" / "template_profile.json"))
    parser.add_argument("--output", default=str(ROOT / "artifacts" / "template_fit_validation.json"))
    parser.add_argument("--fit-plan-output", help="Optional path for artifacts/template_fit_plan.json")
    parser.add_argument("--strict", action="store_true", help="Treat warnings as blocking issues.")
    args = parser.parse_args(argv)
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        renderer_spec = _load_json(Path(args.renderer_spec))
        profile = _load_json(Path(args.template_profile))
        _is_valid, warnings, blocking = _run_fit(renderer_spec, profile)
        fit_plan = _build_fit_plan(renderer_spec, profile)
        for conflict in fit_plan.get("capacity_conflicts") or []:
            message = str(conflict.get("message") or "")
            if message and message not in blocking:
                blocking.append(message)
        if args.strict and warnings:
            blocking.extend(warnings)
            warnings = []
        is_valid = not blocking
        result = {
            "schema_version": "template_fit_v1",
            "is_valid": is_valid,
            "analysis_source": "template_analyzer.py fit",
            "renderer_spec": str(Path(args.renderer_spec)),
            "template_profile": str(Path(args.template_profile)),
            "error_count": len(blocking),
            "warning_count": len(warnings),
            "errors": blocking,
            "warnings": warnings,
            "blocking_issues": blocking,
            "capacity_conflicts": fit_plan.get("capacity_conflicts", []),
            "template_capacity_conflict": bool(fit_plan.get("template_capacity_conflict")),
            "template_fit_plan": args.fit_plan_output or "",
            "fit_checks": {"checked_slides": len(renderer_spec.get("slides") or []), "strict_mode": args.strict},
        }
    except Exception as exc:
        result = {
            "schema_version": "template_fit_v1",
            "is_valid": False,
            "analysis_source": "template_analyzer.py fit",
            "error_count": 1,
            "warning_count": 0,
            "errors": [f"{type(exc).__name__}: {exc}"],
            "warnings": ["template_profile: configs/template_profile.json"],
            "blocking_issues": [f"{type(exc).__name__}: {exc}"],
            "capacity_conflicts": [],
            "template_capacity_conflict": False,
        }
        fit_plan = {"schema_version": "template_fit_plan_v1", "analysis_source": "template_analyzer.py fit", "fit_decision": "template_fit_error", "page_assignments": [], "copy_compression_recommendations": [], "capacity_conflicts": [], "errors": result["errors"]}
    output_path.write_text(json.dumps(result, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    if args.fit_plan_output:
        fit_plan_path = Path(args.fit_plan_output)
        fit_plan_path.parent.mkdir(parents=True, exist_ok=True)
        fit_plan_path.write_text(json.dumps(fit_plan, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0 if result["is_valid"] else 1


def registry_cli(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Extract the formal template capability registry from deterministic config files.")
    parser.add_argument("--template", default=str(ROOT / "assets/industry_section_template_master.pptx"))
    parser.add_argument("--slide-registry", default=str(ROOT / "configs/slide_registry.json"))
    parser.add_argument("--page-type-rules", default=str(ROOT / "configs/page_type_rules.json"))
    parser.add_argument("--ppt-mapping", default=str(ROOT / "configs/ppt_mapping.json"))
    parser.add_argument("--layout-budget", default=str(ROOT / "configs/layout_budget.json"))
    parser.add_argument("--text-fit-rules", default=str(ROOT / "configs/text_fit_rules.json"))
    parser.add_argument("--output", required=True)
    args = parser.parse_args(argv)

    registry = build_registry(
        template=Path(args.template),
        slide_registry_path=Path(args.slide_registry),
        page_type_rules_path=Path(args.page_type_rules),
        ppt_mapping_path=Path(args.ppt_mapping),
        layout_budget_path=Path(args.layout_budget),
        text_fit_rules_path=Path(args.text_fit_rules),
    )
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(registry, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(json.dumps({"is_valid": True, "output": str(output_path)}, ensure_ascii=False, indent=2))
    return 0


def profile_cli(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--template", default="")
    parser.add_argument("--template-selection", default="", help="Path to artifacts/template_selection.json. Used when --template is omitted.")
    parser.add_argument("--layout-config", default=str(ROOT / "configs" / "layout_config.json"))
    parser.add_argument("--slide-registry", default=None)
    parser.add_argument("--page-type-rules", default=None)
    parser.add_argument("--ppt-mapping", default=None)
    parser.add_argument("--render-layouts", default=None)
    parser.add_argument("--text-fit-rules", default=None)
    parser.add_argument("--layout-budget", default=None)
    parser.add_argument("--output", default=str(ROOT / "configs" / "template_profile.json"))
    parser.add_argument("--skip-pptextract", action="store_true", help="Skip template parsing and rely on fallback values.")
    args = parser.parse_args(argv)

    if not args.template and args.template_selection:
        selected_template = _template_from_selection(args.template_selection)
        if selected_template is not None:
            args.template = str(selected_template)
    if not args.template:
        args.template = str(ROOT / "assets" / "industry_section_template_master.pptx")

    layout_paths, template_path = _load_paths(args)
    output = Path(args.output)
    output.parent.mkdir(parents=True, exist_ok=True)

    try:
        profile = _build_profile(layout_paths, template_path, output, args.skip_pptextract)
        profile["analysis_source"] = "template_analyzer.py" if not args.skip_pptextract else "template_analyzer.py (fallback)"
    except Exception as exc:
        inventory, inventory_warnings = _collect_template_inventory(template_path) if not args.skip_pptextract else (
            {"slide_count": 0, "slides": [], "analysis_fallback": True},
            ["template PPTX inventory skipped by --skip-pptextract"],
        )
        fallback_render_layouts = _normalize_render_layouts(_load_json(layout_paths["render_layouts"]))
        fallback_layout_budget = _load_json(layout_paths["layout_budget"])
        fallback_source_policy = {
            "source_footer_fields": ["source_footer"],
            "required_source_footer": False,
        }
        profile = {
            "schema_version": "template_profile_v1",
            "template_file": _profile_path(template_path),
            "analysis_source": "template_analyzer.py (fallback)",
            "analysis_timestamp": datetime.now(timezone.utc).isoformat(),
            "analysis_fallback": True,
            "analysis_errors": [str(exc)] + inventory_warnings,
            "render_layouts_source": _profile_path(layout_paths["render_layouts"]),
            "text_fit_rules_source": _profile_path(layout_paths["text_fit_rules"]),
            "layout_budget_source": _profile_path(layout_paths["layout_budget"]),
            "ppt_mapping_source": _profile_path(layout_paths["ppt_mapping"]),
            "visual_style": {
                "colors": DEFAULT_STYLE["colors"],
                "typography": DEFAULT_STYLE["typography"],
            },
            "text_geometry": {"analysis_fallback": True},
            "template_inventory": inventory,
            "render_layouts": fallback_render_layouts,
            "page_type_capability": _build_page_type_capability([], inventory),
            "source_area": _build_source_area(inventory, fallback_source_policy),
            "density_budget": _build_density_budget(fallback_layout_budget, []),
            "dynamic_slots": {
                "slide_count": inventory.get("slide_count", 0),
                "slides": [
                    {
                        "slide_no": item.get("slide_no"),
                        "supports": item.get("supports", {}),
                        "information_density": item.get("information_density", "unknown"),
                        "slot_count": item.get("shape_count", 0),
                    }
                    for item in inventory.get("slides", [])
                    if isinstance(item, dict)
                ],
            },
            "layout": {
                "render_layouts": fallback_render_layouts,
                "text_fit_rules": {
                    "fields": _load_json(layout_paths["text_fit_rules"]).get("fields", {}),
                    "renderer_field_aliases": _load_json(layout_paths["text_fit_rules"]).get("renderer_field_aliases", {}),
                },
                "layout_budget": fallback_layout_budget,
            },
            "slide_variants": [],
            "source_policy": fallback_source_policy,
            "output_path": str(output),
            "template_registry_file": _profile_path(layout_paths["slide_registry"]),
        }

    # Keep warnings in both schema and report channels explicit.
    errors = _validate_profile(profile)
    profile_warnings = profile.get("analysis_errors") or []
    output.write_text(json.dumps(profile, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    report = {
        "schema_version": "template_analyzer_v1",
        "is_valid": not errors,
        "error_count": len(errors),
        "warning_count": len(profile_warnings),
        "errors": errors,
        "warnings": profile_warnings,
        "output": str(output),
    }
    print(json.dumps(report, ensure_ascii=False, indent=2))
    return 0 if not errors else 1


def main(argv: list[str] | None = None) -> int:
    argv = list(_ib_sys.argv[1:] if argv is None else argv)
    if argv and argv[0] == "registry":
        return registry_cli(argv[1:])
    if argv and argv[0] == "fit":
        return fit_cli(argv[1:])
    return profile_cli(argv)


if __name__ == "__main__":
    raise SystemExit(main())
