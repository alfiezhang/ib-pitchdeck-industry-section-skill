#!/usr/bin/env python3
"""Render an explicitly non-client-ready PPT draft from page_argument_pack.json.

This is the official quick-draft path. It exists to prevent agents from writing
ad-hoc render_deck.py scripts when evidence is thin but a visible page-shape
draft is useful for internal review.
"""

from __future__ import annotations

import argparse
import json
import sys
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
except Exception as exc:  # pragma: no cover - dependency diagnostics
    raise SystemExit(f"python-pptx is required for quick draft rendering: {exc}")


ROOT_DIR = next(
    p for p in Path(__file__).resolve().parents
    if (p / "configs").is_dir() and (p / "assets").is_dir()
)
DEFAULT_TEMPLATE = ROOT_DIR / "assets" / "industry_section_template_master.pptx"


def _load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def _write_json(path: Path, payload: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def _selected_template(run_dir: Path, explicit_template: str) -> Path:
    if explicit_template:
        return Path(explicit_template)
    selection = run_dir / "artifacts" / "template_selection.json"
    if selection.exists():
        try:
            selected = _load_json(selection).get("selected_template_path")
            if selected:
                return Path(selected)
        except Exception:
            pass
    manifest = run_dir / "artifacts" / "material_manifest.json"
    if manifest.exists():
        try:
            for item in _load_json(manifest).get("materials", []):
                if not isinstance(item, dict):
                    continue
                if str(item.get("source_type") or "").strip() == "ppt_template":
                    raw_candidate = Path(str(item.get("file_path_or_url") or ""))
                    candidates = [raw_candidate] if raw_candidate.is_absolute() else [run_dir / raw_candidate, raw_candidate]
                    for candidate in candidates:
                        if candidate.exists():
                            return candidate
        except Exception:
            pass
    return DEFAULT_TEMPLATE


def _clear_slides(prs: Presentation) -> None:
    slide_id_list = prs.slides._sldIdLst  # noqa: SLF001 - python-pptx has no public delete API
    for slide_id in list(slide_id_list):
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        slide_id_list.remove(slide_id)


def _blank_layout(prs: Presentation):
    if len(prs.slide_layouts) >= 7:
        return prs.slide_layouts[6]
    return prs.slide_layouts[-1]


def _text(shape, value: str, *, font_size: int = 18, bold: bool = False, color: tuple[int, int, int] = (34, 34, 34)) -> None:
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = value
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)


def _add_bullets(slide, bullets: list[str], x, y, w, h, *, font_size: int = 15) -> None:
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    for idx, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(48, 48, 48)


def _argument_rows(pack: dict[str, Any]) -> list[dict[str, Any]]:
    rows = pack.get("page_arguments")
    if isinstance(rows, list):
        return [row for row in rows if isinstance(row, dict)]
    return []


def render_quick_draft(run_dir: Path, *, template_path: Path, output: Path) -> dict[str, Any]:
    page_argument_path = run_dir / "artifacts" / "page_argument_pack.json"
    if not page_argument_path.exists():
        raise SystemExit(f"missing page argument pack: {page_argument_path}")
    pack = _load_json(page_argument_path)
    rows = _argument_rows(pack)
    if not rows:
        raise SystemExit("page_argument_pack.json has no page_arguments to render")

    prs = Presentation(str(template_path))
    _clear_slides(prs)
    layout = _blank_layout(prs)

    slide = prs.slides.add_slide(layout)
    _text(
        slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.6)),
        "Industry Section Draft (Not Client Ready)",
        font_size=28,
        bold=True,
        color=(20, 83, 88),
    )
    _text(
        slide.shapes.add_textbox(Inches(0.62), Inches(1.15), Inches(11.5), Inches(0.55)),
        "Generated from page_argument_pack.json for internal page-shape review only.",
        font_size=15,
        color=(85, 85, 85),
    )
    first_bullets = [
        "This draft is not a formal client-ready deliverable.",
        "It may omit charts, source formatting, and final template fit.",
        "Use it to review story shape before formal evidence/QC/render pipeline work.",
    ]
    _add_bullets(slide, first_bullets, Inches(0.75), Inches(2.0), Inches(11.2), Inches(1.8))

    for idx, row in enumerate(rows[:12], start=1):
        slide = prs.slides.add_slide(layout)
        title = str(row.get("page_argument") or row.get("client_question") or f"Page Argument {idx}").strip()
        if len(title) > 115:
            title = title[:112].rstrip() + "..."
        _text(
            slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.3), Inches(0.8)),
            title or f"Page Argument {idx}",
            font_size=23,
            bold=True,
            color=(20, 83, 88),
        )
        meta = [
            f"Source issue: {row.get('source_issue_analysis_id', '')}",
            f"Evidence status: {row.get('evidence_status', '')}",
            f"Allowed usage: {row.get('allowed_deck_usage', '')}",
        ]
        evidence_ids = row.get("evidence_ids") if isinstance(row.get("evidence_ids"), list) else []
        metric_ids = row.get("metric_ids") if isinstance(row.get("metric_ids"), list) else []
        if evidence_ids:
            meta.append("Evidence IDs: " + ", ".join(map(str, evidence_ids[:6])))
        if metric_ids:
            meta.append("Metric IDs: " + ", ".join(map(str, metric_ids[:6])))
        caveat = str(row.get("caveat_or_diligence_question") or "").strip()
        if caveat:
            meta.append("Caveat / diligence: " + caveat)
        _add_bullets(slide, meta, Inches(0.75), Inches(1.45), Inches(11.5), Inches(3.8), font_size=15)
        _text(
            slide.shapes.add_textbox(Inches(0.75), Inches(6.65), Inches(11.4), Inches(0.35)),
            "DRAFT_NOT_CLIENT_READY | page-argument sketch only",
            font_size=10,
            color=(132, 132, 132),
        )

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    marker = run_dir / "DRAFT_NOT_CLIENT_READY.txt"
    marker.write_text(
        "This run contains an internal evidence-limited quick draft. It is not client-ready final delivery.\n",
        encoding="utf-8",
    )
    manifest = {
        "schema_version": "quick_draft_manifest_v1",
        "created_at": datetime.now(timezone.utc).isoformat(),
        "draft_output_only": True,
        "client_ready": False,
        "draft_ppt": str(output),
        "source": str(page_argument_path),
        "template": str(template_path),
        "policy": "Official quick draft from page_argument_pack.json. Do not present as formal delivery.",
    }
    _write_json(run_dir / "artifacts" / "quick_draft_manifest.json", manifest)
    return manifest


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--run-dir", required=True)
    parser.add_argument("--template", default="")
    parser.add_argument("--output", default="")
    args = parser.parse_args()

    run_dir = Path(args.run_dir).resolve()
    template_path = _selected_template(run_dir, args.template)
    output = Path(args.output).resolve() if args.output else run_dir / "industry_section_QUICK_DRAFT_NOT_CLIENT_READY.pptx"
    manifest = render_quick_draft(run_dir, template_path=template_path, output=output)
    print(json.dumps({"is_valid": True, **manifest}, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
