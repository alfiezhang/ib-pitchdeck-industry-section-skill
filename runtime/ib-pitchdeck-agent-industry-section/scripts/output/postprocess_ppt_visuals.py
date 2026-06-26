#!/usr/bin/env python3
# Runtime scripts can be run directly. Shared helpers remain in runtime
# `scripts/`; production tools live under role scripts; validators live under QC.
import sys as _ib_sys
from pathlib import Path as _IbPath
_IB_ROLE_SCRIPT_DIR = _IbPath(__file__).resolve().parent
_IB_RUNTIME_ROOT = next(
    _p for _p in _IbPath(__file__).resolve().parents
    if (_p / 'configs').is_dir() and (_p / 'scripts').is_dir()
)
_IB_SHARED_SCRIPT_DIR = _IB_RUNTIME_ROOT / "scripts" / "_lib"
_IB_ROLE_SCRIPT_DIRS = sorted(_p for _p in (_IB_RUNTIME_ROOT / 'scripts').iterdir() if _p.is_dir())
_IB_QC_VALIDATOR_DIRS = sorted((_IB_RUNTIME_ROOT / 'scripts' / 'qc' / 'validators').glob('*'))
_IB_IMPORT_PATHS = [str(_IB_ROLE_SCRIPT_DIR)]
for _ib_dir in [*_IB_ROLE_SCRIPT_DIRS, *_IB_QC_VALIDATOR_DIRS]:
    _ib_text = str(_ib_dir)
    if _ib_text not in _IB_IMPORT_PATHS:
        _IB_IMPORT_PATHS.append(_ib_text)
_IB_IMPORT_PATHS.append(str(_IB_SHARED_SCRIPT_DIR))
for _ib_path in list(_IB_IMPORT_PATHS):
    if _ib_path in _ib_sys.path:
        _ib_sys.path.remove(_ib_path)
for _ib_path in reversed(_IB_IMPORT_PATHS):
    _ib_sys.path.insert(0, _ib_path)

import argparse
import json
import os
import re
import shutil
import tempfile
from pathlib import Path
from typing import Any, Optional, Union
from zipfile import ZIP_DEFLATED, ZipFile

from runtime_utils import load_json_file

try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
except ImportError as exc:
    raise SystemExit(
        "python-pptx is required for postprocess_ppt_visuals.py. "
        "Run this script with a Python environment that has the `pptx` package installed, "
        "such as the project virtualenv created by `./setup.sh`."
    ) from exc


DEBUG_MARKER = "DEBUG_OUTPUT_ONLY.txt"


def _default_max_repair_cycles() -> int:
    path = _IB_RUNTIME_ROOT / "configs" / "workflow_policy.json"
    try:
        payload = load_json_file(path)
        gate_retry = payload.get("gate_retry") if isinstance(payload, dict) else {}
        return int(gate_retry.get("default_max_repair_cycles") or 3) if isinstance(gate_retry, dict) else 3
    except Exception:
        return 3


DEFAULT_MAX_REPAIR_CYCLES = _default_max_repair_cycles()


def _gate_state_path(run_dir: Path) -> Path:
    return run_dir / "artifacts" / "gate_retry_state.json"


def _load_gate_state(run_dir: Path) -> dict[str, Any]:
    path = _gate_state_path(run_dir)
    if not path.exists():
        return {"schema_version": "gate_retry_state_v1", "gates": {}}
    data = load_json_file(path)
    if not isinstance(data, dict):
        return {"schema_version": "gate_retry_state_v1", "gates": {}}
    data.setdefault("schema_version", "gate_retry_state_v1")
    data.setdefault("gates", {})
    return data


def _check_gate(run_dir: Path, gate: str, *, max_repair_cycles: int = DEFAULT_MAX_REPAIR_CYCLES) -> dict[str, Any]:
    state = _load_gate_state(run_dir)
    gate_state = (state.get("gates") or {}).get(gate) or {}
    failed_count = int(gate_state.get("failed_validation_count") or 0)
    blocked = gate_state.get("status") == "blocked" or failed_count > max_repair_cycles
    return {
        "is_blocked": blocked,
        "gate": gate,
        "run_dir": str(run_dir),
        "failed_validation_count": failed_count,
        "max_repair_cycles": max_repair_cycles,
        "state": gate_state,
    }


def _mark_ungated_debug_run(run_dir: Path) -> None:
    run_dir.mkdir(parents=True, exist_ok=True)
    marker = run_dir / DEBUG_MARKER
    if not marker.exists():
        marker.write_text(
            "This run used --allow-ungated-debug / IB_SKILL_ALLOW_UNGATED_DEBUG=1.\n"
            "It is not a formal delivery package. Do not copy generated PPTX files to final-looking names,\n"
            "do not update LATEST_FINAL_PPT.txt, and do not describe it as client-ready.\n",
            encoding="utf-8",
        )


def require_debug_output_name(output_path: Path) -> None:
    if output_path.suffix.lower() != ".pptx":
        return
    if "DEBUG_NOT_FOR_DELIVERY" not in output_path.name:
        raise RuntimeError(
            "ungated PPT output must include 'DEBUG_NOT_FOR_DELIVERY' in the filename. "
            "Debug PPTs must not use final-looking names."
        )


def _looks_like_formal_run(run_dir: Path) -> bool:
    formal_foundation_markers = (
        run_dir / "input_card.json",
        run_dir / "artifacts/research_evidence_db.json",
        run_dir / "industry_research_pack.md",
        run_dir / "banker_page_pack.json",
        run_dir / "industry_section_filled.pptx",
        run_dir / "artifacts/run_flags.json",
    )
    formal_core_artifacts = (
        run_dir / "artifacts/research_evidence_db.json",
        run_dir / "industry_research_pack.md",
        run_dir / "banker_page_pack.json",
        run_dir / "deck_blueprint.json",
        run_dir / "template_registry.json",
        run_dir / "renderer_spec.json",
        run_dir / "replacement_dict.json",
        run_dir / "artifacts/industry_scope_pack.json",
    )
    evidence_chain = (
        run_dir / "artifacts/formal_search_plan.json",
        run_dir / "artifacts/source_archive/source_archive_index.json",
        run_dir / "artifacts/formal_research_execution_report.json",
    )
    foundation = any(path.exists() for path in formal_foundation_markers)
    core = any(path.exists() for path in formal_core_artifacts)
    return foundation and core and sum(path.exists() for path in evidence_chain) >= 1


def _blocked_retry_gates(run_dir: Path) -> list[str]:
    state = _load_gate_state(run_dir)
    gates = state.get("gates") if isinstance(state, dict) else {}
    if not isinstance(gates, dict):
        return []
    return [
        str(gate)
        for gate, gate_state in gates.items()
        if isinstance(gate_state, dict) and gate_state.get("status") == "blocked"
    ]


def _pre_ppt_gate_is_passing(run_dir: Path) -> bool:
    gate_path = run_dir / "artifacts" / "stage_gate_pre_ppt_validation.json"
    if not gate_path.exists():
        return False
    try:
        gate = load_json_file(gate_path)
    except Exception:
        return False
    return isinstance(gate, dict) and gate.get("is_valid") is True


def _reject_debug_on_formal_run_if_needed(run_dir: Path) -> None:
    if not _looks_like_formal_run(run_dir):
        return
    blocked = _blocked_retry_gates(run_dir)
    if blocked:
        raise RuntimeError(
            "ungated debug output is not allowed for this formal run package because "
            f"gate(s) are blocked after repeated failures: {', '.join(blocked)}. "
            "Run scripts/pipeline.py next and report the blocker instead of generating downstream artifacts."
        )
    if not _pre_ppt_gate_is_passing(run_dir):
        raise RuntimeError(
            "ungated debug output is not allowed for a formal run package without a passing pre-PPT gate. "
            "Use an isolated temporary directory for template/render diagnostics, or fix the formal package first."
        )


def require_pre_ppt_gate(run_dir: Path, *, allow_ungated_debug: bool = False) -> None:
    if allow_ungated_debug:
        if os.environ.get("IB_SKILL_ALLOW_UNGATED_DEBUG") == "1":
            _reject_debug_on_formal_run_if_needed(run_dir)
            _mark_ungated_debug_run(run_dir)
            return
        raise RuntimeError(
            "--allow-ungated-debug was requested, but IB_SKILL_ALLOW_UNGATED_DEBUG=1 is not set. "
            "This bypass is reserved for explicit local diagnostics and must not be used for delivery."
        )

    retry_state = _check_gate(run_dir, "pre_ppt", max_repair_cycles=DEFAULT_MAX_REPAIR_CYCLES)
    if retry_state.get("is_blocked"):
        failed_count = retry_state.get("failed_validation_count", 0)
        max_cycles = retry_state.get("max_repair_cycles", DEFAULT_MAX_REPAIR_CYCLES)
        raise RuntimeError(
            "pre-PPT gate is blocked after repeated failures; refusing PPT output. "
            f"failed_validation_count={failed_count}, max_repair_cycles={max_cycles}. "
            f"State: {run_dir / 'artifacts' / 'gate_retry_state.json'}"
        )

    gate_path = run_dir / "artifacts" / "stage_gate_pre_ppt_validation.json"
    if not gate_path.exists():
        raise RuntimeError(
            f"missing required pre-PPT gate artifact: {gate_path}. "
            "Run scripts/pipeline.py validate --artifact pre_ppt first, or use --allow-ungated-debug only for local diagnostics."
        )
    try:
        gate = load_json_file(gate_path)
    except Exception as exc:
        raise RuntimeError(f"cannot read pre-PPT gate artifact {gate_path}: {exc}") from exc
    if not isinstance(gate, dict) or gate.get("is_valid") is not True:
        errors = gate.get("errors", []) if isinstance(gate, dict) else []
        preview = "; ".join(str(item) for item in errors[:5])
        if len(errors) > 5:
            preview += f"; plus {len(errors) - 5} more"
        raise RuntimeError(
            "pre-PPT gate is not passing; refusing to generate or mutate PPT output. "
            f"Gate: {gate_path}. {preview}"
        )


SCAFFOLD_LABELS = {
    "PRIMARY CHART",
    "CHART / VISUAL",
    "MINI TABLE / SEGMENT CUT",
    "POINT 1",
    "POINT 2",
    "POINT 3",
    "STANDARD",
    "SUMMARY_PAGE",
    "CHART_PAGE",
    "CHART_PLUS_MINI_TABLE_PAGE",
    "DRIVER_CARD_PAGE",
    "VALUE_CHAIN_PAGE",
    "MOAT_PAGE",
    "COMPARE_TABLE_PAGE",
    "MATRIX_PAGE",
    "TREND_PAGE",
    "TIMELINE_PAGE",
    "THESIS_SUMMARY_PAGE",
    "DRIVER_CARDS_PAGE",
    "VALUE_CHAIN_ARCHITECTURE_PAGE",
    "MOAT_BARRIER_PAGE",
    "PRIORITY_TREND_PAGE",
    "TRANSACTION_IMPLICATION_PAGE",
    "PEER COMPARE TABLE",
    "CRx / STRUCTURE",
    "COMPETITION DIMENSIONS",
    "PRIORITY TREND",
    "SECONDARY TREND",
    "INDUSTRY ATTRACTIVENESS",
    "KEY INDUSTRY CHANGES BENEFITING TARGET",
    "KEY MESSAGES",
    "TARGET RELATIVE POSITIONING",
    "WATCHLIST",
    "OPEN " + "D" + "D QUESTIONS",
    "UPSTREAM",
    "MIDSTREAM",
    "DOWNSTREAM",
    "PROFIT POOL",
    "KEY BARRIERS",
    "TARGET POSITIONING",
    "BUYER LOGIC",
    "DILI" + "GENCE FOCUS",
    "EVIDENCE " + "GAPS",
    "KEY TAKEAWAYS",
    "KEY TAKEAWAYS FOR TARGET",
    "DRIVER 1",
    "DRIVER 2",
    "DRIVER 3",
    "DRIVER 4",
    "DRIVER 5",
    "DRIVER 6",
    "TREND 1",
    "TREND 2",
    "TREND 3",
    "TREND 4",
    "TREND 5",
    "TREND 6",
    "BARRIER 1",
    "BARRIER 2",
    "BARRIER 3",
    "industry_overview",
    "market_size_segmentation",
    "key_industry_drivers",
    "value_chain_profit_pool",
    "key_barriers_value_drivers",
    "competitive_landscape",
    "industry_trends_future_evolution",
    "industry_takeaways_for_project",
}


def normalize_compare_table_payload(slide_data: dict) -> tuple[list[str], list[list[str]]]:
    compare_table = slide_data.get("compare_table_data")
    if isinstance(compare_table, dict):
        headers = [str(item).strip() for item in compare_table.get("headers") or []]
        rows: list[list[str]] = []
        for row in compare_table.get("rows") or []:
            if not isinstance(row, dict):
                continue
            label = str(row.get("label") or "").strip()
            cells = [str(item).strip() for item in row.get("cells") or []]
            rows.append([label] + cells)
        return headers, rows
    return [], []

DEFAULT_RENDER_LAYOUTS_PATH = _IB_RUNTIME_ROOT / "configs" / "render_layouts.json"
DEFAULT_TEMPLATE_PROFILE_PATH = _IB_RUNTIME_ROOT / "configs" / "template_profile.json"

DEFAULT_BRAND_BLUE = RGBColor(0x0D, 0x57, 0xAA)
DEFAULT_GRID_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
DEFAULT_TEXT_GRAY = RGBColor(0x55, 0x55, 0x55)
DEFAULT_ACCENT_RED = RGBColor(0xC0, 0x3C, 0x28)

BRAND_BLUE = DEFAULT_BRAND_BLUE
GRID_GRAY = DEFAULT_GRID_GRAY
TEXT_GRAY = DEFAULT_TEXT_GRAY
ACCENT_RED = DEFAULT_ACCENT_RED
LEGEND_FONT_SIZE = 8
BODY_FONT = "Microsoft YaHei"
TABLE_HEADER_FONT_SIZE = 10.0
TABLE_BODY_FONT_SIZE = 10.0


def _parse_hex_rgb(value: object, fallback: RGBColor) -> RGBColor:
    if not value:
        return fallback
    text = str(value).strip().lstrip("#")
    if len(text) != 6:
        return fallback
    if any(ch not in "0123456789ABCDEFabcdef" for ch in text):
        return fallback
    return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))


def _rgb_components(color: RGBColor) -> tuple[int, int, int]:
    return int(color[0]), int(color[1]), int(color[2])


def _relative_luminance(color: RGBColor) -> float:
    def channel(value: int) -> float:
        scaled = value / 255
        return scaled / 12.92 if scaled <= 0.03928 else ((scaled + 0.055) / 1.055) ** 2.4

    red, green, blue = _rgb_components(color)
    return 0.2126 * channel(red) + 0.7152 * channel(green) + 0.0722 * channel(blue)


def _too_light_for_white_background(color: RGBColor) -> bool:
    return _relative_luminance(color) >= 0.78


def _parse_font_size(value: object, fallback: float) -> float:
    try:
        parsed = float(value)
    except (TypeError, ValueError):
        return fallback
    return parsed


def _safe_font_name(value: object, fallback: str) -> str:
    text = str(value or fallback).strip()
    return text or fallback


def _apply_template_profile_style(profile_path: Path, warnings: Optional[list[str]] = None) -> None:
    global BRAND_BLUE, GRID_GRAY, TEXT_GRAY, ACCENT_RED
    global LEGEND_FONT_SIZE, BODY_FONT, TABLE_HEADER_FONT_SIZE, TABLE_BODY_FONT_SIZE

    profile_warnings = warnings if warnings is not None else []
    if not profile_path.exists():
        if profile_warnings is not None:
            profile_warnings.append(f"template profile missing, using static fallback style: {profile_path}")
        return

    data = load_json(profile_path)
    visual = data.get("visual_style", {}) if isinstance(data, dict) else {}
    colors = visual.get("colors", {}) if isinstance(visual, dict) else {}
    typography = visual.get("typography", {}) if isinstance(visual, dict) else {}

    brand_primary = _parse_hex_rgb(colors.get("brand_primary"), DEFAULT_BRAND_BLUE)
    text_gray = _parse_hex_rgb(colors.get("text_gray"), DEFAULT_TEXT_GRAY)
    if _too_light_for_white_background(brand_primary):
        profile_warnings.append(
            "template profile brand_primary is too light for a white style-guided page; using fallback title color"
        )
        BRAND_BLUE = DEFAULT_BRAND_BLUE
    else:
        BRAND_BLUE = brand_primary
    if _too_light_for_white_background(text_gray):
        profile_warnings.append(
            "template profile text_gray is too light for a white style-guided page; using fallback body color"
        )
        TEXT_GRAY = DEFAULT_TEXT_GRAY
    else:
        TEXT_GRAY = text_gray
    ACCENT_RED = _parse_hex_rgb(colors.get("accent_red"), DEFAULT_ACCENT_RED)
    GRID_GRAY = _parse_hex_rgb(colors.get("grid_gray"), DEFAULT_GRID_GRAY)

    LEGEND_FONT_SIZE = _parse_font_size(typography.get("legend_pt"), LEGEND_FONT_SIZE)
    TABLE_HEADER_FONT_SIZE = _parse_font_size(typography.get("table_header_pt"), TABLE_HEADER_FONT_SIZE)
    TABLE_BODY_FONT_SIZE = _parse_font_size(typography.get("table_body_pt"), TABLE_BODY_FONT_SIZE)
    BODY_FONT = _safe_font_name(typography.get("body"), BODY_FONT)

    if visual:
        if not typography.get("body"):
            profile_warnings.append("template profile missing body font; using fallback Microsoft YaHei")
        if not typography.get("table_header"):
            profile_warnings.append("template profile missing table_header font; using fallback Microsoft YaHei")
            typography.setdefault("table_header", BODY_FONT)
        if not typography.get("table_body"):
            profile_warnings.append("template profile missing table_body font; using fallback Microsoft YaHei")
            typography.setdefault("table_body", BODY_FONT)
        if not colors.get("brand_primary"):
            profile_warnings.append("template profile missing brand_primary; using fallback color")
        if not colors.get("accent_red"):
            profile_warnings.append("template profile missing accent_red; using fallback color")
        if not colors.get("grid_gray"):
            profile_warnings.append("template profile missing grid_gray; using fallback color")
        if not colors.get("text_gray"):
            profile_warnings.append("template profile missing text_gray; using fallback color")


def load_json(path: Path) -> dict:
    return load_json_file(path)


def load_render_layouts(path: Path = DEFAULT_RENDER_LAYOUTS_PATH) -> dict[int, dict]:
    data = load_json(path)
    slides = data.get("slides")
    if not isinstance(slides, dict):
        raise ValueError(f"Invalid render layout file {path}: missing object field 'slides'.")
    normalized: dict[int, dict] = {}
    for slide_no_raw, variants in slides.items():
        try:
            slide_no = int(slide_no_raw)
        except (TypeError, ValueError) as exc:
            raise ValueError(f"Invalid slide number in render layout file {path}: {slide_no_raw}") from exc
        if not isinstance(variants, dict):
            raise ValueError(f"Invalid render layout file {path}: slide {slide_no} variants must be an object.")
        normalized[slide_no] = {}
        for page_type, boxes in variants.items():
            if not isinstance(boxes, dict):
                raise ValueError(f"Invalid render layout file {path}: slide {slide_no}/{page_type} boxes must be an object.")
            normalized[slide_no][page_type] = {}
            for box_name, box in boxes.items():
                if box_name == "column_weight_overrides":
                    if not isinstance(box, dict):
                        raise ValueError(
                            f"Invalid render layout file {path}: slide {slide_no}/{page_type}/{box_name} "
                            "must be an object keyed by column count."
                        )
                    normalized_weights = {}
                    for col_count, weights in box.items():
                        try:
                            col_int = int(col_count)
                        except (TypeError, ValueError) as exc:
                            raise ValueError(
                                f"Invalid render layout file {path}: slide {slide_no}/{page_type}/{box_name} "
                                f"has invalid column count {col_count!r}."
                            ) from exc
                        if col_int < 3 or col_int > 6:
                            raise ValueError(
                                f"Invalid render layout file {path}: slide {slide_no}/{page_type}/{box_name} "
                                f"only supports 3-6 columns, found {col_int}."
                            )
                        if not (
                            isinstance(weights, list)
                            and len(weights) == col_int
                            and all(isinstance(value, (int, float)) and value > 0 for value in weights)
                        ):
                            raise ValueError(
                                f"Invalid render layout file {path}: slide {slide_no}/{page_type}/{box_name}/{col_count} "
                                f"must be a positive numeric list with {col_int} entries."
                            )
                        normalized_weights[str(col_int)] = [float(value) for value in weights]
                    normalized[slide_no][page_type][box_name] = normalized_weights
                    continue
                if isinstance(box, bool):
                    normalized[slide_no][page_type][box_name] = box
                    continue
                if not (
                    isinstance(box, list)
                    and len(box) == 4
                    and all(isinstance(value, int) for value in box)
                ):
                    raise ValueError(
                        f"Invalid render layout file {path}: slide {slide_no}/{page_type}/{box_name} "
                        "must be [left, top, width, height] integer EMUs."
                    )
                normalized[slide_no][page_type][box_name] = tuple(box)
    return normalized


def save_json(data: dict, path: Path) -> None:
    with path.open("w", encoding="utf-8") as handle:
        json.dump(data, handle, ensure_ascii=False, indent=2)


def clear_text(shape) -> None:
    if not hasattr(shape, "text_frame"):
        return
    text_frame = shape.text_frame
    text_frame.clear()


def remove_shape(shape) -> None:
    element = shape._element
    element.getparent().remove(element)


def intersects(shape, box: tuple[int, int, int, int]) -> bool:
    left, top, width, height = box
    right = left + width
    bottom = top + height
    shape_left = int(shape.left)
    shape_top = int(shape.top)
    shape_right = shape_left + int(shape.width)
    shape_bottom = shape_top + int(shape.height)
    return not (
        shape_right < left
        or shape_left > right
        or shape_bottom < top
        or shape_top > bottom
    )


def set_single_paragraph(shape, text: str) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.name = BODY_FONT


def is_scaffold_label(text: str) -> bool:
    stripped = " ".join(str(text or "").strip().split())
    if not stripped:
        return False
    upper = stripped.upper()
    scaffold_upper = {label.upper() for label in SCAFFOLD_LABELS}
    if stripped in SCAFFOLD_LABELS or upper in scaffold_upper:
        return True
    if upper.endswith("_PAGE"):
        return True
    if re.fullmatch(r"[a-z][a-z0-9]+(?:_[a-z0-9]+){1,5}", stripped):
        return True
    return False


def remove_scaffold_labels(prs: Presentation) -> list[dict]:
    removed = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if is_scaffold_label(text):
                clear_text(shape)
                removed.append({"slide_no": slide_idx, "label": text, "shape_name": shape.name})
    return removed


def is_footer_page_number_candidate(shape, slide_width: int, slide_height: int) -> bool:
    if not hasattr(shape, "text"):
        return False
    text = shape.text.strip()
    if not text.isdigit():
        return False
    return shape.left >= slide_width * 0.55 and shape.top >= slide_height * 0.85


def rewrite_page_numbers(prs: Presentation) -> list[dict]:
    updates = []
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if is_footer_page_number_candidate(shape, slide_width, slide_height):
                old = shape.text.strip()
                set_single_paragraph(shape, str(slide_idx))
                updates.append(
                    {
                        "slide_no": slide_idx,
                        "shape_name": shape.name,
                        "old": old,
                        "new": str(slide_idx),
                    }
                )
                break
    return updates


def find_slide_data(renderer_spec: dict, slide_no: int) -> Optional[dict]:
    for slide in renderer_spec.get("slides", []):
        if slide.get("slide_no") == slide_no:
            return slide
    return None


def apply_chart_title(slide, text: str, layout: dict) -> bool:
    if not text:
        return False
    left, top, width, height = layout["title_box"]
    candidates = []
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
        if abs(shape.left - left) < 20000 and abs(shape.top - top) < 20000:
            candidates.append(shape)
    if not candidates:
        return False
    target = sorted(candidates, key=lambda shp: (abs(shp.width - width), abs(shp.height - height)))[0]
    set_single_paragraph(target, text)
    paragraph = target.text_frame.paragraphs[0]
    if paragraph.runs:
        font = paragraph.runs[0].font
        font.name = BODY_FONT
        font.size = Pt(10)
        font.color.rgb = TEXT_GRAY
        font.bold = False
    return True


def chart_type_for(chart_data: dict):
    chart_type = str(chart_data.get("chart_type") or "bar").lower()
    if chart_type in {"bar", "column", "clustered_bar", "clustered_column"}:
        return XL_CHART_TYPE.COLUMN_CLUSTERED
    if chart_type in {"stacked_bar", "stacked_column"}:
        return XL_CHART_TYPE.COLUMN_STACKED
    if chart_type in {"line", "line_chart"}:
        return XL_CHART_TYPE.LINE_MARKERS
    return None


def resolve_chart_number_format(chart_data: dict) -> str:
    unit = str(chart_data.get("unit") or "").strip().lower()
    values = []
    for chart_series in chart_data.get("series") or []:
        if isinstance(chart_series, dict):
            values.extend(chart_series.get("values") or [])
    if "%" in unit or "％" in unit or any(token in unit for token in ("percent", "percentage", "share", "rate", "ratio", "占比", "份额", "比例", "率")):
        numeric_values = []
        for value in values:
            try:
                numeric_values.append(abs(float(value)))
            except (TypeError, ValueError):
                continue
        if numeric_values and max(numeric_values) <= 1.0:
            return "0.0%"
        return '0.0"%"'
    if unit in {"x", "倍", "turn", "turns"} or "multiple" in unit:
        return '0.0"x"'
    if any(token in unit for token in ("亿元", "亿", "bn", "billion")):
        return '#,##0.0'
    if any(token in unit for token in ("万元", "万", "mn", "million")):
        return '#,##0'
    return '#,##0'


def short_label(value: str, max_units: float = 14.0) -> str:
    text = str(value or "").strip()
    units = 0.0
    result = []
    for ch in text:
        if ch.isspace():
            char_units = 0.3
        elif ord(ch) < 128:
            char_units = 0.55
        else:
            char_units = 1.0
        if units + char_units > max_units:
            return "".join(result).rstrip() + "..."
        result.append(ch)
        units += char_units
    return text


def format_chart_legend(chart) -> None:
    if not chart.has_legend:
        return
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(LEGEND_FONT_SIZE)
    chart.legend.font.name = BODY_FONT
    chart.legend.font.color.rgb = TEXT_GRAY
    chart.legend.font.bold = False


def build_chart(slide, slide_data: dict, layout: dict) -> dict:
    chart_data = slide_data.get("chart_data") or {}
    series = chart_data.get("series") or []
    categories = chart_data.get("categories") or []
    slide_no = slide_data.get("slide_no")
    page_type = slide_data.get("selected_page_type")
    if not isinstance(categories, list):
        return {
            "rendered": False,
            "reason": "invalid chart_data.categories",
            "slide_no": slide_no,
            "selected_page_type": page_type,
            "path": "chart_data.categories",
            "expected": "list[str]",
            "actual": type(categories).__name__,
            "repair_hint": "Compile deck_blueprint chart data into categories + series[{name, values}] before postprocess.",
        }
    if not isinstance(series, list):
        return {
            "rendered": False,
            "reason": "invalid chart_data.series",
            "slide_no": slide_no,
            "selected_page_type": page_type,
            "path": "chart_data.series",
            "expected": "list[object] where each object has name and values",
            "actual": type(series).__name__,
            "repair_hint": "Use scripts/pipeline.py compile to normalize natural data_series into renderer chart series.",
        }
    non_object_series = [idx for idx, item in enumerate(series, start=1) if not isinstance(item, dict)]
    if non_object_series:
        return {
            "rendered": False,
            "reason": "invalid chart_data.series item",
            "slide_no": slide_no,
            "selected_page_type": page_type,
            "path": f"chart_data.series[{non_object_series[0]}]",
            "expected": "object with name and values",
            "actual": type(series[non_object_series[0] - 1]).__name__,
            "repair_hint": "Do not pass a list of legend strings to postprocess. Compiler should output series: [{name, values}].",
        }
    if not series or not categories:
        return {"rendered": False, "reason": "missing chart_data series/categories"}

    ppt_chart_type = chart_type_for(chart_data)
    if ppt_chart_type is None:
        return {
            "rendered": False,
            "reason": f"unsupported chart_type for deterministic chart renderer: {chart_data.get('chart_type')}",
        }

    chart_title = chart_data.get("title") or ""
    apply_chart_title(slide, chart_title, layout)

    chart_payload = CategoryChartData()
    chart_payload.categories = categories
    for chart_series in series:
        values = chart_series.get("values") or []
        if not isinstance(values, list):
            return {
                "rendered": False,
                "reason": "invalid chart series values",
                "slide_no": slide_no,
                "selected_page_type": page_type,
                "path": "chart_data.series[].values",
                "expected": "list[number]",
                "actual": type(values).__name__,
                "repair_hint": "Compiler should emit numeric values in each chart series.",
            }
        if len(values) != len(categories):
            return {
                "rendered": False,
                "reason": "series/category length mismatch",
                "slide_no": slide_no,
                "selected_page_type": page_type,
                "path": "chart_data.series[].values",
                "expected": f"{len(categories)} values to match categories",
                "actual": f"{len(values)} values",
                "repair_hint": "Every chart series must have exactly one value per category.",
            }
        chart_payload.add_series(short_label(chart_series.get("name") or ""), values)

    chart_box = layout["chart_box"]
    removed = []
    if not layout.get("preserve_existing_shapes"):
        removed = remove_text_shapes_in_box(slide, chart_box)
    left, top, width, height = chart_box
    if len(series) > 1:
        height = int(height * 0.88)
    graphic_frame = slide.shapes.add_chart(
        ppt_chart_type,
        Emu(left),
        Emu(top),
        Emu(width),
        Emu(height),
        chart_payload,
    )
    chart = graphic_frame.chart
    chart.has_title = False
    chart.has_legend = len(series) > 1
    format_chart_legend(chart)

    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(9)
    category_axis.tick_labels.font.name = BODY_FONT
    category_axis.tick_labels.font.color.rgb = TEXT_GRAY

    value_axis = chart.value_axis
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = GRID_GRAY
    value_axis.tick_labels.font.size = Pt(9)
    value_axis.tick_labels.font.name = BODY_FONT
    value_axis.tick_labels.font.color.rgb = TEXT_GRAY
    value_axis.format.line.color.rgb = GRID_GRAY
    number_format = resolve_chart_number_format(chart_data)
    value_axis.tick_labels.number_format = number_format

    normalized_chart_type = str(chart_data.get("chart_type") or "").lower()
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    if normalized_chart_type in {"stacked_bar", "stacked_column"}:
        data_labels.position = XL_DATA_LABEL_POSITION.CENTER
    else:
        data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    data_labels.number_format = number_format
    data_labels.font.size = Pt(9)
    data_labels.font.name = BODY_FONT
    data_labels.font.bold = True

    palette = [BRAND_BLUE, ACCENT_RED, RGBColor(0x6C, 0x8E, 0xB8)]
    for idx, s in enumerate(chart.series):
        fill = s.format.fill
        fill.solid()
        fill.fore_color.rgb = palette[idx % len(palette)]
        s.format.line.color.rgb = palette[idx % len(palette)]

    return {
        "rendered": True,
        "chart_title": chart_title,
        "chart_type": chart_data.get("chart_type"),
        "categories": categories,
        "series_count": len(series),
        "removed_text_shapes": removed,
    }


def _coerce_float(value, default: Optional[float] = None) -> Optional[float]:
    try:
        if value is None or value == "":
            return default
        return float(value)
    except (TypeError, ValueError):
        return default


def _matrix_points_from_chart_data(chart_data: dict) -> list[dict]:
    rows = chart_data.get("source_rows") or []
    points = []
    for row in rows:
        if not isinstance(row, dict):
            continue
        label = str(row.get("label") or "").strip()
        x = _coerce_float(row.get("x"))
        y = _coerce_float(row.get("y"))
        if x is None or y is None:
            value = row.get("value")
            if isinstance(value, dict):
                x = _coerce_float(value.get("x"))
                y = _coerce_float(value.get("y"))
            elif isinstance(value, (list, tuple)) and len(value) >= 2:
                x = _coerce_float(value[0])
                y = _coerce_float(value[1])
        if label and x is not None and y is not None:
            points.append({"label": label, "x": x, "y": y, "note": row.get("note", "")})

    if points:
        return points

    categories = chart_data.get("categories") or []
    series = chart_data.get("series") or []
    if len(series) < 2:
        return []
    if not isinstance(series[0], dict) or not isinstance(series[1], dict):
        return []
    x_values = series[0].get("values") or []
    y_values = series[1].get("values") or []
    for idx, label in enumerate(categories):
        if idx >= len(x_values) or idx >= len(y_values):
            continue
        x = _coerce_float(x_values[idx])
        y = _coerce_float(y_values[idx])
        if label and x is not None and y is not None:
            points.append({"label": str(label), "x": x, "y": y, "note": ""})
    return points


def _normalize(value: float, values: list[float]) -> float:
    min_value = min(values)
    max_value = max(values)
    if max_value == min_value:
        return 0.5
    return (value - min_value) / (max_value - min_value)


def _add_textbox(slide, left: int, top: int, width: int, height: int, text: str, font_size: int = 8, bold: bool = False) -> None:
    textbox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.name = BODY_FONT
    run.font.bold = bold
    run.font.color.rgb = TEXT_GRAY


def render_matrix_slide(slide, slide_data: dict, layout: dict) -> dict:
    chart_data = slide_data.get("chart_data") or {}
    body_copy = slide_data.get("body_copy") or {}
    points = _matrix_points_from_chart_data(chart_data)
    if len(points) < 2:
        return {"rendered": False, "reason": "matrix_page needs at least two points with x/y values"}

    cleanup_box = layout.get("cleanup_box")
    removed_shapes = []
    if cleanup_box:
        for shape in list(slide.shapes):
            if intersects(shape, cleanup_box):
                text = getattr(shape, "text", "").strip() if hasattr(shape, "text") else ""
                remove_shape(shape)
                removed_shapes.append({"shape_name": shape.name, "text": text[:80]})

    left, top, width, height = layout["matrix_box"]
    axis_label_x = body_copy.get("matrix_label_x") or chart_data.get("x_axis_label") or "Axis X"
    axis_label_y = body_copy.get("matrix_label_y") or chart_data.get("y_axis_label") or "Axis Y"

    # Background panel.
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Emu(left),
        Emu(top),
        Emu(width),
        Emu(height),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFC)
    panel.line.color.rgb = GRID_GRAY
    panel.line.width = Pt(1)

    mid_x = left + width // 2
    mid_y = top + height // 2
    for line_left, line_top, line_width, line_height in [
        (mid_x, top, 0, height),
        (left, mid_y, width, 0),
    ]:
        line = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Emu(line_left),
            Emu(line_top),
            Emu(max(line_width, 1)),
            Emu(max(line_height, 1)),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = GRID_GRAY
        line.line.color.rgb = GRID_GRAY

    _add_textbox(slide, left, top + height + 80000, width, 180000, axis_label_x, 8, True)
    _add_textbox(slide, left - 240000, top + height // 2 - 120000, 220000, 260000, axis_label_y, 8, True)

    x_values = [point["x"] for point in points]
    y_values = [point["y"] for point in points]
    bubble_size = 175000
    plotted = []
    for idx, point in enumerate(points[:8]):
        x_norm = _normalize(point["x"], x_values)
        y_norm = _normalize(point["y"], y_values)
        cx = left + int(260000 + x_norm * (width - 520000))
        cy = top + int(height - 260000 - y_norm * (height - 520000))

        bubble = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Emu(cx - bubble_size // 2),
            Emu(cy - bubble_size // 2),
            Emu(bubble_size),
            Emu(bubble_size),
        )
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = BRAND_BLUE if idx else ACCENT_RED
        bubble.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        bubble.line.width = Pt(1)

        label_left = min(max(cx - 360000, left + 20000), left + width - 720000)
        label_top = min(max(cy + 90000, top + 20000), top + height - 220000)
        _add_textbox(slide, label_left, label_top, 720000, 180000, point["label"], 7, idx == 0)
        plotted.append({"label": point["label"], "x": point["x"], "y": point["y"]})

    return {
        "rendered": True,
        "chart_title": chart_data.get("title") or body_copy.get("matrix_title") or "",
        "chart_type": "matrix",
        "points": plotted,
        "truncated_points": max(0, len(points) - len(plotted)),
        "removed_existing_matrix_shapes": removed_shapes,
    }


def add_metric_card(slide, left: int, top: int, width: int, height: int, label: str, value: str, accent: RGBColor) -> None:
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Emu(left),
        Emu(top),
        Emu(width),
        Emu(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape.line.color.rgb = GRID_GRAY
    shape.line.width = Pt(1)

    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True

    p1 = text_frame.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = label
    r1.font.size = Pt(10)
    r1.font.name = BODY_FONT
    r1.font.bold = True
    r1.font.color.rgb = TEXT_GRAY

    p2 = text_frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = value
    r2.font.size = Pt(24)
    r2.font.name = BODY_FONT
    r2.font.bold = True
    r2.font.color.rgb = accent


def add_supporting_note(slide, left: int, top: int, width: int, height: int, text: str) -> None:
    from pptx.enum.text import PP_ALIGN

    textbox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.name = BODY_FONT
    r.font.color.rgb = TEXT_GRAY


def add_left_textbox(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    font_size: int = 9,
    bold: bool = False,
    color: RGBColor = TEXT_GRAY,
) -> None:
    textbox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.05
    run = p.add_run()
    run.text = str(text or "")
    run.font.size = Pt(font_size)
    run.font.name = BODY_FONT
    run.font.bold = bold
    run.font.color.rgb = color


def add_panel_box(slide, left: int, top: int, width: int, height: int) -> None:
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Emu(left),
        Emu(top),
        Emu(width),
        Emu(height),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFC)
    panel.line.color.rgb = GRID_GRAY
    panel.line.width = Pt(1)


def remove_think_cell_ole_shapes(slide) -> list[dict]:
    removed = []
    for shape in list(slide.shapes):
        name = str(getattr(shape, "name", "") or "")
        if "think-cell" not in name.lower():
            continue
        remove_shape(shape)
        removed.append({"shape_name": name, "shape_type": str(getattr(shape, "shape_type", ""))})
    return removed


def format_metric_value(label: str, value: Union[float, int, str], unit: str) -> str:
    if isinstance(value, str):
        text = value.strip()
        if not text:
            return ""
        if any(token in text for token in ("%", "亿", "万", "元", "RMB", "USD", "$")):
            return text
        try:
            value = float(text.replace(",", ""))
        except ValueError:
            return text
    is_integer = isinstance(value, int) or (isinstance(value, float) and value.is_integer())
    if "CAGR" in label.upper():
        return f"{int(value) if is_integer else f'{value:.1f}'}%"
    if "%" in label or "同比" in label or "增速" in label:
        return f"{value:+.0f}%" if is_integer else f"{value:+.1f}%"
    if "份额" in label or "占比" in label or "比例" in label or "渗透率" in label:
        return f"{int(value)}%" if is_integer else f"{value:.1f}%"
    if "%" in unit and not any(token in unit for token in ("亿", "万", "元", "RMB", "USD", "$")):
        return f"{int(value)}%" if is_integer else f"{value:.1f}%"
    if "亿元" in label or "规模" in label or "亿元" in unit or "人民币" in unit:
        if is_integer:
            return f"{int(value):,}亿元"
        return f"{value:,.1f}亿元"
    if is_integer:
        return f"{int(value):,}"
    return f"{value:,.1f}"


def render_secondary_module(slide, chart_data: dict, layout: dict) -> dict:
    module = chart_data.get("secondary_module") or {}
    if not isinstance(module, dict):
        module = {}
    module_type = str(module.get("module_type") or "metric_cards").lower()
    rows = module.get("rows") or []
    rows = [row for row in rows if isinstance(row, dict)]
    if not rows:
        return {
            "rendered": False,
            "reason": "missing secondary module rows; refused to reuse primary chart source_rows as fallback",
        }

    left, top, width, height = layout["side_box"]
    removed = remove_text_shapes_in_box(slide, layout["side_box"])
    add_panel_box(slide, left, top, width, height)
    title = str(module.get("title") or chart_data.get("secondary_title") or "").strip()
    title_height = 260000 if title else 0
    if title:
        add_left_textbox(slide, left + 120000, top + 70000, width - 240000, title_height, title, 9, True, BRAND_BLUE)

    content_top = top + title_height + 110000
    content_height = height - title_height - 180000
    if module_type in {"mini_table", "table", "segmentation_table"}:
        headers = module.get("headers") or ["Metric", "Read-through"]
        col_count = min(max(len(headers), 2), 3)
        row_count = min(len(rows), 4) + 1
        table_shape = slide.shapes.add_table(
            row_count,
            col_count,
            Emu(left + 90000),
            Emu(content_top),
            Emu(width - 180000),
            Emu(content_height),
        )
        table = table_shape.table
        for col_idx in range(col_count):
            table.columns[col_idx].width = Emu(int((width - 180000) / col_count))
            set_cell_text(table.cell(0, col_idx), headers[col_idx] if col_idx < len(headers) else "", 8.5, True, RGBColor(0xFF, 0xFF, 0xFF))
        for row_idx, row in enumerate(rows[: row_count - 1], start=1):
            values = [
                row.get("label", ""),
                format_metric_value(str(row.get("label") or ""), row.get("value", ""), row.get("unit") or row.get("value_unit") or chart_data.get("unit") or ""),
                row.get("note", ""),
            ]
            for col_idx in range(col_count):
                set_cell_text(table.cell(row_idx, col_idx), values[col_idx] if col_idx < len(values) else "", 8.5, False)
        style_table_shape(table_shape)
        return {"rendered": True, "module_type": module_type, "rows": row_count, "removed_text_shapes": removed}

    card_count = min(max(len(rows), 2), 3)
    gap = 90000
    card_height = (content_height - gap * (card_count - 1)) // card_count
    accents = [ACCENT_RED, BRAND_BLUE, RGBColor(0x4D, 0x7C, 0x3A)]
    for idx, row in enumerate(rows[:card_count]):
        label = row.get("label", "")
        row_unit = row.get("unit") or row.get("value_unit") or chart_data.get("unit") or ""
        add_metric_card(
            slide,
            left + 160000,
            content_top + idx * (card_height + gap),
            width - 320000,
            card_height,
            label,
            format_metric_value(label, row.get("value", ""), row_unit),
            accents[idx],
        )
    return {"rendered": True, "module_type": "metric_cards", "cards": card_count, "removed_text_shapes": removed}


def render_bottom_takeaways(slide, slide_data: dict, layout: dict) -> dict:
    body = slide_data.get("body_copy") or {}
    takeaways = [
        str(body.get("bullet_1") or "").strip(),
        str(body.get("bullet_2") or "").strip(),
    ]
    takeaways = [item for item in takeaways if item]
    if not takeaways:
        return {"rendered": False, "reason": "missing bottom takeaway bullets"}
    left, top, width, height = layout["bottom_box"]
    removed = remove_text_shapes_in_box(slide, layout["bottom_box"])
    line_height = height // max(len(takeaways), 1)
    for idx, text in enumerate(takeaways[:2]):
        add_left_textbox(slide, left, top + idx * line_height, width, line_height, f"• {text}", 9, False)
    return {"rendered": True, "takeaways": len(takeaways[:2]), "removed_text_shapes": removed}


def render_slide1_dynamic_overview(slide, slide_data: dict, layout: dict) -> dict:
    chart_data = slide_data.get("chart_data") or {}
    chart_type = str(chart_data.get("chart_type") or "").lower()
    cleanup_box = layout.get("cleanup_box")
    cleared = []
    if cleanup_box:
        cleared = clear_text_shapes_in_box(slide, cleanup_box)

    if chart_type in {"bar", "column", "clustered_bar", "clustered_column", "stacked_bar", "stacked_column", "line", "line_chart"}:
        chart_result = build_chart(slide, slide_data, layout)
    else:
        chart_result = {
            "rendered": False,
            "reason": f"dynamic overview requires a chart type supported by the deterministic renderer, found '{chart_type}'",
        }
    return {
        "rendered": bool(chart_result.get("rendered")),
        "mode": "dynamic_overview",
        "chart": chart_result,
        "left_key_messages_preserved": True,
        "right_visual_only": True,
        "cleared_existing_visual_text_shapes": cleared,
    }


def render_slide1_visual(slide, slide_data: dict, layout: dict) -> dict:
    chart_data = slide_data.get("chart_data") or {}
    chart_type = str(chart_data.get("chart_type") or "").lower()
    if chart_type in {"none", "no_chart", "text"}:
        return {
            "rendered": False,
            "required_render": False,
            "reason": "slide 1 chart_type is none; visual area intentionally left as text/context",
        }
    if chart_type in {"bar", "column", "clustered_bar", "clustered_column", "stacked_bar", "stacked_column", "line", "line_chart"}:
        return build_chart(slide, slide_data, layout)

    rows = chart_data.get("source_rows") or []
    if not isinstance(rows, list):
        return {
            "rendered": False,
            "reason": "invalid chart_data.source_rows",
            "slide_no": slide_data.get("slide_no"),
            "selected_page_type": slide_data.get("selected_page_type"),
            "path": "chart_data.source_rows",
            "expected": "list[object]",
            "actual": type(rows).__name__,
            "repair_hint": "Use scripts/pipeline.py compile to normalize chart data before postprocess.",
        }
    if len(rows) < 2:
        return {
            "rendered": False,
            "reason": "slide 1 needs chart_data series/categories for charts or at least two source_rows for metric cards",
        }
    non_object_rows = [idx for idx, row in enumerate(rows, start=1) if not isinstance(row, dict)]
    if non_object_rows:
        return {
            "rendered": False,
            "reason": "invalid chart_data.source_rows item",
            "slide_no": slide_data.get("slide_no"),
            "selected_page_type": slide_data.get("selected_page_type"),
            "path": f"chart_data.source_rows[{non_object_rows[0]}]",
            "expected": "object with label/value/unit",
            "actual": type(rows[non_object_rows[0] - 1]).__name__,
            "repair_hint": "Do not pass source_rows as strings; compile deck_blueprint natural data_series into renderer rows.",
        }

    chart_title = chart_data.get("title") or ""
    apply_chart_title(slide, chart_title, layout)

    left, top, width, height = layout["visual_box"]
    card_count = min(max(len(rows), 2), 3)
    gap = 180000
    card_width = (width - gap * (card_count - 1)) // card_count
    display_note = chart_data.get("display_note") or chart_data.get("on_slide_note") or ""
    card_height = int(height * (0.58 if display_note else 0.68))
    note_height = 480000 if display_note else 0
    card_top = top + max(0, int((height - card_height - note_height) / 2))

    unit = chart_data.get("unit") or ""
    accents = [ACCENT_RED, BRAND_BLUE, RGBColor(0x4D, 0x7C, 0x3A)]

    for idx, row in enumerate(rows[:card_count]):
        label = row.get("label", "")
        row_unit = row.get("unit") or row.get("value_unit") or unit
        add_metric_card(
            slide,
            left + idx * (card_width + gap),
            card_top,
            card_width,
            card_height,
            label,
            format_metric_value(label, row.get("value", 0), row_unit),
            accents[idx],
        )

    if display_note:
        add_supporting_note(slide, left, card_top + card_height + 160000, width, note_height, display_note)

    return {
        "rendered": True,
        "chart_title": chart_title,
        "chart_type": chart_data.get("chart_type"),
        "mode": "metric_cards",
    }


def split_table_cells(text: str, expected_cols: Optional[int] = None) -> list[str]:
    value = str(text or "").strip()
    if not value:
        return [""] * expected_cols if expected_cols else []
    if "｜" in value:
        cells = [part.strip() for part in value.split("｜")]
    elif "|" in value:
        cells = [part.strip() for part in value.split("|")]
    elif " / " in value:
        cells = [part.strip() for part in value.split(" / ")]
    else:
        cells = [value]
    if expected_cols is None:
        return cells
    if len(cells) > expected_cols:
        cells = cells[: expected_cols - 1] + [" / ".join(cells[expected_cols - 1 :])]
    return cells + [""] * (expected_cols - len(cells))


def estimate_column_weights(headers: list[str], rows: list[list[str]]) -> list[float]:
    weights: list[float] = []
    col_count = len(headers)
    for col_idx in range(col_count):
        values = [headers[col_idx]] + [row[col_idx] for row in rows if col_idx < len(row)]
        max_len = max((len(str(value)) for value in values), default=1)
        weights.append(max(0.8, min(3.2, max_len / 8.0)))
    if weights:
        weights[0] = min(weights[0], 1.3)
    return weights


def count_table_cells(text: str) -> int:
    value = str(text or "").strip()
    if not value:
        return 0
    if "｜" in value:
        return len([part for part in value.split("｜")])
    if "|" in value:
        return len([part for part in value.split("|")])
    if " / " in value:
        return len([part for part in value.split(" / ")])
    return 1


def remove_text_shapes_in_box(slide, box: tuple[int, int, int, int]) -> list[dict]:
    removed = []
    for shape in list(slide.shapes):
        if not hasattr(shape, "text_frame"):
            continue
        if not intersects(shape, box):
            continue
        text = getattr(shape, "text", "").strip() if hasattr(shape, "text") else ""
        remove_shape(shape)
        removed.append({"shape_name": shape.name, "text": text[:80]})
    return removed


def clear_text_shapes_in_box(slide, box: tuple[int, int, int, int]) -> list[dict]:
    cleared = []
    for shape in list(slide.shapes):
        if not hasattr(shape, "text_frame"):
            continue
        if not intersects(shape, box):
            continue
        text = getattr(shape, "text", "").strip() if hasattr(shape, "text") else ""
        if not text:
            continue
        shape.text_frame.clear()
        cleared.append({"shape_name": shape.name, "text": text[:80]})
    return cleared


def remove_shapes_in_box(slide, box: tuple[int, int, int, int]) -> list[dict]:
    removed = []
    for shape in list(slide.shapes):
        if not intersects(shape, box):
            continue
        text = getattr(shape, "text", "").strip() if hasattr(shape, "text") else ""
        remove_shape(shape)
        removed.append({"shape_name": shape.name, "text": text[:80]})
    return removed


def set_cell_text(cell, text: str, font_size: float, bold: bool = False, color: Optional[RGBColor] = None) -> None:
    if color is None:
        color = TEXT_GRAY
    cell.text = ""
    cell.margin_left = Emu(30000)
    cell.margin_right = Emu(30000)
    cell.margin_top = Emu(18000)
    cell.margin_bottom = Emu(18000)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.text_frame.word_wrap = True
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.line_spacing = 1.0
    run = paragraph.add_run()
    run.text = str(text or "")
    run.font.name = BODY_FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def style_table_shape(table_shape, header_fill: Optional[RGBColor] = None) -> None:
    if header_fill is None:
        header_fill = BRAND_BLUE
    table = table_shape.table
    for row_idx, row in enumerate(table.rows):
        for cell in row.cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if row_idx == 0 else RGBColor(0xFF, 0xFF, 0xFF)


def render_slide2_mini_table(slide, slide_data: dict, layout: dict) -> dict:
    body = slide_data.get("body_copy") or {}
    headers = [body.get("table_header_1", ""), body.get("table_header_2", "")]
    rows = [body.get(f"table_row_{idx}", "") for idx in range(1, 4)]
    if not any(str(item).strip() for item in headers + rows):
        return {"rendered": False, "reason": "missing slide 2 table body_copy fields"}

    table_box = layout["table_box"]
    removed = remove_text_shapes_in_box(slide, table_box)
    left, top, width, height = table_box
    table_shape = slide.shapes.add_table(4, 2, Emu(left), Emu(top), Emu(width), Emu(height))
    table = table_shape.table
    table.columns[0].width = Emu(int(width * 0.42))
    table.columns[1].width = Emu(int(width * 0.58))
    for idx, header in enumerate(headers):
        set_cell_text(table.cell(0, idx), header, TABLE_HEADER_FONT_SIZE, True, RGBColor(0xFF, 0xFF, 0xFF))
    for row_idx, row_text in enumerate(rows, start=1):
        for col_idx, cell_text in enumerate(split_table_cells(row_text, 2)):
            set_cell_text(table.cell(row_idx, col_idx), cell_text, TABLE_BODY_FONT_SIZE, False)
    style_table_shape(table_shape)
    return {"rendered": True, "rows": 4, "columns": 2, "removed_text_shapes": removed}


def render_slide2_chart_plus_table(slide, slide_data: dict, layout: dict) -> dict:
    chart_result = build_chart(slide, slide_data, layout)
    table_result = render_slide2_mini_table(slide, slide_data, layout)
    rendered = chart_result.get("rendered", False) and table_result.get("rendered", False)
    return {
        "rendered": rendered,
        "chart": chart_result,
        "table": table_result,
        "chart_title": chart_result.get("chart_title", ""),
        "chart_type": chart_result.get("chart_type", ""),
        "table_rows": table_result.get("rows", 0),
        "table_columns": table_result.get("columns", 0),
    }


def render_slide6_compare_table(slide, slide_data: dict, layout: dict) -> dict:
    header_cells, rows = normalize_compare_table_payload(slide_data)
    col_count = len(header_cells)
    slide_no = slide_data.get("slide_no")
    if col_count < 1 or col_count > 8:
        return {
            "rendered": False,
            "reason": f"slide {slide_no} compare table requires 1-8 explicit header columns; found {col_count}",
        }

    if not header_cells and not rows:
        return {"rendered": False, "reason": f"missing slide {slide_no} compare_table_data payload"}
    if len(rows) < 1:
        return {"rendered": False, "reason": f"slide {slide_no} compare table needs at least 1 populated row"}

    normalized_rows = []
    for row_idx, cells in enumerate(rows, start=1):
        cells = [str(cell or "").strip() for cell in cells]
        if len(cells) > col_count:
            cells = cells[: col_count - 1] + [" / ".join(cells[col_count - 1 :])]
        if len(cells) < col_count:
            cells = cells + [""] * (col_count - len(cells))
        normalized_rows.append(cells)

    table_box = layout["table_box"]
    removed = remove_text_shapes_in_box(slide, table_box)
    left, top, width, height = table_box
    table_shape = slide.shapes.add_table(len(normalized_rows) + 1, col_count, Emu(left), Emu(top), Emu(width), Emu(height))
    table = table_shape.table
    col_weights = layout.get("column_weight_overrides", {}).get(str(col_count)) if isinstance(layout.get("column_weight_overrides"), dict) else None
    if not col_weights:
        col_weights = estimate_column_weights(header_cells, normalized_rows)
    total = sum(col_weights)
    for idx, weight in enumerate(col_weights):
        table.columns[idx].width = Emu(int(width * weight / total))
    header_font = TABLE_HEADER_FONT_SIZE
    body_font = TABLE_BODY_FONT_SIZE
    for idx, header in enumerate(header_cells):
        set_cell_text(table.cell(0, idx), header, header_font, True, RGBColor(0xFF, 0xFF, 0xFF))
    for row_idx, row_cells in enumerate(normalized_rows, start=1):
        for col_idx, cell_text in enumerate(row_cells):
            set_cell_text(table.cell(row_idx, col_idx), cell_text, body_font, False)
    style_table_shape(table_shape)
    return {"rendered": True, "rows": len(normalized_rows) + 1, "columns": col_count, "removed_text_shapes": removed}


def generic_visual_layout(prs: Presentation, slide_data: dict) -> dict:
    width = int(prs.slide_width)
    height = int(prs.slide_height)
    margin_x = int(width * 0.065)
    title_top = int(height * 0.055)
    title_height = int(height * 0.10)
    message_top = title_top + title_height
    message_height = int(height * 0.065)
    content_top = message_top + message_height + int(height * 0.035)
    footer_height = int(height * 0.045)
    content_height = height - content_top - footer_height - int(height * 0.045)
    content_width = width - margin_x * 2
    has_body = bool(slide_data.get("body_copy"))
    if has_body and (slide_data.get("chart_data") or slide_data.get("compare_table_data")):
        visual_left = margin_x + int(content_width * 0.40)
        visual_width = int(content_width * 0.60)
        body_width = int(content_width * 0.36)
    else:
        visual_left = margin_x
        visual_width = content_width
        body_width = content_width
    return {
        "title_box": (margin_x, content_top - int(height * 0.045), content_width, int(height * 0.035)),
        "chart_box": (visual_left, content_top, visual_width, content_height),
        "table_box": (visual_left, content_top, visual_width, content_height),
        "body_box": (margin_x, content_top, body_width, content_height),
        "footer_box": (margin_x, height - footer_height - int(height * 0.02), content_width, footer_height),
        "preserve_existing_shapes": True,
    }


def clear_slides(prs: Presentation) -> None:
    slide_id_list = prs.slides._sldIdLst  # python-pptx internal API; no public delete API exists.
    for slide_id in list(slide_id_list):
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        slide_id_list.remove(slide_id)


def text_units(text: str) -> float:
    units = 0.0
    for char in str(text or ""):
        if "\u4e00" <= char <= "\u9fff":
            units += 1.0
        elif char.isspace():
            units += 0.3
        else:
            units += 0.55
    return units


def fit_font_size(text: str, base: float, minimum: float, max_units: float) -> float:
    units = text_units(text)
    if max_units <= 0 or units <= max_units:
        return base
    return max(minimum, round(base * max_units / units, 1))


def add_text_box(slide, box: tuple[int, int, int, int], text: str, *, font_size: float, bold: bool = False, color: Optional[RGBColor] = None, fill: Optional[RGBColor] = None) -> None:
    if not str(text or "").strip():
        return
    left, top, width, height = box
    shape = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = Emu(45000)
    text_frame.margin_right = Emu(45000)
    text_frame.margin_top = Emu(25000)
    text_frame.margin_bottom = Emu(20000)
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = str(text or "").strip()
    run.font.name = BODY_FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or TEXT_GRAY


def add_card(slide, box: tuple[int, int, int, int], text: str, *, idx: int) -> None:
    left, top, width, height = box
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFC)
    shape.line.color.rgb = GRID_GRAY
    prefix = f"{idx}. " if idx else ""
    add_text_box(slide, box, prefix + str(text or "").strip(), font_size=fit_font_size(text, 10.5, 7.5, 80), color=TEXT_GRAY)


def body_items(slide_data: dict) -> list[str]:
    body = slide_data.get("body_copy") if isinstance(slide_data.get("body_copy"), dict) else {}
    values = [str(value or "").strip() for value in body.values() if str(value or "").strip()]
    if values:
        return values
    blocks = slide_data.get("body_blocks") if isinstance(slide_data.get("body_blocks"), list) else []
    return [str(block.get("copy") or "").strip() for block in blocks if isinstance(block, dict) and str(block.get("copy") or "").strip()]


def render_body_cards(slide, slide_data: dict, box: tuple[int, int, int, int]) -> dict:
    items = body_items(slide_data)
    if not items:
        return {"rendered": False, "reason": "no body copy"}
    left, top, width, height = box
    count = min(len(items), 8)
    cols = 1 if width < 4_000_000 or count <= 3 else 2
    rows = (count + cols - 1) // cols
    gap = 90_000
    card_w = int((width - gap * (cols - 1)) / cols)
    card_h = int((height - gap * (rows - 1)) / rows)
    for idx, item in enumerate(items[:count], start=1):
        row = (idx - 1) // cols
        col = (idx - 1) % cols
        add_card(slide, (left + col * (card_w + gap), top + row * (card_h + gap), card_w, card_h), item, idx=idx)
    return {"rendered": True, "items": count}


def render_style_guided_deck(template_ppt: Path, renderer_spec_path: Path, output_ppt: Path, template_profile_path: Path) -> dict:
    template_profile_warnings: list[str] = []
    _apply_template_profile_style(template_profile_path, template_profile_warnings)
    renderer_spec = load_json(renderer_spec_path)
    prs = Presentation(str(template_ppt))
    clear_slides(prs)
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    visual_results = []
    for slide_data in renderer_spec.get("slides", []):
        slide = prs.slides.add_slide(blank_layout)
        width = int(prs.slide_width)
        height = int(prs.slide_height)
        margin_x = int(width * 0.065)
        add_text_box(
            slide,
            (margin_x, int(height * 0.052), int(width * 0.87), int(height * 0.095)),
            str(slide_data.get("headline") or ""),
            font_size=fit_font_size(str(slide_data.get("headline") or ""), 24, 15, 34),
            bold=True,
            color=BRAND_BLUE,
        )
        add_text_box(
            slide,
            (margin_x, int(height * 0.145), int(width * 0.87), int(height * 0.075)),
            str(slide_data.get("main_message") or ""),
            font_size=fit_font_size(str(slide_data.get("main_message") or ""), 12.5, 8.5, 88),
            color=TEXT_GRAY,
        )
        layout = generic_visual_layout(prs, slide_data)
        body_result = render_body_cards(slide, slide_data, layout["body_box"])
        chart_result = {"rendered": False, "required_render": False, "reason": "no visual payload"}
        if slide_data.get("compare_table_data"):
            chart_result = render_slide6_compare_table(slide, slide_data, layout)
        elif slide_data.get("chart_data"):
            chart_result = build_chart(slide, slide_data, layout)
        add_text_box(
            slide,
            layout["footer_box"],
            str(slide_data.get("source_note") or ""),
            font_size=7.0,
            color=TEXT_GRAY,
        )
        visual_results.append(
            {
                "slide_no": slide_data.get("slide_no"),
                "selected_page_type": slide_data.get("selected_page_type"),
                "style_guided": True,
                "body": body_result,
                "visual": chart_result,
                "rendered": bool(body_result.get("rendered") or chart_result.get("rendered")),
                "required_render": False,
            }
        )
    save_presentation(prs, output_ppt)
    return {
        "input_ppt": str(template_ppt),
        "renderer_spec": str(renderer_spec_path),
        "template_profile": str(template_profile_path),
        "template_profile_warnings": template_profile_warnings,
        "output_ppt": str(output_ppt),
        "style_guided_render": True,
        "chart_rendering": visual_results,
    }


def render_quant_slide(prs: Presentation, renderer_spec: dict, slide_no: int, render_layouts: dict[int, dict]) -> dict:
    slide_data = find_slide_data(renderer_spec, slide_no)
    if not slide_data:
        return {"slide_no": slide_no, "rendered": False, "reason": f"slide {slide_no} not found in renderer_spec"}

    page_type = slide_data.get("selected_page_type")
    slide_layouts = render_layouts.get(slide_no, {})
    layout = slide_layouts.get(page_type)
    if not layout:
        return {"slide_no": slide_no, "rendered": False, "reason": f"unsupported page type: {page_type}"}

    if len(prs.slides) < slide_no:
        return {"slide_no": slide_no, "rendered": False, "reason": "clean deck has fewer slides than expected"}

    slide = prs.slides[slide_no - 1]
    try:
        if slide_no == 1 and page_type == "industry_overview_dynamic_page":
            result = render_slide1_dynamic_overview(slide, slide_data, layout)
        elif page_type == "matrix_page":
            result = render_matrix_slide(slide, slide_data, layout)
        elif slide_no == 2 and page_type == "chart_plus_mini_table_page":
            result = render_slide2_chart_plus_table(slide, slide_data, layout)
        elif slide_data.get("compare_table_data") and "table_box" in layout:
            result = render_slide6_compare_table(slide, slide_data, layout)
        elif slide_no == 1:
            result = render_slide1_visual(slide, slide_data, layout)
        else:
            result = build_chart(slide, slide_data, layout)
    except Exception as exc:
        return {
            "slide_no": slide_no,
            "selected_page_type": page_type,
            "rendered": False,
            "reason": f"{type(exc).__name__}: {exc}",
            "repair_hint": "Check renderer_spec chart/table payload for this slide. Expected canonical chart series or compare_table_data rows.",
        }
    result["slide_no"] = slide_no
    result["selected_page_type"] = page_type
    return result


def has_postprocess_renderer(slide_data: dict, render_layouts: dict[int, dict]) -> bool:
    slide_no = int(slide_data["slide_no"])
    page_type = slide_data.get("selected_page_type")
    return page_type in render_layouts.get(slide_no, {})


def skipped_non_rendered_slide(slide_data: dict) -> dict:
    return {
        "slide_no": int(slide_data["slide_no"]),
        "selected_page_type": slide_data.get("selected_page_type"),
        "rendered": False,
        "required_render": False,
        "skipped": True,
        "reason": "no deterministic postprocess renderer for this page type; text/table layout remains authoritative",
    }


def save_presentation(prs: Presentation, output_path: Path) -> None:
    if output_path.exists():
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = Path(tmp.name)
        try:
            prs.save(tmp_path)
            shutil.move(str(tmp_path), output_path)
        finally:
            if tmp_path.exists():
                os.unlink(tmp_path)
        return
    prs.save(output_path)


def _remove_ole_graphic_frames(xml_text: str) -> tuple[str, int]:
    removed_count = 0
    frame_pattern = re.compile(r"<p:graphicFrame\b.*?</p:graphicFrame>", re.DOTALL)

    def replace_frame(match: re.Match) -> str:
        nonlocal removed_count
        frame = match.group(0)
        lowered = frame.lower()
        if "<p:oleobj" in lowered or "tclayout.activedocument" in lowered or "think-cell" in lowered:
            removed_count += 1
            return ""
        return frame

    cleaned = frame_pattern.sub(replace_frame, xml_text)
    ole_obj_pattern = re.compile(r"<p:oleObj\b.*?</p:oleObj>", re.DOTALL)
    cleaned, fallback_removed = ole_obj_pattern.subn("", cleaned)
    removed_count += fallback_removed
    return cleaned, removed_count


def _remove_ole_relationships(rels_text: str) -> tuple[str, int]:
    rel_pattern = re.compile(r"<Relationship\b[^>]*/>", re.DOTALL)
    removed_count = 0

    def replace_rel(match: re.Match) -> str:
        nonlocal removed_count
        rel = match.group(0)
        lowered = rel.lower()
        if "/oleobject" in lowered or "target=\"../embeddings/oleobject" in lowered:
            removed_count += 1
            return ""
        return rel

    return rel_pattern.sub(replace_rel, rels_text), removed_count


def _remove_ole_content_types(content_types_text: str) -> tuple[str, int]:
    override_pattern = re.compile(r"<Override\b[^>]*PartName=\"/ppt/embeddings/oleObject[^\"]*\"[^>]*/>", re.DOTALL)
    cleaned, removed_count = override_pattern.subn("", content_types_text)
    return cleaned, removed_count


def sanitize_ole_artifacts(pptx_path: Path) -> dict:
    with ZipFile(pptx_path, "r") as archive:
        original_items = [(name, archive.read(name)) for name in archive.namelist()]

    removed_parts = []
    updated_parts = []
    sanitized_items = []

    for name, payload in original_items:
        lower_name = name.lower()
        if lower_name.startswith("ppt/embeddings/oleobject") or (
            lower_name.startswith("ppt/embeddings/") and lower_name.endswith(".bin")
        ):
            removed_parts.append(name)
            continue

        if lower_name.endswith(".rels") or lower_name.endswith(".xml"):
            text = payload.decode("utf-8", errors="ignore")
            updated_text = text
            removed_count = 0

            if lower_name.endswith(".rels"):
                updated_text, removed_count = _remove_ole_relationships(updated_text)
            elif lower_name == "[content_types].xml":
                updated_text, removed_count = _remove_ole_content_types(updated_text)
            elif lower_name.startswith("ppt/slides/") or lower_name.startswith("ppt/slidelayouts/") or lower_name.startswith("ppt/slidemasters/"):
                updated_text, removed_count = _remove_ole_graphic_frames(updated_text)

            if "think-cell" in updated_text.lower():
                updated_text = re.sub(r"think-cell[^<]*", "removed embedded object", updated_text, flags=re.IGNORECASE)
                removed_count += 1
            if "tclayout.activedocument" in updated_text.lower():
                updated_text = re.sub(r"TCLayout\.ActiveDocument\.\d+", "removed.embedded.object", updated_text, flags=re.IGNORECASE)
                removed_count += 1

            if removed_count:
                updated_parts.append({"part": name, "removed_count": removed_count})
                payload = updated_text.encode("utf-8")

        sanitized_items.append((name, payload))

    if not removed_parts and not updated_parts:
        return {"removed_parts": [], "updated_parts": []}

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp_path = Path(tmp.name)
    try:
        with ZipFile(tmp_path, "w", ZIP_DEFLATED) as archive:
            for name, payload in sanitized_items:
                archive.writestr(name, payload)
        shutil.move(str(tmp_path), pptx_path)
    finally:
        if tmp_path.exists():
            os.unlink(tmp_path)

    return {"removed_parts": removed_parts, "updated_parts": updated_parts}


def postprocess(
    input_ppt: Path,
    renderer_spec_path: Path,
    output_ppt: Path,
    render_layouts_path: Path = DEFAULT_RENDER_LAYOUTS_PATH,
    template_profile_path: Optional[Path] = None,
) -> dict:
    template_profile_warnings: list[str] = []
    if template_profile_path is None:
        template_profile_path = DEFAULT_TEMPLATE_PROFILE_PATH
    _apply_template_profile_style(template_profile_path, template_profile_warnings)

    renderer_spec = load_json(renderer_spec_path)
    render_layouts = load_render_layouts(render_layouts_path)
    prs = Presentation(str(input_ppt))

    removed_labels = remove_scaffold_labels(prs)
    removed_think_cell_ole = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        removed = remove_think_cell_ole_shapes(slide)
        if removed:
            removed_think_cell_ole.append({"slide_no": slide_idx, "removed": removed})
    page_number_updates = rewrite_page_numbers(prs)
    chart_results = []
    for slide_data in renderer_spec.get("slides", []):
        if slide_data.get("chart_data") or slide_data.get("compare_table_data") or (
            int(slide_data.get("slide_no", 0)) in {2, 6}
            and has_postprocess_renderer(slide_data, render_layouts)
        ):
            if has_postprocess_renderer(slide_data, render_layouts):
                chart_results.append(render_quant_slide(prs, renderer_spec, int(slide_data["slide_no"]), render_layouts))
            else:
                chart_results.append(skipped_non_rendered_slide(slide_data))

    save_presentation(prs, output_ppt)
    ole_sanitization = sanitize_ole_artifacts(output_ppt)

    return {
        "input_ppt": str(input_ppt),
        "renderer_spec": str(renderer_spec_path),
        "render_layouts": str(render_layouts_path),
        "template_profile": str(template_profile_path),
        "template_profile_warnings": template_profile_warnings,
        "output_ppt": str(output_ppt),
        "removed_scaffold_labels": removed_labels,
        "removed_think_cell_ole": removed_think_cell_ole,
        "ole_sanitization": ole_sanitization,
        "page_number_updates": page_number_updates,
        "chart_rendering": chart_results,
    }


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Post-process filled PPT with object-level visual cleanup and chart rendering."
    )
    parser.add_argument("--input-ppt", required=True, help="Path to cleaned PPTX input.")
    parser.add_argument("--renderer-spec", dest="renderer_spec", required=True, help="Path to renderer_spec.json.")
    parser.add_argument("--output", required=True, help="Path to write the post-processed PPTX.")
    parser.add_argument(
        "--render-layouts",
        default=str(DEFAULT_RENDER_LAYOUTS_PATH),
        help="Path to configs/render_layouts.json with deterministic renderer coordinates.",
    )
    parser.add_argument(
        "--template-profile",
        default=str(DEFAULT_TEMPLATE_PROFILE_PATH),
        help="Path to template_profile.json with visual style and fit rules.",
    )
    parser.add_argument("--log", help="Optional path to write a JSON log.")
    parser.add_argument(
        "--fail-on-unrendered",
        action="store_true",
        help="Exit non-zero if a required deterministic visual renderer fails.",
    )
    parser.add_argument(
        "--allow-ungated-debug",
        action="store_true",
        help="Bypass the pre-PPT stage gate. Use only for local diagnostics, never delivery.",
    )
    parser.add_argument(
        "--style-guided-render",
        action="store_true",
        help="Create an editable PPT from renderer_spec using the input PPT only as a style/size reference.",
    )
    args = parser.parse_args()

    output_path = Path(args.output)
    try:
        if args.allow_ungated_debug:
            require_debug_output_name(output_path)
        require_pre_ppt_gate(output_path.parent, allow_ungated_debug=args.allow_ungated_debug)
        if args.style_guided_render:
            result = render_style_guided_deck(
                Path(args.input_ppt),
                Path(args.renderer_spec),
                output_path,
                Path(args.template_profile),
            )
        else:
            result = postprocess(
                Path(args.input_ppt),
                Path(args.renderer_spec),
                output_path,
                Path(args.render_layouts),
                Path(args.template_profile),
            )
    except Exception as exc:
        result = {
            "input_ppt": args.input_ppt,
            "renderer_spec": args.renderer_spec,
            "render_layouts": args.render_layouts,
            "template_profile": args.template_profile,
            "output_ppt": args.output,
            "is_valid": False,
            "error": f"{type(exc).__name__}: {exc}",
            "repair_hint": "Run via scripts/pipeline.py render --run-dir <attempt_dir>, or pass absolute --input-ppt, --renderer-spec, --output, and --render-layouts paths for a single diagnostic step. Use configs/render_layouts.json, not template_registry.json.",
            "chart_rendering": [],
        }
        if args.log:
            save_json(result, Path(args.log))
        print(json.dumps(result, ensure_ascii=False, indent=2))
        raise SystemExit(1) from exc
    if args.log:
        save_json(result, Path(args.log))
    print(json.dumps(result, ensure_ascii=False, indent=2))
    if args.fail_on_unrendered:
        failed = [
            item
            for item in result["chart_rendering"]
            if not item.get("rendered") and item.get("required_render", True)
        ]
        if failed:
            raise SystemExit(1)


if __name__ == "__main__":
    main()
