#!/usr/bin/env python3
"""Analyze a PPT template and emit a deterministic template profile."""

from __future__ import annotations

import argparse
import json
import re
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from extract_template_registry import build_registry
from json_utils import load_json_file
from layout_config import layout_config_paths

try:
    from pptx import Presentation
except Exception:
    Presentation = None


ROOT = Path(__file__).resolve().parents[1]
TOKEN_PATTERN = re.compile(r"\{\{[^{}]+\}\}")

DEFAULT_STYLE = {
    "colors": {
        "brand_primary": "#0D57AA",
        "accent_red": "#C03C28",
        "grid_gray": "#D9D9D9",
        "text_gray": "#555555",
        "panel_fill": "#FAFBFC",
        "matrix_axis": "#FFFFFF",
    },
    "typography": {
        "body": "Microsoft YaHei",
        "table_header": "Microsoft YaHei",
        "table_body": "Microsoft YaHei",
        "legend_pt": 8.0,
        "table_header_pt": 10.0,
        "table_body_pt": 10.0,
    },
}


def _load_json(path: Path | str) -> dict[str, Any]:
    data = load_json_file(Path(path))
    if not isinstance(data, dict):
        raise ValueError(f"{path} must be a JSON object")
    return data


def _as_str(value: Any) -> str:
    text = str(value or "").strip()
    return text


def _as_color(value: Any, fallback: str) -> str:
    if value is None:
        return fallback
    if isinstance(value, str):
        raw = value.strip().lstrip("#")
        if len(raw) == 6 and all(ch in "0123456789ABCDEFabcdef" for ch in raw):
            return f"#{raw.upper()}"
    if isinstance(value, int):
        if 0 <= value <= 0xFFFFFF:
            return f"#{value:06X}"
    if isinstance(value, (list, tuple)) and len(value) >= 3:
        try:
            return f"#{int(value[0]):02X}{int(value[1]):02X}{int(value[2]):02X}"
        except Exception:
            return fallback
    if isinstance(value, dict):
        for key in ("hex", "color", "rgb", "value", "hex_value"):
            nested = value.get(key)
            converted = _as_color(nested, fallback)
            if converted != fallback:
                return converted
    return fallback


def _as_font_size(value: Any, fallback: float) -> float:
    try:
        value = float(value)
    except (TypeError, ValueError):
        return fallback
    if value <= 0:
        return fallback
    return value


def _most_common(values: list[Any], fallback: str) -> str:
    if not values:
        return fallback
    counts = Counter([str(v).strip() for v in values if str(v).strip()])
    if not counts:
        return fallback
    return counts.most_common(1)[0][0]


def _emu_to_inches(value: Any) -> float:
    try:
        return round(float(value) / 914400.0, 3)
    except Exception:
        return 0.0


def _shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    try:
        return str(shape.text or "").strip()
    except Exception:
        return ""


def _placeholder_type(shape: Any) -> str:
    if not getattr(shape, "is_placeholder", False):
        return ""
    try:
        return str(shape.placeholder_format.type)
    except Exception:
        return "unknown"


def _shape_role(shape: Any, slide_height: float) -> str:
    text = _shape_text(shape).lower()
    y = _emu_to_inches(getattr(shape, "top", 0))
    if "source" in text or "来源" in text or "资料来源" in text:
        return "source_footer"
    if slide_height and y >= slide_height * 0.82 and len(text) <= 180:
        return "footer_or_source_area"
    if getattr(shape, "has_chart", False):
        return "chart_slot"
    if getattr(shape, "has_table", False):
        return "table_slot"
    if getattr(shape, "shape_type", None) is not None and "PICTURE" in str(shape.shape_type):
        return "image_slot"
    if getattr(shape, "has_text_frame", False):
        return "text_slot"
    return "shape"


def _slot_capacity(shape: Any) -> dict[str, Any]:
    width = _emu_to_inches(getattr(shape, "width", 0))
    height = _emu_to_inches(getattr(shape, "height", 0))
    # Conservative text density approximation used only for fit planning.
    return {
        "width_in": width,
        "height_in": height,
        "area_in2": round(width * height, 3),
        "estimated_text_units": max(12, int(width * height * 18)) if width and height else 0,
    }


def _collect_template_inventory(template_path: Path) -> tuple[dict[str, Any], list[str]]:
    warnings: list[str] = []
    if Presentation is None:
        return {
            "slide_count": 0,
            "slides": [],
            "analysis_fallback": True,
        }, ["python-pptx unavailable; template inventory fell back to deterministic config only"]

    try:
        prs = Presentation(str(template_path))
    except Exception as exc:
        return {
            "slide_count": 0,
            "slides": [],
            "analysis_fallback": True,
            "analysis_error": str(exc),
        }, [f"could not open PPTX template for inventory: {exc}"]

    page_width = _emu_to_inches(prs.slide_width)
    page_height = _emu_to_inches(prs.slide_height)
    slides: list[dict[str, Any]] = []
    aggregate = Counter()
    for idx, slide in enumerate(prs.slides, start=1):
        shapes: list[dict[str, Any]] = []
        roles = Counter()
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            role = _shape_role(shape, page_height)
            roles[role] += 1
            aggregate[role] += 1
            text = _shape_text(shape)
            shapes.append(
                {
                    "shape_index": shape_idx,
                    "name": str(getattr(shape, "name", "") or ""),
                    "role": role,
                    "placeholder_type": _placeholder_type(shape),
                    "has_text": bool(text),
                    "text_sample": text[:160],
                    "token_count": len(TOKEN_PATTERN.findall(text)),
                    "geometry": {
                        "x_in": _emu_to_inches(getattr(shape, "left", 0)),
                        "y_in": _emu_to_inches(getattr(shape, "top", 0)),
                        **_slot_capacity(shape),
                    },
                }
            )
        slides.append(
            {
                "slide_no": idx,
                "shape_count": len(shapes),
                "roles": dict(sorted(roles.items())),
                "slots": shapes,
                "supports": {
                    "chart": roles.get("chart_slot", 0) > 0,
                    "table": roles.get("table_slot", 0) > 0,
                    "image": roles.get("image_slot", 0) > 0,
                    "source_footer": roles.get("source_footer", 0) > 0 or roles.get("footer_or_source_area", 0) > 0,
                    "text": roles.get("text_slot", 0) > 0,
                },
                "information_density": "high" if len(shapes) >= 18 else "medium" if len(shapes) >= 9 else "low",
            }
        )

    return {
        "slide_count": len(slides),
        "page_size": {"width_in": page_width, "height_in": page_height},
        "slides": slides,
        "aggregate_roles": dict(sorted(aggregate.items())),
        "analysis_fallback": False,
    }, warnings


def _normalize_render_layouts(raw: dict[str, Any]) -> dict[str, Any]:
    if isinstance(raw, dict) and isinstance(raw.get("slides"), dict):
        return raw
    if isinstance(raw, dict):
        return {"slides": raw}
    return {"slides": {}}


def _collect_template_style(template_path: Path) -> tuple[dict[str, Any], bool]:
    if Presentation is None:
        return {
            "colors": DEFAULT_STYLE["colors"],
            "typography": DEFAULT_STYLE["typography"],
            "analysis_fallback": True,
        }, True

    fill_colors: list[str] = []
    line_colors: list[str] = []
    text_colors: list[str] = []
    fonts_body: list[str] = []
    fonts_headers: list[str] = []
    fonts_table: list[str] = []
    body_sizes: list[float] = []
    header_sizes: list[float] = []
    table_sizes: list[float] = []
    fallback = False

    try:
        prs = Presentation(str(template_path))
    except Exception as exc:
        return {
            "colors": DEFAULT_STYLE["colors"],
            "typography": DEFAULT_STYLE["typography"],
            "analysis_fallback": True,
            "analysis_error": str(exc),
        }, True

    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                fill = getattr(shape, "fill", None)
                if fill is not None and getattr(fill, "fore_color", None) is not None and getattr(fill.fore_color, "type", None) is not None:
                    fill_colors.append(_as_color(fill.fore_color.rgb, DEFAULT_STYLE["colors"]["panel_fill"]))
            except Exception:
                pass
            try:
                line = getattr(shape, "line", None)
                if line is not None and getattr(line, "color", None) is not None and getattr(line.color, "type", None) is not None:
                    line_colors.append(_as_color(line.color.rgb, DEFAULT_STYLE["colors"]["grid_gray"]))
            except Exception:
                pass

            text_frame = getattr(shape, "text_frame", None)
            if not text_frame:
                continue
            for paragraph in text_frame.paragraphs:
                is_token = bool(TOKEN_PATTERN.search(paragraph.text or ""))
                for run in paragraph.runs:
                    run_font = run.font
                    if run_font is None:
                        continue
                    if run_font.name:
                        fonts_body.append(_as_str(run_font.name))
                        if is_token:
                            fonts_table.append(_as_str(run_font.name))
                    if run_font.size:
                        pt = _as_font_size(run_font.size.pt, DEFAULT_STYLE["typography"]["table_body_pt"])
                        body_sizes.append(pt)
                        if is_token and run_font.bold:
                            header_sizes.append(pt)
                            if run_font.name:
                                fonts_headers.append(_as_str(run_font.name))
                        elif is_token:
                            table_sizes.append(pt)
                    try:
                        if run_font.color is not None and getattr(run_font.color, "type", None) is not None:
                            text_colors.append(_as_color(run_font.color.rgb, DEFAULT_STYLE["colors"]["text_gray"]))
                    except Exception:
                        pass

                if is_token and paragraph.font is not None and paragraph.font.name:
                    fonts_headers.append(_as_str(paragraph.font.name))
                if paragraph.font is not None and paragraph.font.size is not None:
                    pt = _as_font_size(paragraph.font.size.pt, DEFAULT_STYLE["typography"]["table_body_pt"])
                    body_sizes.append(pt)

                try:
                    if paragraph.font is not None and paragraph.font.color is not None and getattr(paragraph.font.color, "type", None) is not None:
                        text_colors.append(_as_color(paragraph.font.color.rgb, DEFAULT_STYLE["colors"]["text_gray"]))
                except Exception:
                    pass

                if not is_token and paragraph.font is not None and paragraph.font.name:
                    fonts_body.append(_as_str(paragraph.font.name))

    if not (fill_colors or line_colors or text_colors or fonts_body or fonts_headers):
        fallback = True

    body_font = _most_common(fonts_body, DEFAULT_STYLE["typography"]["body"])
    header_font = _most_common(fonts_headers, body_font)
    table_font = _most_common(fonts_table, header_font or body_font)

    return {
        "colors": {
            "brand_primary": _most_common(fill_colors, DEFAULT_STYLE["colors"]["brand_primary"]),
            "accent_red": _most_common(
                [c for c in text_colors + line_colors if c != _most_common(fill_colors, DEFAULT_STYLE["colors"]["brand_primary"])],
                DEFAULT_STYLE["colors"]["accent_red"],
            ),
            "grid_gray": _most_common(line_colors, DEFAULT_STYLE["colors"]["grid_gray"]),
            "text_gray": _most_common(text_colors, DEFAULT_STYLE["colors"]["text_gray"]),
            "panel_fill": DEFAULT_STYLE["colors"]["panel_fill"],
            "matrix_axis": DEFAULT_STYLE["colors"]["matrix_axis"],
        },
        "typography": {
            "body": body_font,
            "table_header": header_font,
            "table_body": table_font,
            "legend_pt": _as_font_size(sum(body_sizes) / len(body_sizes), DEFAULT_STYLE["typography"]["legend_pt"]) if body_sizes else DEFAULT_STYLE["typography"]["legend_pt"],
            "table_header_pt": _as_font_size(sum(header_sizes) / len(header_sizes), DEFAULT_STYLE["typography"]["table_header_pt"]) if header_sizes else DEFAULT_STYLE["typography"]["table_header_pt"],
            "table_body_pt": _as_font_size(
                sum(table_sizes) / len(table_sizes),
                _as_font_size(sum(body_sizes) / len(body_sizes), DEFAULT_STYLE["typography"]["table_body_pt"]) if body_sizes else DEFAULT_STYLE["typography"]["table_body_pt"],
            ) if table_sizes else _as_font_size(
                sum(body_sizes) / len(body_sizes), DEFAULT_STYLE["typography"]["table_body_pt"]
            ) if body_sizes else DEFAULT_STYLE["typography"]["table_body_pt"],
        },
        "analysis_fallback": fallback,
    }, fallback


def _extract_source_footer_fields(template_registry: dict[str, Any]) -> list[str]:
    fields: list[str] = []
    for slide in template_registry.get("slides") or []:
        for variant in slide.get("variants") or []:
            for key in (variant.get("field_roles") or {}):
                if str(key).lower() == "source_footer":
                    fields.append("source_footer")
                    break
    return sorted(set(fields))


def _build_variants(template_registry: dict[str, Any], layout_budget: dict[str, Any], source_footer_fields: list[str]) -> list[dict[str, Any]]:
    budget_by_key = layout_budget.get("slide_budgets", {})
    if not isinstance(budget_by_key, dict):
        budget_by_key = {}

    variants: list[dict[str, Any]] = []
    for slide in template_registry.get("slides") or []:
        if not isinstance(slide, dict):
            continue
        slide_no = int(slide.get("slide_no") or 0)
        for variant in slide.get("variants") or []:
            if not isinstance(variant, dict):
                continue
            page_type = str(variant.get("page_type") or "").strip()
            if not page_type:
                continue
            supports = variant.get("supports") if isinstance(variant.get("supports"), dict) else {}
            field_roles = variant.get("field_roles") if isinstance(variant.get("field_roles"), dict) else {}
            budget = budget_by_key.get(f"{slide_no}:{page_type}", {})
            body_fields = [str(item) for item in (variant.get("required_body_fields") or []) if str(item)]
            variants.append(
                {
                    "slide_no": slide_no,
                    "page_type": page_type,
                    "render_layout": str(variant.get("render_layout_key") or page_type),
                    "supports": {
                        "chart": bool(supports.get("chart")),
                        "table": bool(supports.get("table")),
                        "matrix": bool(supports.get("matrix")),
                        "cards": bool(supports.get("cards")),
                    },
                    "required_body_fields": body_fields,
                    "field_roles": {str(k): str(v) for k, v in field_roles.items() if str(k)},
                    "capacity_rules": {
                        "body_fields_max_units": dict((budget or {}).get("body_fields_max_units", {})),
                        "layout_rules": {
                            "table": dict((budget or {}).get("table", {})),
                            "matrix": dict((budget or {}).get("matrix", {})),
                        },
                    },
                    "source_footer_required": bool("source_footer" in field_roles or "source_footer" in source_footer_fields),
                }
            )
    return variants


def _load_paths(args) -> tuple[dict[str, Path], Path]:
    layout_paths = layout_config_paths(args.layout_config)
    if args.slide_registry is not None:
        layout_paths["slide_registry"] = Path(args.slide_registry)
    if args.page_type_rules is not None:
        layout_paths["page_type_rules"] = Path(args.page_type_rules)
    if args.ppt_mapping is not None:
        layout_paths["ppt_mapping"] = Path(args.ppt_mapping)
    layout_paths["render_layouts"] = Path(args.render_layouts) if args.render_layouts else layout_paths["render_layouts"]
    layout_paths["text_fit_rules"] = Path(args.text_fit_rules) if args.text_fit_rules else layout_paths["text_fit_rules"]
    layout_paths["layout_budget"] = Path(args.layout_budget) if args.layout_budget else layout_paths["layout_budget"]
    template_path = Path(args.template)
    return layout_paths, template_path


def _build_profile(layout_paths: dict[str, Path], template_path: Path, output_path: Path, skip_pptextract: bool) -> dict[str, Any]:
    template_registry = build_registry(
        template=template_path,
        slide_registry_path=layout_paths["slide_registry"],
        page_type_rules_path=layout_paths["page_type_rules"],
        ppt_mapping_path=layout_paths["ppt_mapping"],
        layout_budget_path=layout_paths["layout_budget"],
        text_fit_rules_path=layout_paths["text_fit_rules"],
    )
    style_payload, style_fallback = _collect_template_style(template_path) if not skip_pptextract else (
        {"colors": DEFAULT_STYLE["colors"], "typography": DEFAULT_STYLE["typography"], "analysis_fallback": True},
        True,
    )

    if style_payload.get("analysis_fallback"):
        style_fallback = True
    inventory, inventory_warnings = _collect_template_inventory(template_path) if not skip_pptextract else (
        {"slide_count": 0, "slides": [], "analysis_fallback": True},
        ["template PPTX inventory skipped by --skip-pptextract"],
    )

    layout_budget = _load_json(layout_paths["layout_budget"])
    text_fit_rules = _load_json(layout_paths["text_fit_rules"])

    source_footer_fields = _extract_source_footer_fields(template_registry)
    render_layouts = _normalize_render_layouts(_load_json(layout_paths["render_layouts"]))
    variant_payload = _build_variants(template_registry, layout_budget, source_footer_fields)

    return {
        "schema_version": "template_profile_v1",
        "template_file": str(template_path),
        "analysis_source": "template_analyzer.py",
        "analysis_timestamp": datetime.now(timezone.utc).isoformat(),
        "analysis_fallback": bool(style_fallback),
        "analysis_errors": inventory_warnings,
        "render_layouts_source": str(layout_paths["render_layouts"]),
        "text_fit_rules_source": str(layout_paths["text_fit_rules"]),
        "layout_budget_source": str(layout_paths["layout_budget"]),
        "ppt_mapping_source": str(layout_paths["ppt_mapping"]),
        "visual_style": {
            "colors": style_payload["colors"],
            "typography": style_payload["typography"],
        },
        "text_geometry": {
            "analysis_fallback": bool(style_fallback),
            "font_samples": {
                "body": style_payload["typography"]["body"],
                "table_header": style_payload["typography"]["table_header"],
                "table_body": style_payload["typography"]["table_body"],
            },
        },
        "template_inventory": inventory,
        "dynamic_slots": {
            "slide_count": inventory.get("slide_count", 0),
            "slides": [
                {
                    "slide_no": item.get("slide_no"),
                    "supports": item.get("supports", {}),
                    "information_density": item.get("information_density", "unknown"),
                    "slot_count": item.get("shape_count", 0),
                }
                for item in inventory.get("slides", [])
                if isinstance(item, dict)
            ],
        },
        "layout": {
            "render_layouts": render_layouts,
            "text_fit_rules": {
                "fields": text_fit_rules.get("fields", {}),
                "renderer_field_aliases": text_fit_rules.get("renderer_field_aliases", {}),
            },
            "layout_budget": layout_budget,
        },
        "slide_variants": variant_payload,
        "source_policy": {
            "source_footer_fields": source_footer_fields or ["source_footer"],
            "required_source_footer": True if source_footer_fields else bool(
                any((item.get("supports") or {}).get("source_footer") for item in inventory.get("slides", []) if isinstance(item, dict))
            ),
        },
        "output_path": str(output_path),
        "template_registry_file": str(layout_paths["slide_registry"]),
    }


def _validate_profile(profile: dict[str, Any]) -> list[str]:
    errors = []
    if profile.get("schema_version") != "template_profile_v1":
        errors.append("schema_version must be template_profile_v1")
    if "visual_style" not in profile:
        errors.append("visual_style is required")
    if "layout" not in profile:
        errors.append("layout is required")
    return errors


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--template", default=str(ROOT / "assets" / "industry_section_template_master.pptx"))
    parser.add_argument("--layout-config", default=str(ROOT / "templates" / "layout_config.json"))
    parser.add_argument("--slide-registry", default=None)
    parser.add_argument("--page-type-rules", default=None)
    parser.add_argument("--ppt-mapping", default=None)
    parser.add_argument("--render-layouts", default=None)
    parser.add_argument("--text-fit-rules", default=None)
    parser.add_argument("--layout-budget", default=None)
    parser.add_argument("--output", default=str(ROOT / "templates" / "template_profile.json"))
    parser.add_argument("--skip-pptextract", action="store_true", help="Skip template parsing and rely on fallback values.")
    args = parser.parse_args()

    layout_paths, template_path = _load_paths(args)
    output = Path(args.output)
    output.parent.mkdir(parents=True, exist_ok=True)

    try:
        profile = _build_profile(layout_paths, template_path, output, args.skip_pptextract)
        profile["analysis_source"] = "template_analyzer.py" if not args.skip_pptextract else "template_analyzer.py (fallback)"
    except Exception as exc:
        inventory, inventory_warnings = _collect_template_inventory(template_path) if not args.skip_pptextract else (
            {"slide_count": 0, "slides": [], "analysis_fallback": True},
            ["template PPTX inventory skipped by --skip-pptextract"],
        )
        profile = {
            "schema_version": "template_profile_v1",
            "template_file": str(template_path),
            "analysis_source": "template_analyzer.py (fallback)",
            "analysis_timestamp": datetime.now(timezone.utc).isoformat(),
            "analysis_fallback": True,
            "analysis_errors": [str(exc)] + inventory_warnings,
            "render_layouts_source": str(layout_paths["render_layouts"]),
            "text_fit_rules_source": str(layout_paths["text_fit_rules"]),
            "layout_budget_source": str(layout_paths["layout_budget"]),
            "ppt_mapping_source": str(layout_paths["ppt_mapping"]),
            "visual_style": {
                "colors": DEFAULT_STYLE["colors"],
                "typography": DEFAULT_STYLE["typography"],
            },
            "text_geometry": {"analysis_fallback": True},
            "template_inventory": inventory,
            "dynamic_slots": {
                "slide_count": inventory.get("slide_count", 0),
                "slides": [
                    {
                        "slide_no": item.get("slide_no"),
                        "supports": item.get("supports", {}),
                        "information_density": item.get("information_density", "unknown"),
                        "slot_count": item.get("shape_count", 0),
                    }
                    for item in inventory.get("slides", [])
                    if isinstance(item, dict)
                ],
            },
            "layout": {
                "render_layouts": _normalize_render_layouts(_load_json(layout_paths["render_layouts"])),
                "text_fit_rules": {
                    "fields": _load_json(layout_paths["text_fit_rules"]).get("fields", {}),
                    "renderer_field_aliases": _load_json(layout_paths["text_fit_rules"]).get("renderer_field_aliases", {}),
                },
                "layout_budget": _load_json(layout_paths["layout_budget"]),
            },
            "slide_variants": [],
            "source_policy": {
                "source_footer_fields": ["source_footer"],
                "required_source_footer": False,
            },
            "output_path": str(output),
            "template_registry_file": str(layout_paths["slide_registry"]),
        }

    # Keep warnings in both schema and report channels explicit.
    errors = _validate_profile(profile)
    profile_warnings = profile.get("analysis_errors") or []
    output.write_text(json.dumps(profile, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    report = {
        "schema_version": "template_analyzer_v1",
        "is_valid": not errors,
        "error_count": len(errors),
        "warning_count": len(profile_warnings),
        "errors": errors,
        "warnings": profile_warnings,
        "output": str(output),
    }
    print(json.dumps(report, ensure_ascii=False, indent=2))
    return 0 if not errors else 1


if __name__ == "__main__":
    raise SystemExit(main())
