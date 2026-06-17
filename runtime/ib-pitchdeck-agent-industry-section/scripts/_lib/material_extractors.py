"""Material text extractors used by the Material Intake role."""

from __future__ import annotations

import csv
import html
import re
import urllib.error
import zipfile
import zlib
import xml.etree.ElementTree as ET
from html.parser import HTMLParser
from pathlib import Path
from urllib.request import Request, urlopen

from material_intake_common import clean_text_block, read_text_file


PDF_TEXT_RE = re.compile(rb"\((.*?)\)")
PDF_STREAM_TEXT_RE = re.compile(rb"BT(.*?)ET", re.S)
PDF_STREAM_FILTER_RE = re.compile(rb"/Filter\s*/([A-Za-z0-9]+)")
ASCII_HEX_RE = re.compile(r"#[0-9A-Fa-f]{2}")
PPT_TEXT_PATH_RE = re.compile(r"\[slide\s+(\d+)\]")


def _decode_hex(encoded: str) -> str:
    text = encoded.replace("\\", "")
    chunks = text.split("#")
    if len(chunks) == 1:
        return text
    decoded = [chunks[0]]
    for chunk in chunks[1:]:
        pair = chunk[:2]
        rest = chunk[2:]
        if len(pair) == 2 and re.fullmatch(r"[0-9A-Fa-f]{2}", pair):
            decoded.append(chr(int(pair, 16)))
            decoded.append(rest)
        else:
            decoded.append("#" + chunk)
    return "".join(decoded)


def _decode_pdf_text_bytes(raw: bytes) -> str:
    text_chunks: list[str] = []
    for segment in PDF_STREAM_TEXT_RE.finditer(raw):
        block = segment.group(0)
        data = block
        if b"/Filter" in data:
            filter_match = PDF_STREAM_FILTER_RE.search(data)
            if filter_match and filter_match.group(1).lower() == b"flatedecode":
                start = data.find(b"stream") + 6
                end = data.find(b"endstream")
                if start > 6 and end > start:
                    compressed = data[start:end].strip(b"\r\n")
                    try:
                        data = zlib.decompress(compressed)
                    except Exception:
                        pass
        for raw_text in PDF_TEXT_RE.finditer(data):
            raw_fragment = raw_text.group(1)
            try:
                text_chunks.append(raw_fragment.decode("latin1", errors="ignore"))
            except Exception:
                text_chunks.append(str(raw_fragment))
    if text_chunks:
        return clean_text_block("\n".join(text_chunks))
    fallback = raw.decode("latin1", errors="ignore")
    return clean_text_block(fallback)


def extract_pdf_text(pdf_path: str, output_path: str | None = None) -> tuple[str, list[str], int]:
    warnings: list[str] = []
    source = Path(pdf_path)
    if not source.exists():
        return "", [f"source path not found: {source}"], 1

    try:
        from PyPDF2 import PdfReader  # type: ignore

        reader = PdfReader(str(source))
        text = []
        for page in reader.pages:
            try:
                value = page.extract_text() or ""
            except Exception:
                continue
            if value:
                text.append(value)
        text_value = clean_text_block("\n".join(text))
        if not text_value:
            raise RuntimeError("PyPDF2 returned empty text")
        if output_path:
            Path(output_path).write_text(text_value, encoding="utf-8")
        return text_value, warnings, 0
    except Exception:
        pass

    text_value = _decode_pdf_text_bytes(source.read_bytes())
    if not text_value.strip():
        warnings.append("PDF text extraction fallback found no readable text stream; manual review required.")
        return "", warnings, 1
    if output_path:
        Path(output_path).write_text(text_value, encoding="utf-8")
    return text_value, warnings, 0


def _walk_shapes(shape, prefix: str, lines: list[str]) -> None:
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            _walk_shapes(child, prefix, lines)
        return
    if hasattr(shape, "text") and str(shape.text).strip():
        lines.append(f"{prefix} {str(shape.text).strip()}")
        return
    if hasattr(shape, "table"):
        table = shape.table
        for row in table.rows:
            row_values = [cell.text_frame.text.strip() if cell.text_frame else "" for cell in row.cells]
            filtered = [value for value in row_values if value]
            if filtered:
                lines.append(f"{prefix} {' | '.join(filtered)}")


def _extract_pptx_xml(xml_bytes: bytes) -> str:
    root = ET.fromstring(xml_bytes)
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    a_t_nodes = root.findall(".//a:t", ns)
    t_vals = [node.text.strip() for node in a_t_nodes if node.text and node.text.strip()]
    if t_vals:
        return " | ".join(t_vals)
    text_nodes = [node.text.strip() for node in root.iter() if node.text and node.text.strip()]
    return " | ".join(text_nodes)


def _extract_pptx_with_python_pptx(source: Path, output_path: str | None = None) -> tuple[str, list[str], int]:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:
        return "", [f"python-pptx unavailable ({exc})"], 1

    try:
        presentation = Presentation(str(source))
    except Exception as exc:
        return "", [f"cannot open PPTX: {exc}"], 1

    lines: list[str] = []
    for slide_idx, slide in enumerate(presentation.slides, start=1):
        prefix = f"[slide {slide_idx}]"
        for shape in slide.shapes:
            _walk_shapes(shape, prefix, lines)

    text_value = "\n".join([line.strip() for line in lines if line.strip()])
    if not text_value:
        return "", ["no readable text extracted from this pptx"], 1
    if output_path:
        Path(output_path).write_text(text_value, encoding="utf-8")
    return text_value, [], 0


def _extract_pptx_with_zip(source: Path, output_path: str | None = None) -> tuple[str, list[str], int]:
    lines: list[str] = []
    try:
        with zipfile.ZipFile(source, "r") as archive:
            for name in sorted(
                item
                for item in archive.namelist()
                if item.startswith("ppt/slides/slide") and item.endswith(".xml")
            ):
                m = PPT_TEXT_PATH_RE.search(name.split("/")[-1])
                slide_no = m.group(1) if m else ""
                prefix = f"[slide {slide_no}]" if slide_no else "[slide]"
                text = _extract_pptx_xml(archive.read(name))
                if text:
                    lines.append(f"{prefix} {text}")
    except Exception as exc:
        return "", [f"failed to read pptx package: {exc}"], 1

    text_value = "\n".join(line for line in lines if line.strip())
    if not text_value:
        return "", ["no readable text extracted from this pptx"], 1
    if output_path:
        Path(output_path).write_text(text_value, encoding="utf-8")
    return text_value, [], 0


def extract_pptx_text(pptx_path: str, output_path: str | None = None, material_id: str = "") -> tuple[str, list[str], int]:
    source = Path(pptx_path)
    if not source.exists():
        return "", [f"source path not found: {source}"], 1
    if source.suffix.lower() == ".ppt":
        return "", [".ppt binary text extraction not supported without external tools"], 1

    text_value, warnings, status = _extract_pptx_with_python_pptx(source, output_path)
    if status == 0:
        return text_value, warnings, status
    return _extract_pptx_with_zip(source, output_path)


def _read_csv_text(path: Path) -> str:
    lines: list[str] = []
    with path.open("r", encoding="utf-8", errors="replace") as fh:
        reader = csv.reader(fh)
        for row in reader:
            if row:
                line = " | ".join([cell.strip() for cell in row if str(cell).strip()])
                if line:
                    lines.append(line)
    return "\n".join(lines)


def _read_xlsx_text(path: Path) -> str:
    try:
        with zipfile.ZipFile(path, "r") as zf:
            names = zf.namelist()
            shared_strings = []
            if "xl/sharedStrings.xml" in names:
                shared_xml = ET.fromstring(zf.read("xl/sharedStrings.xml"))
                for item in shared_xml.findall(".//{http://schemas.openxmlformats.org/spreadsheetml/2006/main}t"):
                    shared_strings.append("".join(item.itertext()).strip())

            ns = {"ns": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
            sheets = [name for name in names if name.startswith("xl/worksheets/") and name.endswith(".xml")]
            rows: list[str] = []
            for sheet_name in sorted(sheets):
                sheet = ET.fromstring(zf.read(sheet_name))
                for row in sheet.findall(".//ns:sheetData/ns:row", ns):
                    values: list[str] = []
                    for cell in row.findall("ns:c", ns):
                        type_attr = cell.attrib.get("t", "")
                        if type_attr == "inlineStr":
                            inline_text = " ".join(
                                "".join(node.itertext()).strip()
                                for node in cell.findall(".//ns:is//ns:t", ns)
                                if "".join(node.itertext()).strip()
                            )
                            if inline_text:
                                values.append(inline_text.strip())
                            continue

                        value_node = cell.find("ns:v", ns)
                        if value_node is None:
                            continue
                        value = value_node.text or ""
                        if type_attr == "s":
                            try:
                                value = shared_strings[int(value)]
                            except Exception:
                                pass
                        values.append(str(value).strip())
                    line = " | ".join([item for item in values if item])
                    if line:
                        rows.append(line)
            return "\n".join(rows)
    except Exception:
        return ""


def extract_excel_text(excel_path: str, output_path: str | None = None, material_id: str = "") -> tuple[str, list[str], int]:
    source = Path(excel_path)
    if not source.exists():
        return "", [f"source path not found: {source}"], 1
    suffix = source.suffix.lower()
    if suffix == ".csv":
        text_value = _read_csv_text(source)
    elif suffix in {".xls", ".xlsx"}:
        if suffix == ".xls":
            return "", ["xls binary format is unsupported without external dependencies"], 1
        text_value = _read_xlsx_text(source)
    else:
        text_value = read_text_file(source)
        if not text_value.strip():
            return "", [f"unsupported file extension for structured extract: {suffix}"], 1

    if not text_value.strip():
        return "", ["no readable table rows extracted from file"], 1

    if output_path:
        Path(output_path).write_text(text_value, encoding="utf-8")
    return text_value, [], 0


class _TextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self._chunks: list[str] = []

    def handle_data(self, data: str) -> None:
        if data:
            self._chunks.append(data)

    def text(self) -> str:
        return clean_text_block(html.unescape(" ".join(self._chunks)))


def _strip_scripts(html_text: str) -> str:
    html_text = re.sub(r"<script[\s\S]*?</script>", " ", html_text, flags=re.IGNORECASE)
    html_text = re.sub(r"<style[\s\S]*?</style>", " ", html_text, flags=re.IGNORECASE)
    return html_text


def extract_web_url(url: str, output_path: str | None = None, material_id: str = "") -> tuple[str, list[str], int]:
    req = Request(url, headers={"User-Agent": "Mozilla/5.0 (compatible; ib-pitchdeck-agent/1.0)"})
    try:
        with urlopen(req, timeout=20) as response:
            raw = response.read()
            mime = response.headers.get_content_charset() or "utf-8"
    except urllib.error.URLError as exc:
        return "", [f"url fetch failed: {exc}"], 1
    except Exception as exc:
        return "", [f"unexpected fetch error: {exc}"], 1

    try:
        html_text = raw.decode(mime, errors="replace")
    except Exception:
        html_text = raw.decode("utf-8", errors="replace")
    if "<html" in html_text.lower():
        parser = _TextExtractor()
        parser.feed(_strip_scripts(html_text))
        parser.close()
        text_content = parser.text()
    else:
        text_content = html_text

    if not text_content.strip():
        return "", ["url content empty or non-text"], 1

    cleaned = clean_text_block(text_content)
    if output_path:
        Path(output_path).write_text(cleaned, encoding="utf-8")
    return cleaned, [], 0
