#!/usr/bin/env python3

import argparse
import json
import re
from pathlib import Path
from typing import Optional
from zipfile import ZipFile
import xml.etree.ElementTree as ET

from json_utils import load_json_file

try:
    from pptx import Presentation
except ImportError as exc:
    raise SystemExit(
        "python-pptx is required for validate_filled_ppt.py. "
        "Run this script with the project virtualenv created by ./setup.sh."
    ) from exc

TOKEN_PATTERN = re.compile(r"\{\{[^{}]+\}\}")
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
PKG_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
SLIDE_LAYOUT_LIBRARY_PATH = Path(__file__).resolve().parents[1] / "templates" / "slide_layout_library.json"

DISALLOWED_VISIBLE_LABELS = {
    "STANDARD",
    "PRIMARY CHART",
    "CHART / VISUAL",
    "MINI TABLE / SEGMENT CUT",
    "POINT 1",
    "POINT 2",
    "POINT 3",
    "DRIVER 1",
    "DRIVER 2",
    "DRIVER 3",
    "DRIVER 4",
    "DRIVER 5",
    "DRIVER 6",
    "TREND 1",
    "TREND 2",
    "TREND 3",
    "TREND 4",
    "TREND 5",
    "TREND 6",
    "BARRIER 1",
    "BARRIER 2",
    "BARRIER 3",
    "UPSTREAM",
    "MIDSTREAM",
    "DOWNSTREAM",
    "PROFIT POOL",
    "KEY BARRIERS",
    "KEY MESSAGES",
    "PEER COMPARE TABLE",
    "CRX / STRUCTURE",
    "COMPETITION DIMENSIONS",
    "TARGET RELATIVE POSITIONING",
    "TARGET POSITIONING",
    "PRIORITY TREND",
    "SECONDARY TREND",
    "WATCHLIST",
    "INDUSTRY ATTRACTIVENESS",
    "KEY INDUSTRY CHANGES BENEFITING TARGET",
    "OPEN DD QUESTIONS",
    "SUMMARY_PAGE",
    "CHART_PAGE",
    "CHART_PLUS_MINI_TABLE_PAGE",
    "DRIVER_CARD_PAGE",
    "VALUE_CHAIN_PAGE",
    "MOAT_PAGE",
    "COMPARE_TABLE_PAGE",
    "MATRIX_PAGE",
    "TREND_PAGE",
    "TIMELINE_PAGE",
    "THESIS_SUMMARY_PAGE",
    "DRIVER_CARDS_PAGE",
    "VALUE_CHAIN_ARCHITECTURE_PAGE",
    "MOAT_BARRIER_PAGE",
    "PRIORITY_TREND_PAGE",
    "TRANSACTION_IMPLICATION_PAGE",
}


def load_json(path: Path):
    try:
        return load_json_file(path)
    except FileNotFoundError as exc:
        raise FileNotFoundError(f"JSON file not found: {path}") from exc


def normalize_visible_text(text: str) -> str:
    return " ".join(text.strip().split())


def is_disallowed_visible_label(text: str) -> bool:
    normalized = normalize_visible_text(text)
    upper = normalized.upper()
    if upper.endswith("_PAGE"):
        return True
    if upper in DISALLOWED_VISIBLE_LABELS:
        return True
    if re.fullmatch(r"(SUMMARY|CHART|DRIVER_CARD|VALUE_CHAIN|MOAT|COMPARE_TABLE|MATRIX|TREND|TIMELINE)_PAGE", upper):
        return True
    return False


def collect_visible_text_issues(pptx_path: Path) -> list[dict]:
    prs = Presentation(str(pptx_path))
    issues = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            text = getattr(shape, "text", "") if hasattr(shape, "text") else ""
            if not text or not text.strip():
                continue
            normalized = normalize_visible_text(text)
            if is_disallowed_visible_label(normalized):
                issues.append(
                    {
                        "slide_no": slide_idx,
                        "shape_name": shape.name,
                        "text": normalized,
                    }
                )
    return issues


def is_footer_page_number_candidate(shape, slide_width: int, slide_height: int) -> bool:
    text = getattr(shape, "text", "") if hasattr(shape, "text") else ""
    normalized = normalize_visible_text(text)
    if not re.fullmatch(r"\d{1,3}", normalized):
        return False
    return shape.left >= slide_width * 0.55 and shape.top >= slide_height * 0.85


def collect_page_number_issues(pptx_path: Path) -> dict:
    prs = Presentation(str(pptx_path))
    page_numbers = []
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    for slide_idx, slide in enumerate(prs.slides, start=1):
        candidates = []
        for shape in slide.shapes:
            if is_footer_page_number_candidate(shape, slide_width, slide_height):
                candidates.append(normalize_visible_text(shape.text))
        page_numbers.append({"slide_no": slide_idx, "numbers": candidates})

    expected = [str(i) for i in range(1, len(prs.slides) + 1)]
    actual = [item["numbers"][0] if len(item["numbers"]) == 1 else "" for item in page_numbers]
    issues = []
    for idx, item in enumerate(page_numbers):
        expected_no = expected[idx]
        if item["numbers"] != [expected_no]:
            issues.append(
                {
                    "slide_no": item["slide_no"],
                    "expected": expected_no,
                    "found": item["numbers"],
                }
            )
    return {
        "expected": expected,
        "actual": actual,
        "issues": issues,
        "is_valid": not issues,
    }


def load_slide_layout_library(path: Path = SLIDE_LAYOUT_LIBRARY_PATH) -> dict[int, dict]:
    data = load_json(path)
    slides = data.get("slides")
    if not isinstance(slides, list):
        raise ValueError(f"Invalid slide layout library in {path}: missing list field 'slides'.")

    library = {}
    for item in slides:
        slide_no = item.get("slide_no")
        slide_key = item.get("slide_key")
        page_type_to_slide = item.get("page_type_to_slide")
        if not isinstance(slide_no, int) or not slide_key or not isinstance(page_type_to_slide, dict):
            raise ValueError(
                f"Invalid slide layout entry in {path}: "
                f"slide_no={slide_no}, slide_key={slide_key}, page_type_to_slide={page_type_to_slide}."
            )
        library[slide_no] = {
            "slide_key": slide_key,
            "page_type_to_slide": page_type_to_slide,
        }
    return library


def open_pptx(path: Path) -> ZipFile:
    try:
        return ZipFile(path, "r")
    except FileNotFoundError as exc:
        raise FileNotFoundError(f"PPTX file not found: {path}") from exc
    except Exception as exc:
        raise ValueError(f"Failed to open PPTX {path}: {exc}") from exc


def _extract_slides(data: dict, source_path: Optional[Path] = None) -> list[dict]:
    """Extract slide list from ppt_copy (ppt_copy_slides) or storyboard (slides)."""
    for key in ("ppt_copy_slides", "slides"):
        slides = data.get(key)
        if isinstance(slides, list) and slides:
            return slides
    raise ValueError(
        f"Cannot find slide data in {source_path or 'input'}: "
        f"expected 'ppt_copy_slides' or 'slides' key with a non-empty array."
    )


def expected_slide_selection(control_data: dict, control_file_path: Path) -> list[dict]:
    slide_layout_library = load_slide_layout_library()
    slides = _extract_slides(control_data, control_file_path)
    by_no = {int(page["slide_no"]): page for page in slides}
    expected = []
    for slide_no, config in slide_layout_library.items():
        page = by_no.get(slide_no)
        if not page:
            raise ValueError(
                f"Missing slide in control_file={control_file_path}: "
                f"slide_no={slide_no}, slide_key={config['slide_key']}."
            )
        selected_page_type = page.get("selected_page_type")
        physical_slide = config["page_type_to_slide"].get(selected_page_type)
        if not physical_slide:
            allowed = ", ".join(config["page_type_to_slide"].keys())
            raise ValueError(
                f"Invalid selected_page_type in control_file={control_file_path}: "
                f"slide_no={slide_no}, slide_key={config['slide_key']}, found='{selected_page_type}', allowed={allowed}."
            )
        expected.append(
            {
                "slide_no": slide_no,
                "slide_key": config["slide_key"],
                "selected_page_type": selected_page_type,
                "expected_physical_slide": physical_slide,
            }
        )
    return expected


def collect_remaining_placeholders(pptx_path: Path) -> list[dict]:
    remaining = []
    with open_pptx(pptx_path) as archive:
        for name in archive.namelist():
            if not (name.startswith("ppt/slides/slide") and name.endswith(".xml")):
                continue
            text = archive.read(name).decode("utf-8")
            for placeholder in sorted(set(TOKEN_PATTERN.findall(text))):
                remaining.append({"slide_xml": Path(name).name, "placeholder": placeholder})
    return remaining


def collect_referenced_slide_files(pptx_path: Path) -> list[str]:
    with open_pptx(pptx_path) as archive:
        try:
            presentation_xml = ET.fromstring(archive.read("ppt/presentation.xml"))
            rels_xml = ET.fromstring(archive.read("ppt/_rels/presentation.xml.rels"))
        except KeyError as exc:
            raise ValueError(f"Missing presentation relationship file in {pptx_path}: {exc}") from exc

        rel_targets = {
            rel.attrib["Id"]: Path(rel.attrib["Target"]).name
            for rel in rels_xml.findall(f"{{{PKG_NS}}}Relationship")
            if rel.attrib.get("Type", "").endswith("/slide")
        }

        sld_id_lst = presentation_xml.find(f"{{{P_NS}}}sldIdLst")
        if sld_id_lst is None:
            raise ValueError(f"presentation.xml is missing p:sldIdLst in {pptx_path}.")

        referenced = []
        for sld_id in list(sld_id_lst):
            rid = sld_id.attrib.get(f"{{{R_NS}}}id")
            referenced.append(rel_targets.get(rid, ""))
        return referenced


def collect_active_placeholders(ppt_mapping: dict, control_data: dict, control_file_path: Path) -> list[dict]:
    expected = {item["slide_no"]: item for item in expected_slide_selection(control_data, control_file_path)}
    active = []
    for mapping_slide in ppt_mapping.get("slides", []):
        slide_no = int(mapping_slide["slide_no"])
        slide_key = mapping_slide.get("slide_key", "")
        if "tokens" in mapping_slide:
            for token in mapping_slide.get("tokens", []):
                active.append(
                    {
                        "slide_no": slide_no,
                        "slide_key": slide_key,
                        "selected_page_type": mapping_slide.get("selected_page_type", ""),
                        "placeholder": token["placeholder"],
                        "field_name": token.get("field_name", ""),
                    }
                )
            continue

        selected_page_type = expected[slide_no]["selected_page_type"]
        variant = mapping_slide.get("controlled_variants", {}).get(selected_page_type)
        if variant is None:
            allowed = ", ".join(mapping_slide.get("controlled_variants", {}).keys())
            raise ValueError(
                f"Mapping mismatch in templates/ppt_mapping.json: slide_no={slide_no}, slide_key={slide_key}, "
                f"selected_page_type='{selected_page_type}', allowed={allowed}."
            )
        for token in variant.get("tokens", []):
            active.append(
                {
                    "slide_no": slide_no,
                    "slide_key": slide_key,
                    "selected_page_type": selected_page_type,
                    "placeholder": token["placeholder"],
                    "field_name": token.get("field_name", ""),
                }
            )
    return active


def collect_suspicious_missing_values(replacement_dict: dict, active_placeholders: list[dict]) -> list[dict]:
    suspicious = []
    for item in active_placeholders:
        if (
            item["slide_no"] == 1
            and item["selected_page_type"] == "summary_page"
            and item["field_name"] == "chart_title"
        ):
            continue
        value = replacement_dict.get(item["placeholder"])
        if value is None or str(value).strip() == "":
            suspicious.append(
                {
                    "slide_no": item["slide_no"],
                    "slide_key": item["slide_key"],
                    "selected_page_type": item["selected_page_type"],
                    "placeholder": item["placeholder"],
                    "field_name": item["field_name"],
                    "value_status": "missing_key" if value is None else "empty_string",
                }
            )
    return suspicious


def build_report(
    filled_ppt_path: Path,
    clean_ppt_path: Path,
    control_file_path: Path,
    replacement_dict_path: Path,
    ppt_mapping_path: Path,
) -> dict:
    control_file = load_json(control_file_path)
    replacement_dict = load_json(replacement_dict_path)
    ppt_mapping = load_json(ppt_mapping_path)

    expected = expected_slide_selection(control_file, control_file_path)
    expected_physical_slides = [item["expected_physical_slide"] for item in expected]
    actual_kept_slides = collect_referenced_slide_files(clean_ppt_path)
    remaining_placeholders = collect_remaining_placeholders(filled_ppt_path)
    active_placeholders = collect_active_placeholders(ppt_mapping, control_file, control_file_path)
    suspicious_missing_values = collect_suspicious_missing_values(replacement_dict, active_placeholders)
    visible_text_issues = collect_visible_text_issues(clean_ppt_path)
    page_number_check = collect_page_number_issues(clean_ppt_path)

    kept_slide_count_ok = len(actual_kept_slides) == len(expected_physical_slides) == 8
    renumbered_after_resave = actual_kept_slides == [f"slide{i}.xml" for i in range(1, len(actual_kept_slides) + 1)]
    kept_slide_selection_ok = actual_kept_slides == expected_physical_slides or (
        kept_slide_count_ok and renumbered_after_resave
    )
    placeholders_ok = not remaining_placeholders
    suspicious_values_ok = not suspicious_missing_values
    visible_text_ok = not visible_text_issues
    page_numbers_ok = page_number_check["is_valid"]

    return {
        "summary": {
            "filled_ppt": str(filled_ppt_path),
            "clean_ppt": str(clean_ppt_path),
            "control_file": str(control_file_path),
            "replacement_dict": str(replacement_dict_path),
            "ppt_mapping": str(ppt_mapping_path),
            "remaining_placeholder_count": len(remaining_placeholders),
            "expected_kept_slide_count": len(expected_physical_slides),
            "actual_kept_slide_count": len(actual_kept_slides),
            "suspicious_missing_active_value_count": len(suspicious_missing_values),
            "visible_scaffold_label_count": len(visible_text_issues),
            "page_number_issue_count": len(page_number_check["issues"]),
            "placeholders_ok": placeholders_ok,
            "kept_slide_count_ok": kept_slide_count_ok,
            "kept_slide_selection_ok": kept_slide_selection_ok,
            "kept_slide_files_renumbered_after_resave": renumbered_after_resave,
            "suspicious_values_ok": suspicious_values_ok,
            "visible_text_ok": visible_text_ok,
            "page_numbers_ok": page_numbers_ok,
            "is_valid": (
                placeholders_ok
                and kept_slide_count_ok
                and kept_slide_selection_ok
                and suspicious_values_ok
                and visible_text_ok
                and page_numbers_ok
            ),
        },
        "expected_kept_slides": expected,
        "actual_kept_physical_slides": actual_kept_slides,
        "remaining_placeholders_in_filled_ppt": remaining_placeholders,
        "suspicious_missing_active_values": suspicious_missing_values,
        "visible_scaffold_label_issues": visible_text_issues,
        "page_number_check": page_number_check,
    }


def main():
    parser = argparse.ArgumentParser(
        description="Validate filled and cleaned PPT outputs against placeholder coverage and slide control file."
    )
    parser.add_argument("--filled-ppt", default="industry_section_filled.pptx")
    parser.add_argument("--clean-ppt", default="industry_section_filled_clean.pptx")
    parser.add_argument("--control-file", default="industry_storyboard.json",
                        help="Path to slide control file: industry_storyboard.json or industry_section_ppt_copy.json.")
    parser.add_argument("--replacement-dict", default="replacement_dict.json")
    parser.add_argument("--ppt-mapping", default="templates/ppt_mapping.json")
    parser.add_argument("--output", default="filled_ppt_validation.json")
    parser.add_argument("--fail-on-issue", action="store_true")
    args = parser.parse_args()

    try:
        report = build_report(
            Path(args.filled_ppt),
            Path(args.clean_ppt),
            Path(args.control_file),
            Path(args.replacement_dict),
            Path(args.ppt_mapping),
        )
    except Exception as exc:
        raise SystemExit(str(exc)) from exc

    with Path(args.output).open("w", encoding="utf-8") as f:
        json.dump(report, f, ensure_ascii=False, indent=2)
        f.write("\n")

    if args.fail_on_issue and not report["summary"]["is_valid"]:
        raise SystemExit(1)


if __name__ == "__main__":
    main()
