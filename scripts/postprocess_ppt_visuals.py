#!/usr/bin/env python3

import argparse
import json
import os
import re
import shutil
import tempfile
from pathlib import Path
from typing import Optional, Union
from zipfile import ZIP_DEFLATED, ZipFile

from gate_guard import require_pre_ppt_gate
from json_utils import load_json_file

try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
except ImportError as exc:
    raise SystemExit(
        "python-pptx is required for postprocess_ppt_visuals.py. "
        "Run this script with a Python environment that has the `pptx` package installed, "
        "such as the project virtualenv created by `./setup.sh`."
    ) from exc


SCAFFOLD_LABELS = {
    "PRIMARY CHART",
    "CHART / VISUAL",
    "MINI TABLE / SEGMENT CUT",
    "POINT 1",
    "POINT 2",
    "POINT 3",
    "STANDARD",
    "SUMMARY_PAGE",
    "CHART_PAGE",
    "CHART_PLUS_MINI_TABLE_PAGE",
    "DRIVER_CARD_PAGE",
    "VALUE_CHAIN_PAGE",
    "MOAT_PAGE",
    "COMPARE_TABLE_PAGE",
    "MATRIX_PAGE",
    "TREND_PAGE",
    "TIMELINE_PAGE",
    "THESIS_SUMMARY_PAGE",
    "DRIVER_CARDS_PAGE",
    "VALUE_CHAIN_ARCHITECTURE_PAGE",
    "MOAT_BARRIER_PAGE",
    "PRIORITY_TREND_PAGE",
    "TRANSACTION_IMPLICATION_PAGE",
    "PEER COMPARE TABLE",
    "CRx / STRUCTURE",
    "COMPETITION DIMENSIONS",
    "TARGET RELATIVE POSITIONING",
    "TARGET POSITIONING",
    "PRIORITY TREND",
    "SECONDARY TREND",
    "WATCHLIST",
    "INDUSTRY ATTRACTIVENESS",
    "KEY INDUSTRY CHANGES BENEFITING TARGET",
    "OPEN DD QUESTIONS",
    "DRIVER 1",
    "DRIVER 2",
    "DRIVER 3",
    "DRIVER 4",
    "DRIVER 5",
    "DRIVER 6",
    "TREND 1",
    "TREND 2",
    "TREND 3",
    "TREND 4",
    "TREND 5",
    "TREND 6",
    "BARRIER 1",
    "BARRIER 2",
    "BARRIER 3",
    "UPSTREAM",
    "MIDSTREAM",
    "DOWNSTREAM",
    "PROFIT POOL",
    "KEY BARRIERS",
    "industry_overview",
    "market_size_segmentation",
    "key_industry_drivers",
    "value_chain_profit_pool",
    "key_barriers_value_drivers",
    "competitive_landscape",
    "industry_trends_future_evolution",
    "key_takeaways_for_target",
}

DEFAULT_RENDER_LAYOUTS_PATH = Path(__file__).resolve().parents[1] / "templates" / "render_layouts.json"

BRAND_BLUE = RGBColor(0x0D, 0x57, 0xAA)
GRID_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
TEXT_GRAY = RGBColor(0x55, 0x55, 0x55)
ACCENT_RED = RGBColor(0xC0, 0x3C, 0x28)
LEGEND_FONT_SIZE = 8
BODY_FONT = "Microsoft YaHei"
TABLE_HEADER_FONT_SIZE = 10.0
TABLE_BODY_FONT_SIZE = 10.0


def load_json(path: Path) -> dict:
    return load_json_file(path)


def load_render_layouts(path: Path = DEFAULT_RENDER_LAYOUTS_PATH) -> dict[int, dict]:
    data = load_json(path)
    slides = data.get("slides")
    if not isinstance(slides, dict):
        raise ValueError(f"Invalid render layout file {path}: missing object field 'slides'.")
    normalized: dict[int, dict] = {}
    for slide_no_raw, variants in slides.items():
        try:
            slide_no = int(slide_no_raw)
        except (TypeError, ValueError) as exc:
            raise ValueError(f"Invalid slide number in render layout file {path}: {slide_no_raw}") from exc
        if not isinstance(variants, dict):
            raise ValueError(f"Invalid render layout file {path}: slide {slide_no} variants must be an object.")
        normalized[slide_no] = {}
        for page_type, boxes in variants.items():
            if not isinstance(boxes, dict):
                raise ValueError(f"Invalid render layout file {path}: slide {slide_no}/{page_type} boxes must be an object.")
            normalized[slide_no][page_type] = {}
            for box_name, box in boxes.items():
                if isinstance(box, bool):
                    normalized[slide_no][page_type][box_name] = box
                    continue
                if not (
                    isinstance(box, list)
                    and len(box) == 4
                    and all(isinstance(value, int) for value in box)
                ):
                    raise ValueError(
                        f"Invalid render layout file {path}: slide {slide_no}/{page_type}/{box_name} "
                        "must be [left, top, width, height] integer EMUs."
                    )
                normalized[slide_no][page_type][box_name] = tuple(box)
    return normalized


def save_json(data: dict, path: Path) -> None:
    with path.open("w", encoding="utf-8") as handle:
        json.dump(data, handle, ensure_ascii=False, indent=2)


def clear_text(shape) -> None:
    if not hasattr(shape, "text_frame"):
        return
    text_frame = shape.text_frame
    text_frame.clear()


def remove_shape(shape) -> None:
    element = shape._element
    element.getparent().remove(element)


def intersects(shape, box: tuple[int, int, int, int]) -> bool:
    left, top, width, height = box
    right = left + width
    bottom = top + height
    shape_left = int(shape.left)
    shape_top = int(shape.top)
    shape_right = shape_left + int(shape.width)
    shape_bottom = shape_top + int(shape.height)
    return not (
        shape_right < left
        or shape_left > right
        or shape_bottom < top
        or shape_top > bottom
    )


def set_single_paragraph(shape, text: str) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.name = BODY_FONT


def remove_scaffold_labels(prs: Presentation) -> list[dict]:
    removed = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if text in SCAFFOLD_LABELS or text.upper().endswith("_PAGE"):
                clear_text(shape)
                removed.append({"slide_no": slide_idx, "label": text, "shape_name": shape.name})
    return removed


def is_footer_page_number_candidate(shape, slide_width: int, slide_height: int) -> bool:
    if not hasattr(shape, "text"):
        return False
    text = shape.text.strip()
    if not text.isdigit():
        return False
    return shape.left >= slide_width * 0.55 and shape.top >= slide_height * 0.85


def rewrite_page_numbers(prs: Presentation) -> list[dict]:
    updates = []
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if is_footer_page_number_candidate(shape, slide_width, slide_height):
                old = shape.text.strip()
                set_single_paragraph(shape, str(slide_idx))
                updates.append(
                    {
                        "slide_no": slide_idx,
                        "shape_name": shape.name,
                        "old": old,
                        "new": str(slide_idx),
                    }
                )
                break
    return updates


def find_slide_data(storyboard: dict, slide_no: int) -> Optional[dict]:
    for slide in storyboard.get("slides", []):
        if slide.get("slide_no") == slide_no:
            return slide
    return None


def apply_chart_title(slide, text: str, layout: dict) -> bool:
    if not text:
        return False
    left, top, width, height = layout["title_box"]
    candidates = []
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
        if abs(shape.left - left) < 20000 and abs(shape.top - top) < 20000:
            candidates.append(shape)
    if not candidates:
        return False
    target = sorted(candidates, key=lambda shp: (abs(shp.width - width), abs(shp.height - height)))[0]
    set_single_paragraph(target, text)
    paragraph = target.text_frame.paragraphs[0]
    if paragraph.runs:
        font = paragraph.runs[0].font
        font.name = BODY_FONT
        font.size = Pt(10)
        font.color.rgb = TEXT_GRAY
        font.bold = False
    return True


def chart_type_for(chart_data: dict):
    chart_type = str(chart_data.get("chart_type") or "bar").lower()
    if chart_type in {"bar", "column", "clustered_bar", "clustered_column"}:
        return XL_CHART_TYPE.COLUMN_CLUSTERED
    if chart_type in {"stacked_bar", "stacked_column"}:
        return XL_CHART_TYPE.COLUMN_STACKED
    if chart_type in {"line", "line_chart"}:
        return XL_CHART_TYPE.LINE_MARKERS
    return None


def short_label(value: str, max_units: float = 14.0) -> str:
    text = str(value or "").strip()
    units = 0.0
    result = []
    for ch in text:
        if ch.isspace():
            char_units = 0.3
        elif ord(ch) < 128:
            char_units = 0.55
        else:
            char_units = 1.0
        if units + char_units > max_units:
            return "".join(result).rstrip() + "..."
        result.append(ch)
        units += char_units
    return text


def format_chart_legend(chart) -> None:
    if not chart.has_legend:
        return
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(LEGEND_FONT_SIZE)
    chart.legend.font.name = BODY_FONT
    chart.legend.font.color.rgb = TEXT_GRAY
    chart.legend.font.bold = False


def build_chart(slide, slide_data: dict, layout: dict) -> dict:
    chart_data = slide_data.get("chart_data") or {}
    series = chart_data.get("series") or []
    categories = chart_data.get("categories") or []
    if not series or not categories:
        return {"rendered": False, "reason": "missing chart_data series/categories"}

    ppt_chart_type = chart_type_for(chart_data)
    if ppt_chart_type is None:
        return {
            "rendered": False,
            "reason": f"unsupported chart_type for deterministic chart renderer: {chart_data.get('chart_type')}",
        }

    chart_title = chart_data.get("title") or ""
    apply_chart_title(slide, chart_title, layout)

    chart_payload = CategoryChartData()
    chart_payload.categories = categories
    for chart_series in series:
        values = chart_series.get("values") or []
        if len(values) != len(categories):
            return {"rendered": False, "reason": "series/category length mismatch"}
        chart_payload.add_series(short_label(chart_series.get("name") or ""), values)

    chart_box = layout["chart_box"]
    removed = []
    if not layout.get("preserve_existing_shapes"):
        removed = remove_text_shapes_in_box(slide, chart_box)
    left, top, width, height = chart_box
    if len(series) > 1:
        height = int(height * 0.88)
    graphic_frame = slide.shapes.add_chart(
        ppt_chart_type,
        Emu(left),
        Emu(top),
        Emu(width),
        Emu(height),
        chart_payload,
    )
    chart = graphic_frame.chart
    chart.has_title = False
    chart.has_legend = len(series) > 1
    format_chart_legend(chart)

    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(9)
    category_axis.tick_labels.font.name = BODY_FONT
    category_axis.tick_labels.font.color.rgb = TEXT_GRAY

    value_axis = chart.value_axis
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = GRID_GRAY
    value_axis.tick_labels.font.size = Pt(9)
    value_axis.tick_labels.font.name = BODY_FONT
    value_axis.tick_labels.font.color.rgb = TEXT_GRAY
    value_axis.format.line.color.rgb = GRID_GRAY

    normalized_chart_type = str(chart_data.get("chart_type") or "").lower()
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    if normalized_chart_type in {"stacked_bar", "stacked_column"}:
        data_labels.position = XL_DATA_LABEL_POSITION.CENTER
        data_labels.number_format = '0.0'
    else:
        data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
        data_labels.number_format = '#,##0'
    data_labels.font.size = Pt(9)
    data_labels.font.name = BODY_FONT
    data_labels.font.bold = True

    palette = [BRAND_BLUE, ACCENT_RED, RGBColor(0x6C, 0x8E, 0xB8)]
    for idx, s in enumerate(chart.series):
        fill = s.format.fill
        fill.solid()
        fill.fore_color.rgb = palette[idx % len(palette)]
        s.format.line.color.rgb = palette[idx % len(palette)]

    return {
        "rendered": True,
        "chart_title": chart_title,
        "chart_type": chart_data.get("chart_type"),
        "categories": categories,
        "series_count": len(series),
        "removed_text_shapes": removed,
    }


def _coerce_float(value, default: Optional[float] = None) -> Optional[float]:
    try:
        if value is None or value == "":
            return default
        return float(value)
    except (TypeError, ValueError):
        return default


def _matrix_points_from_chart_data(chart_data: dict) -> list[dict]:
    rows = chart_data.get("source_rows") or []
    points = []
    for row in rows:
        if not isinstance(row, dict):
            continue
        label = str(row.get("label") or "").strip()
        x = _coerce_float(row.get("x"))
        y = _coerce_float(row.get("y"))
        if x is None or y is None:
            value = row.get("value")
            if isinstance(value, dict):
                x = _coerce_float(value.get("x"))
                y = _coerce_float(value.get("y"))
            elif isinstance(value, (list, tuple)) and len(value) >= 2:
                x = _coerce_float(value[0])
                y = _coerce_float(value[1])
        if label and x is not None and y is not None:
            points.append({"label": label, "x": x, "y": y, "note": row.get("note", "")})

    if points:
        return points

    categories = chart_data.get("categories") or []
    series = chart_data.get("series") or []
    if len(series) < 2:
        return []
    x_values = series[0].get("values") or []
    y_values = series[1].get("values") or []
    for idx, label in enumerate(categories):
        if idx >= len(x_values) or idx >= len(y_values):
            continue
        x = _coerce_float(x_values[idx])
        y = _coerce_float(y_values[idx])
        if label and x is not None and y is not None:
            points.append({"label": str(label), "x": x, "y": y, "note": ""})
    return points


def _normalize(value: float, values: list[float]) -> float:
    min_value = min(values)
    max_value = max(values)
    if max_value == min_value:
        return 0.5
    return (value - min_value) / (max_value - min_value)


def _add_textbox(slide, left: int, top: int, width: int, height: int, text: str, font_size: int = 8, bold: bool = False) -> None:
    textbox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.name = BODY_FONT
    run.font.bold = bold
    run.font.color.rgb = TEXT_GRAY


def render_matrix_slide(slide, slide_data: dict, layout: dict) -> dict:
    chart_data = slide_data.get("chart_data") or {}
    body_copy = slide_data.get("body_copy") or {}
    points = _matrix_points_from_chart_data(chart_data)
    if len(points) < 2:
        return {"rendered": False, "reason": "matrix_page needs at least two points with x/y values"}

    cleanup_box = layout.get("cleanup_box")
    removed_shapes = []
    if cleanup_box:
        for shape in list(slide.shapes):
            if intersects(shape, cleanup_box):
                text = getattr(shape, "text", "").strip() if hasattr(shape, "text") else ""
                remove_shape(shape)
                removed_shapes.append({"shape_name": shape.name, "text": text[:80]})

    left, top, width, height = layout["matrix_box"]
    axis_label_x = body_copy.get("matrix_label_x") or chart_data.get("x_axis_label") or "Axis X"
    axis_label_y = body_copy.get("matrix_label_y") or chart_data.get("y_axis_label") or "Axis Y"

    # Background panel.
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Emu(left),
        Emu(top),
        Emu(width),
        Emu(height),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFC)
    panel.line.color.rgb = GRID_GRAY
    panel.line.width = Pt(1)

    mid_x = left + width // 2
    mid_y = top + height // 2
    for line_left, line_top, line_width, line_height in [
        (mid_x, top, 0, height),
        (left, mid_y, width, 0),
    ]:
        line = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Emu(line_left),
            Emu(line_top),
            Emu(max(line_width, 1)),
            Emu(max(line_height, 1)),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = GRID_GRAY
        line.line.color.rgb = GRID_GRAY

    _add_textbox(slide, left, top + height + 80000, width, 180000, axis_label_x, 8, True)
    _add_textbox(slide, left - 240000, top + height // 2 - 120000, 220000, 260000, axis_label_y, 8, True)

    x_values = [point["x"] for point in points]
    y_values = [point["y"] for point in points]
    bubble_size = 175000
    plotted = []
    for idx, point in enumerate(points[:8]):
        x_norm = _normalize(point["x"], x_values)
        y_norm = _normalize(point["y"], y_values)
        cx = left + int(260000 + x_norm * (width - 520000))
        cy = top + int(height - 260000 - y_norm * (height - 520000))

        bubble = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Emu(cx - bubble_size // 2),
            Emu(cy - bubble_size // 2),
            Emu(bubble_size),
            Emu(bubble_size),
        )
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = BRAND_BLUE if idx else ACCENT_RED
        bubble.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        bubble.line.width = Pt(1)

        label_left = min(max(cx - 360000, left + 20000), left + width - 720000)
        label_top = min(max(cy + 90000, top + 20000), top + height - 220000)
        _add_textbox(slide, label_left, label_top, 720000, 180000, point["label"], 7, idx == 0)
        plotted.append({"label": point["label"], "x": point["x"], "y": point["y"]})

    return {
        "rendered": True,
        "chart_title": chart_data.get("title") or body_copy.get("matrix_title") or "",
        "chart_type": "matrix",
        "points": plotted,
        "truncated_points": max(0, len(points) - len(plotted)),
        "removed_existing_matrix_shapes": removed_shapes,
    }


def add_metric_card(slide, left: int, top: int, width: int, height: int, label: str, value: str, accent: RGBColor) -> None:
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Emu(left),
        Emu(top),
        Emu(width),
        Emu(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape.line.color.rgb = GRID_GRAY
    shape.line.width = Pt(1)

    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True

    p1 = text_frame.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = label
    r1.font.size = Pt(10)
    r1.font.name = BODY_FONT
    r1.font.bold = True
    r1.font.color.rgb = TEXT_GRAY

    p2 = text_frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = value
    r2.font.size = Pt(24)
    r2.font.name = BODY_FONT
    r2.font.bold = True
    r2.font.color.rgb = accent


def add_supporting_note(slide, left: int, top: int, width: int, height: int, text: str) -> None:
    from pptx.enum.text import PP_ALIGN

    textbox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.name = BODY_FONT
    r.font.color.rgb = TEXT_GRAY


def add_left_textbox(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    font_size: int = 9,
    bold: bool = False,
    color: RGBColor = TEXT_GRAY,
) -> None:
    textbox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.05
    run = p.add_run()
    run.text = str(text or "")
    run.font.size = Pt(font_size)
    run.font.name = BODY_FONT
    run.font.bold = bold
    run.font.color.rgb = color


def add_panel_box(slide, left: int, top: int, width: int, height: int) -> None:
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Emu(left),
        Emu(top),
        Emu(width),
        Emu(height),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFC)
    panel.line.color.rgb = GRID_GRAY
    panel.line.width = Pt(1)


def remove_think_cell_ole_shapes(slide) -> list[dict]:
    removed = []
    for shape in list(slide.shapes):
        name = str(getattr(shape, "name", "") or "")
        if "think-cell" not in name.lower():
            continue
        remove_shape(shape)
        removed.append({"shape_name": name, "shape_type": str(getattr(shape, "shape_type", ""))})
    return removed


def format_metric_value(label: str, value: Union[float, int, str], unit: str) -> str:
    if isinstance(value, str):
        text = value.strip()
        if not text:
            return ""
        if any(token in text for token in ("%", "亿", "万", "元", "RMB", "USD", "$")):
            return text
        try:
            value = float(text.replace(",", ""))
        except ValueError:
            return text
    is_integer = isinstance(value, int) or (isinstance(value, float) and value.is_integer())
    if "CAGR" in label.upper():
        return f"{int(value) if is_integer else f'{value:.1f}'}%"
    if "%" in label or "同比" in label or "增速" in label:
        return f"{value:+.0f}%" if is_integer else f"{value:+.1f}%"
    if "份额" in label or "占比" in label or "比例" in label or "渗透率" in label:
        return f"{int(value)}%" if is_integer else f"{value:.1f}%"
    if "%" in unit and not any(token in unit for token in ("亿", "万", "元", "RMB", "USD", "$")):
        return f"{int(value)}%" if is_integer else f"{value:.1f}%"
    if "亿元" in label or "规模" in label or "亿元" in unit or "人民币" in unit:
        if is_integer:
            return f"{int(value):,}亿元"
        return f"{value:,.1f}亿元"
    if is_integer:
        return f"{int(value):,}"
    return f"{value:,.1f}"


def render_secondary_module(slide, chart_data: dict, layout: dict) -> dict:
    module = chart_data.get("secondary_module") or {}
    if not isinstance(module, dict):
        module = {}
    module_type = str(module.get("module_type") or "metric_cards").lower()
    rows = module.get("rows") or []
    rows = [row for row in rows if isinstance(row, dict)]
    if not rows:
        return {
            "rendered": False,
            "reason": "missing secondary module rows; refused to reuse primary chart source_rows as fallback",
        }

    left, top, width, height = layout["side_box"]
    removed = remove_text_shapes_in_box(slide, layout["side_box"])
    add_panel_box(slide, left, top, width, height)
    title = str(module.get("title") or chart_data.get("secondary_title") or "").strip()
    title_height = 260000 if title else 0
    if title:
        add_left_textbox(slide, left + 120000, top + 70000, width - 240000, title_height, title, 9, True, BRAND_BLUE)

    content_top = top + title_height + 110000
    content_height = height - title_height - 180000
    if module_type in {"mini_table", "table", "segmentation_table"}:
        headers = module.get("headers") or ["Metric", "Read-through"]
        col_count = min(max(len(headers), 2), 3)
        row_count = min(len(rows), 4) + 1
        table_shape = slide.shapes.add_table(
            row_count,
            col_count,
            Emu(left + 90000),
            Emu(content_top),
            Emu(width - 180000),
            Emu(content_height),
        )
        table = table_shape.table
        for col_idx in range(col_count):
            table.columns[col_idx].width = Emu(int((width - 180000) / col_count))
            set_cell_text(table.cell(0, col_idx), headers[col_idx] if col_idx < len(headers) else "", 8.5, True, RGBColor(0xFF, 0xFF, 0xFF))
        for row_idx, row in enumerate(rows[: row_count - 1], start=1):
            values = [
                row.get("label", ""),
                format_metric_value(str(row.get("label") or ""), row.get("value", ""), row.get("unit") or row.get("value_unit") or chart_data.get("unit") or ""),
                row.get("note", ""),
            ]
            for col_idx in range(col_count):
                set_cell_text(table.cell(row_idx, col_idx), values[col_idx] if col_idx < len(values) else "", 8.5, False)
        style_table_shape(table_shape)
        return {"rendered": True, "module_type": module_type, "rows": row_count, "removed_text_shapes": removed}

    card_count = min(max(len(rows), 2), 3)
    gap = 90000
    card_height = (content_height - gap * (card_count - 1)) // card_count
    accents = [ACCENT_RED, BRAND_BLUE, RGBColor(0x4D, 0x7C, 0x3A)]
    for idx, row in enumerate(rows[:card_count]):
        label = row.get("label", "")
        row_unit = row.get("unit") or row.get("value_unit") or chart_data.get("unit") or ""
        add_metric_card(
            slide,
            left + 160000,
            content_top + idx * (card_height + gap),
            width - 320000,
            card_height,
            label,
            format_metric_value(label, row.get("value", ""), row_unit),
            accents[idx],
        )
    return {"rendered": True, "module_type": "metric_cards", "cards": card_count, "removed_text_shapes": removed}


def render_bottom_takeaways(slide, slide_data: dict, layout: dict) -> dict:
    body = slide_data.get("body_copy") or {}
    takeaways = [
        str(body.get("bullet_1") or "").strip(),
        str(body.get("bullet_2") or "").strip(),
    ]
    takeaways = [item for item in takeaways if item]
    if not takeaways:
        return {"rendered": False, "reason": "missing bottom takeaway bullets"}
    left, top, width, height = layout["bottom_box"]
    removed = remove_text_shapes_in_box(slide, layout["bottom_box"])
    line_height = height // max(len(takeaways), 1)
    for idx, text in enumerate(takeaways[:2]):
        add_left_textbox(slide, left, top + idx * line_height, width, line_height, f"• {text}", 9, False)
    return {"rendered": True, "takeaways": len(takeaways[:2]), "removed_text_shapes": removed}


def render_slide1_dynamic_overview(slide, slide_data: dict, layout: dict) -> dict:
    chart_data = slide_data.get("chart_data") or {}
    chart_type = str(chart_data.get("chart_type") or "").lower()
    cleanup_box = layout.get("cleanup_box")
    cleared = []
    if cleanup_box:
        cleared = clear_text_shapes_in_box(slide, cleanup_box)

    if chart_type in {"bar", "column", "clustered_bar", "clustered_column", "stacked_bar", "stacked_column", "line", "line_chart"}:
        chart_result = build_chart(slide, slide_data, layout)
    else:
        chart_result = {
            "rendered": False,
            "reason": f"dynamic overview requires a chart type supported by the deterministic renderer, found '{chart_type}'",
        }
    return {
        "rendered": bool(chart_result.get("rendered")),
        "mode": "dynamic_overview",
        "chart": chart_result,
        "left_key_messages_preserved": True,
        "right_visual_only": True,
        "cleared_existing_visual_text_shapes": cleared,
    }


def render_slide1_visual(slide, slide_data: dict, layout: dict) -> dict:
    chart_data = slide_data.get("chart_data") or {}
    chart_type = str(chart_data.get("chart_type") or "").lower()
    if chart_type in {"none", "no_chart", "text"}:
        return {
            "rendered": False,
            "required_render": False,
            "reason": "slide 1 chart_type is none; visual area intentionally left as text/context",
        }
    if chart_type in {"bar", "column", "clustered_bar", "clustered_column", "stacked_bar", "stacked_column", "line", "line_chart"}:
        return build_chart(slide, slide_data, layout)

    rows = chart_data.get("source_rows") or []
    if len(rows) < 2:
        return {
            "rendered": False,
            "reason": "slide 1 needs chart_data series/categories for charts or at least two source_rows for metric cards",
        }

    chart_title = chart_data.get("title") or ""
    apply_chart_title(slide, chart_title, layout)

    left, top, width, height = layout["visual_box"]
    card_count = min(max(len(rows), 2), 3)
    gap = 180000
    card_width = (width - gap * (card_count - 1)) // card_count
    display_note = chart_data.get("display_note") or chart_data.get("on_slide_note") or ""
    card_height = int(height * (0.58 if display_note else 0.68))
    note_height = 480000 if display_note else 0
    card_top = top + max(0, int((height - card_height - note_height) / 2))

    unit = chart_data.get("unit") or ""
    accents = [ACCENT_RED, BRAND_BLUE, RGBColor(0x4D, 0x7C, 0x3A)]

    for idx, row in enumerate(rows[:card_count]):
        label = row.get("label", "")
        row_unit = row.get("unit") or row.get("value_unit") or unit
        add_metric_card(
            slide,
            left + idx * (card_width + gap),
            card_top,
            card_width,
            card_height,
            label,
            format_metric_value(label, row.get("value", 0), row_unit),
            accents[idx],
        )

    if display_note:
        add_supporting_note(slide, left, card_top + card_height + 160000, width, note_height, display_note)

    return {
        "rendered": True,
        "chart_title": chart_title,
        "chart_type": chart_data.get("chart_type"),
        "mode": "metric_cards",
    }


def split_table_cells(text: str, expected_cols: Optional[int] = None) -> list[str]:
    value = str(text or "").strip()
    if not value:
        return [""] * expected_cols if expected_cols else []
    if "｜" in value:
        cells = [part.strip() for part in value.split("｜")]
    elif "|" in value:
        cells = [part.strip() for part in value.split("|")]
    elif " / " in value:
        cells = [part.strip() for part in value.split(" / ")]
    else:
        cells = [value]
    if expected_cols is None:
        return cells
    if len(cells) > expected_cols:
        cells = cells[: expected_cols - 1] + [" / ".join(cells[expected_cols - 1 :])]
    return cells + [""] * (expected_cols - len(cells))


def estimate_column_weights(headers: list[str], rows: list[list[str]]) -> list[float]:
    weights: list[float] = []
    col_count = len(headers)
    for col_idx in range(col_count):
        values = [headers[col_idx]] + [row[col_idx] for row in rows if col_idx < len(row)]
        max_len = max((len(str(value)) for value in values), default=1)
        weights.append(max(0.8, min(3.2, max_len / 8.0)))
    if weights:
        weights[0] = min(weights[0], 1.3)
    return weights


def count_table_cells(text: str) -> int:
    value = str(text or "").strip()
    if not value:
        return 0
    if "｜" in value:
        return len([part for part in value.split("｜")])
    if "|" in value:
        return len([part for part in value.split("|")])
    if " / " in value:
        return len([part for part in value.split(" / ")])
    return 1


def remove_text_shapes_in_box(slide, box: tuple[int, int, int, int]) -> list[dict]:
    removed = []
    for shape in list(slide.shapes):
        if not hasattr(shape, "text_frame"):
            continue
        if not intersects(shape, box):
            continue
        text = getattr(shape, "text", "").strip() if hasattr(shape, "text") else ""
        remove_shape(shape)
        removed.append({"shape_name": shape.name, "text": text[:80]})
    return removed


def clear_text_shapes_in_box(slide, box: tuple[int, int, int, int]) -> list[dict]:
    cleared = []
    for shape in list(slide.shapes):
        if not hasattr(shape, "text_frame"):
            continue
        if not intersects(shape, box):
            continue
        text = getattr(shape, "text", "").strip() if hasattr(shape, "text") else ""
        if not text:
            continue
        shape.text_frame.clear()
        cleared.append({"shape_name": shape.name, "text": text[:80]})
    return cleared


def remove_shapes_in_box(slide, box: tuple[int, int, int, int]) -> list[dict]:
    removed = []
    for shape in list(slide.shapes):
        if not intersects(shape, box):
            continue
        text = getattr(shape, "text", "").strip() if hasattr(shape, "text") else ""
        remove_shape(shape)
        removed.append({"shape_name": shape.name, "text": text[:80]})
    return removed


def set_cell_text(cell, text: str, font_size: float, bold: bool = False, color: RGBColor = TEXT_GRAY) -> None:
    cell.text = ""
    cell.margin_left = Emu(30000)
    cell.margin_right = Emu(30000)
    cell.margin_top = Emu(18000)
    cell.margin_bottom = Emu(18000)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.text_frame.word_wrap = True
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.line_spacing = 1.0
    run = paragraph.add_run()
    run.text = str(text or "")
    run.font.name = BODY_FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def style_table_shape(table_shape, header_fill: RGBColor = BRAND_BLUE) -> None:
    table = table_shape.table
    for row_idx, row in enumerate(table.rows):
        for cell in row.cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if row_idx == 0 else RGBColor(0xFF, 0xFF, 0xFF)


def render_slide2_mini_table(slide, slide_data: dict, layout: dict) -> dict:
    body = slide_data.get("body_copy") or {}
    headers = [body.get("table_header_1", ""), body.get("table_header_2", "")]
    rows = [body.get(f"table_row_{idx}", "") for idx in range(1, 4)]
    if not any(str(item).strip() for item in headers + rows):
        return {"rendered": False, "reason": "missing slide 2 table body_copy fields"}

    table_box = layout["table_box"]
    removed = remove_text_shapes_in_box(slide, table_box)
    left, top, width, height = table_box
    table_shape = slide.shapes.add_table(4, 2, Emu(left), Emu(top), Emu(width), Emu(height))
    table = table_shape.table
    table.columns[0].width = Emu(int(width * 0.42))
    table.columns[1].width = Emu(int(width * 0.58))
    for idx, header in enumerate(headers):
        set_cell_text(table.cell(0, idx), header, TABLE_HEADER_FONT_SIZE, True, RGBColor(0xFF, 0xFF, 0xFF))
    for row_idx, row_text in enumerate(rows, start=1):
        for col_idx, cell_text in enumerate(split_table_cells(row_text, 2)):
            set_cell_text(table.cell(row_idx, col_idx), cell_text, TABLE_BODY_FONT_SIZE, False)
    style_table_shape(table_shape)
    return {"rendered": True, "rows": 4, "columns": 2, "removed_text_shapes": removed}


def render_slide2_chart_plus_table(slide, slide_data: dict, layout: dict) -> dict:
    chart_result = build_chart(slide, slide_data, layout)
    table_result = render_slide2_mini_table(slide, slide_data, layout)
    rendered = chart_result.get("rendered", False) and table_result.get("rendered", False)
    return {
        "rendered": rendered,
        "chart": chart_result,
        "table": table_result,
        "chart_title": chart_result.get("chart_title", ""),
        "chart_type": chart_result.get("chart_type", ""),
        "table_rows": table_result.get("rows", 0),
        "table_columns": table_result.get("columns", 0),
    }


def render_slide6_compare_table(slide, slide_data: dict, layout: dict) -> dict:
    body = slide_data.get("body_copy") or {}
    header_cells = split_table_cells(body.get("table_header", ""))
    col_count = len(header_cells)
    if col_count < 3 or col_count > 6:
        return {
            "rendered": False,
            "reason": f"slide 6 compare table requires 3-6 explicit header columns; found {col_count}",
        }

    row_texts = [
        body.get(f"table_row_{idx}", "")
        for idx in range(1, 7)
        if str(body.get(f"table_row_{idx}", "")).strip()
    ]
    if not header_cells and not row_texts:
        return {"rendered": False, "reason": "missing slide 6 compare table body_copy fields"}
    if len(row_texts) < 3:
        return {"rendered": False, "reason": "slide 6 compare table needs at least 3 populated peer rows"}

    rows = []
    for row_idx, row_text in enumerate(row_texts, start=1):
        cells = split_table_cells(row_text)
        if len(cells) != col_count:
            return {"rendered": False, "reason": f"row {row_idx} has {len(cells)} cells; expected {col_count}"}
        if any(not cell.strip() for cell in cells):
            return {"rendered": False, "reason": f"row {row_idx} contains blank cell"}
        rows.append(cells)

    table_box = layout["table_box"]
    removed = remove_text_shapes_in_box(slide, table_box)
    left, top, width, height = table_box
    table_shape = slide.shapes.add_table(len(rows) + 1, col_count, Emu(left), Emu(top), Emu(width), Emu(height))
    table = table_shape.table
    col_weights = layout.get("column_weight_overrides", {}).get(str(col_count)) if isinstance(layout.get("column_weight_overrides"), dict) else None
    if not col_weights:
        col_weights = estimate_column_weights(header_cells, rows)
    total = sum(col_weights)
    for idx, weight in enumerate(col_weights):
        table.columns[idx].width = Emu(int(width * weight / total))
    header_font = TABLE_HEADER_FONT_SIZE
    body_font = TABLE_BODY_FONT_SIZE
    for idx, header in enumerate(header_cells):
        set_cell_text(table.cell(0, idx), header, header_font, True, RGBColor(0xFF, 0xFF, 0xFF))
    for row_idx, row_cells in enumerate(rows, start=1):
        for col_idx, cell_text in enumerate(row_cells):
            set_cell_text(table.cell(row_idx, col_idx), cell_text, body_font, False)
    style_table_shape(table_shape)
    return {"rendered": True, "rows": len(rows) + 1, "columns": col_count, "removed_text_shapes": removed}


def render_quant_slide(prs: Presentation, storyboard: dict, slide_no: int, render_layouts: dict[int, dict]) -> dict:
    slide_data = find_slide_data(storyboard, slide_no)
    if not slide_data:
        return {"slide_no": slide_no, "rendered": False, "reason": f"slide {slide_no} not found in storyboard"}

    page_type = slide_data.get("selected_page_type")
    slide_layouts = render_layouts.get(slide_no, {})
    layout = slide_layouts.get(page_type)
    if not layout:
        return {"slide_no": slide_no, "rendered": False, "reason": f"unsupported page type: {page_type}"}

    if len(prs.slides) < slide_no:
        return {"slide_no": slide_no, "rendered": False, "reason": "clean deck has fewer slides than expected"}

    slide = prs.slides[slide_no - 1]
    if slide_no == 1 and page_type == "industry_overview_dynamic_page":
        result = render_slide1_dynamic_overview(slide, slide_data, layout)
    elif page_type == "matrix_page":
        result = render_matrix_slide(slide, slide_data, layout)
    elif slide_no == 2 and page_type == "chart_plus_mini_table_page":
        result = render_slide2_chart_plus_table(slide, slide_data, layout)
    elif slide_no == 6 and page_type == "compare_table_page":
        result = render_slide6_compare_table(slide, slide_data, layout)
    elif slide_no == 1:
        result = render_slide1_visual(slide, slide_data, layout)
    else:
        result = build_chart(slide, slide_data, layout)
    result["slide_no"] = slide_no
    result["selected_page_type"] = page_type
    return result


def has_postprocess_renderer(slide_data: dict, render_layouts: dict[int, dict]) -> bool:
    slide_no = int(slide_data["slide_no"])
    page_type = slide_data.get("selected_page_type")
    return page_type in render_layouts.get(slide_no, {})


def skipped_non_rendered_slide(slide_data: dict) -> dict:
    return {
        "slide_no": int(slide_data["slide_no"]),
        "selected_page_type": slide_data.get("selected_page_type"),
        "rendered": False,
        "required_render": False,
        "skipped": True,
        "reason": "no deterministic postprocess renderer for this page type; text/table layout remains authoritative",
    }


def save_presentation(prs: Presentation, output_path: Path) -> None:
    if output_path.exists():
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = Path(tmp.name)
        try:
            prs.save(tmp_path)
            shutil.move(str(tmp_path), output_path)
        finally:
            if tmp_path.exists():
                os.unlink(tmp_path)
        return
    prs.save(output_path)


def _remove_ole_graphic_frames(xml_text: str) -> tuple[str, int]:
    removed_count = 0
    frame_pattern = re.compile(r"<p:graphicFrame\b.*?</p:graphicFrame>", re.DOTALL)

    def replace_frame(match: re.Match) -> str:
        nonlocal removed_count
        frame = match.group(0)
        lowered = frame.lower()
        if "<p:oleobj" in lowered or "tclayout.activedocument" in lowered or "think-cell" in lowered:
            removed_count += 1
            return ""
        return frame

    cleaned = frame_pattern.sub(replace_frame, xml_text)
    ole_obj_pattern = re.compile(r"<p:oleObj\b.*?</p:oleObj>", re.DOTALL)
    cleaned, fallback_removed = ole_obj_pattern.subn("", cleaned)
    removed_count += fallback_removed
    return cleaned, removed_count


def _remove_ole_relationships(rels_text: str) -> tuple[str, int]:
    rel_pattern = re.compile(r"<Relationship\b[^>]*/>", re.DOTALL)
    removed_count = 0

    def replace_rel(match: re.Match) -> str:
        nonlocal removed_count
        rel = match.group(0)
        lowered = rel.lower()
        if "/oleobject" in lowered or "target=\"../embeddings/oleobject" in lowered:
            removed_count += 1
            return ""
        return rel

    return rel_pattern.sub(replace_rel, rels_text), removed_count


def _remove_ole_content_types(content_types_text: str) -> tuple[str, int]:
    override_pattern = re.compile(r"<Override\b[^>]*PartName=\"/ppt/embeddings/oleObject[^\"]*\"[^>]*/>", re.DOTALL)
    cleaned, removed_count = override_pattern.subn("", content_types_text)
    return cleaned, removed_count


def sanitize_ole_artifacts(pptx_path: Path) -> dict:
    with ZipFile(pptx_path, "r") as archive:
        original_items = [(name, archive.read(name)) for name in archive.namelist()]

    removed_parts = []
    updated_parts = []
    sanitized_items = []

    for name, payload in original_items:
        lower_name = name.lower()
        if lower_name.startswith("ppt/embeddings/oleobject") or (
            lower_name.startswith("ppt/embeddings/") and lower_name.endswith(".bin")
        ):
            removed_parts.append(name)
            continue

        if lower_name.endswith(".rels") or lower_name.endswith(".xml"):
            text = payload.decode("utf-8", errors="ignore")
            updated_text = text
            removed_count = 0

            if lower_name.endswith(".rels"):
                updated_text, removed_count = _remove_ole_relationships(updated_text)
            elif lower_name == "[content_types].xml":
                updated_text, removed_count = _remove_ole_content_types(updated_text)
            elif lower_name.startswith("ppt/slides/") or lower_name.startswith("ppt/slidelayouts/") or lower_name.startswith("ppt/slidemasters/"):
                updated_text, removed_count = _remove_ole_graphic_frames(updated_text)

            if "think-cell" in updated_text.lower():
                updated_text = re.sub(r"think-cell[^<]*", "removed embedded object", updated_text, flags=re.IGNORECASE)
                removed_count += 1
            if "tclayout.activedocument" in updated_text.lower():
                updated_text = re.sub(r"TCLayout\.ActiveDocument\.\d+", "removed.embedded.object", updated_text, flags=re.IGNORECASE)
                removed_count += 1

            if removed_count:
                updated_parts.append({"part": name, "removed_count": removed_count})
                payload = updated_text.encode("utf-8")

        sanitized_items.append((name, payload))

    if not removed_parts and not updated_parts:
        return {"removed_parts": [], "updated_parts": []}

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp_path = Path(tmp.name)
    try:
        with ZipFile(tmp_path, "w", ZIP_DEFLATED) as archive:
            for name, payload in sanitized_items:
                archive.writestr(name, payload)
        shutil.move(str(tmp_path), pptx_path)
    finally:
        if tmp_path.exists():
            os.unlink(tmp_path)

    return {"removed_parts": removed_parts, "updated_parts": updated_parts}


def postprocess(input_ppt: Path, storyboard_path: Path, output_ppt: Path, render_layouts_path: Path = DEFAULT_RENDER_LAYOUTS_PATH) -> dict:
    storyboard = load_json(storyboard_path)
    render_layouts = load_render_layouts(render_layouts_path)
    prs = Presentation(str(input_ppt))

    removed_labels = remove_scaffold_labels(prs)
    removed_think_cell_ole = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        removed = remove_think_cell_ole_shapes(slide)
        if removed:
            removed_think_cell_ole.append({"slide_no": slide_idx, "removed": removed})
    page_number_updates = rewrite_page_numbers(prs)
    chart_results = []
    for slide_data in storyboard.get("slides", []):
        if slide_data.get("chart_data") or (
            int(slide_data.get("slide_no", 0)) in {2, 6}
            and has_postprocess_renderer(slide_data, render_layouts)
        ):
            if has_postprocess_renderer(slide_data, render_layouts):
                chart_results.append(render_quant_slide(prs, storyboard, int(slide_data["slide_no"]), render_layouts))
            else:
                chart_results.append(skipped_non_rendered_slide(slide_data))

    save_presentation(prs, output_ppt)
    ole_sanitization = sanitize_ole_artifacts(output_ppt)

    return {
        "input_ppt": str(input_ppt),
        "storyboard": str(storyboard_path),
        "render_layouts": str(render_layouts_path),
        "output_ppt": str(output_ppt),
        "removed_scaffold_labels": removed_labels,
        "removed_think_cell_ole": removed_think_cell_ole,
        "ole_sanitization": ole_sanitization,
        "page_number_updates": page_number_updates,
        "chart_rendering": chart_results,
    }


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Post-process filled PPT with object-level visual cleanup and chart rendering."
    )
    parser.add_argument("--input-ppt", required=True, help="Path to cleaned PPTX input.")
    parser.add_argument("--storyboard", required=True, help="Path to industry_storyboard.json.")
    parser.add_argument("--output", required=True, help="Path to write the post-processed PPTX.")
    parser.add_argument(
        "--render-layouts",
        default=str(DEFAULT_RENDER_LAYOUTS_PATH),
        help="Path to templates/render_layouts.json with deterministic renderer coordinates.",
    )
    parser.add_argument("--log", help="Optional path to write a JSON log.")
    parser.add_argument(
        "--fail-on-unrendered",
        action="store_true",
        help="Exit non-zero if a required deterministic visual renderer fails.",
    )
    parser.add_argument(
        "--allow-ungated-debug",
        action="store_true",
        help="Bypass the pre-PPT stage gate. Use only for local diagnostics, never delivery.",
    )
    args = parser.parse_args()

    output_path = Path(args.output)
    try:
        require_pre_ppt_gate(output_path.parent, allow_ungated_debug=args.allow_ungated_debug)
        result = postprocess(Path(args.input_ppt), Path(args.storyboard), output_path, Path(args.render_layouts))
    except Exception as exc:
        raise SystemExit(str(exc)) from exc
    if args.log:
        save_json(result, Path(args.log))
    print(json.dumps(result, ensure_ascii=False, indent=2))
    if args.fail_on_unrendered:
        failed = [
            item
            for item in result["chart_rendering"]
            if not item.get("rendered") and item.get("required_render", True)
        ]
        if failed:
            raise SystemExit(1)


if __name__ == "__main__":
    main()
