#!/usr/bin/env python3
"""Generate template-aware text fit rules for title and takeaway placeholders."""

from __future__ import annotations

import argparse
import json
import re
from pathlib import Path
from typing import Any

from pptx import Presentation

from json_utils import load_json_file


TOKEN_RE = re.compile(r"\{\{slide_(\d{2})([a-z]?|)_(title|takeaway)\}\}")
EMU_PER_INCH = 914400


def token_shape_metrics(template_path: Path) -> dict[str, dict[str, Any]]:
    prs = Presentation(template_path)
    metrics: dict[str, dict[str, Any]] = {}
    for physical_slide_no, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = "\n".join(p.text for p in shape.text_frame.paragraphs)
            match = TOKEN_RE.search(text)
            if not match:
                continue
            token = match.group(0)
            font_sizes = []
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.font.size:
                    font_sizes.append(paragraph.font.size.pt)
                for run in paragraph.runs:
                    if run.font.size:
                        font_sizes.append(run.font.size.pt)
            metrics[token] = {
                "physical_slide_no": physical_slide_no,
                "width_in": round(shape.width / EMU_PER_INCH, 3),
                "height_in": round(shape.height / EMU_PER_INCH, 3),
                "font_size_pt": round(max(font_sizes), 1) if font_sizes else None,
            }
    return metrics


def iter_mapping_tokens(ppt_mapping: dict[str, Any]) -> list[dict[str, Any]]:
    rows: list[dict[str, Any]] = []
    for slide in ppt_mapping.get("slides", []):
        slide_no = slide.get("slide_no")
        selected_page_type = slide.get("selected_page_type")
        for token in slide.get("tokens", []):
            rows.append({
                "slide_no": slide_no,
                "page_type": selected_page_type,
                "placeholder": token.get("placeholder"),
                "field_name": token.get("field_name"),
            })
        variants = slide.get("variants") or slide.get("controlled_variants") or {}
        for page_type, variant in variants.items():
            for token in variant.get("tokens", []):
                rows.append({
                    "slide_no": slide_no,
                    "page_type": page_type,
                    "placeholder": token.get("placeholder"),
                    "field_name": token.get("field_name"),
                })
    return rows


def max_line_units(width_in: float, font_size_pt: float, width_factor: float) -> float:
    width_pt = width_in * 72
    return round(width_pt / (font_size_pt * width_factor), 1)


def build_rules(template_path: Path, ppt_mapping_path: Path, title_font_size: float) -> dict[str, Any]:
    ppt_mapping = load_json_file(ppt_mapping_path)
    metrics = token_shape_metrics(template_path)

    fields: dict[str, dict[str, Any]] = {}
    for row in iter_mapping_tokens(ppt_mapping):
        field_name = row.get("field_name")
        if field_name not in {"slide_title", "main_takeaway"}:
            continue
        placeholder = row.get("placeholder")
        metric = metrics.get(placeholder or "")
        if not metric:
            continue

        font_size = metric["font_size_pt"] or title_font_size
        max_lines = 1 if field_name == "slide_title" else 2
        target_lines = 1
        width_factor = 1.0 if field_name == "slide_title" else 0.95
        field_rule = {
            "placeholder": placeholder,
            "field_name": field_name,
            "width_in": metric["width_in"],
            "height_in": metric["height_in"],
            "font_size_pt": font_size,
            "target_lines": target_lines,
            "max_lines": max_lines,
            "max_line_units": max_line_units(metric["width_in"], font_size, width_factor),
            "block_if_exceeds_max_lines": True,
        }
        key = f"{row['slide_no']}:{row['page_type']}:{field_name}"
        fields[key] = field_rule

    return {
        "_description": (
            "Template-derived line-fit rules for storyboard headline/main_message and "
            "PPT copy slide_title/main_takeaway. Regenerate after changing the PPT template."
        ),
        "template_file": str(template_path),
        "ppt_mapping": str(ppt_mapping_path),
        "estimation": {
            "unit_model": "CJK characters count as 1.0 unit; ASCII letters/digits count as 0.55; spaces and narrow punctuation count as 0.3.",
            "title_font_size_fallback_pt": title_font_size,
            "notes": "This is a deterministic preflight estimate, not a PowerPoint rendering engine.",
        },
        "storyboard_field_aliases": {
            "headline": "slide_title",
            "main_message": "main_takeaway",
        },
        "fields": fields,
    }


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate text-fit rules from a PPT template and ppt_mapping.json.")
    parser.add_argument("--template", default="assets/industry_section_template_master.pptx")
    parser.add_argument("--ppt-mapping", default="templates/ppt_mapping.json")
    parser.add_argument("--output", default="templates/text_fit_rules.json")
    parser.add_argument(
        "--title-font-size",
        type=float,
        default=28.0,
        help="Fallback title font size when the template token run inherits its font from layout/theme.",
    )
    args = parser.parse_args()

    result = build_rules(Path(args.template), Path(args.ppt_mapping), args.title_font_size)
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(result, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
