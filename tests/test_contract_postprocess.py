"""Contract tests: Group 10 - postprocess PPT visuals."""

from __future__ import annotations

import sys
from pathlib import Path

import pytest

SCRIPT_DIR = Path(__file__).resolve().parents[1] / "runtime" / "ib-pitchdeck-agent-industry-section" / "scripts"
sys.path.insert(0, str(SCRIPT_DIR))
sys.path.insert(0, str(SCRIPT_DIR / "output"))
sys.path.insert(0, str(SCRIPT_DIR / "_lib"))


class TestBuildChart:
    def test_invalid_series_returns_not_rendered_with_repair_hint(self):
        """build_chart should return rendered=False with repair_hint for invalid series."""
        try:
            from postprocess_ppt_visuals import build_chart
        except SystemExit:
            pytest.skip("postprocess_ppt_visuals not importable")

        result = build_chart(
            None,
            {
                "slide_no": 1,
                "selected_page_type": "industry_overview_dynamic_page",
                "chart_data": {
                    "chart_type": "bar",
                    "categories": ["2024"],
                    "series": ["Market size"],  # string instead of dict
                },
            },
            {},
        )
        assert result.get("rendered") is False, result
        assert result.get("reason") == "invalid chart_data.series item", result
        assert "repair_hint" in result, result


def test_style_guided_title_color_falls_back_when_template_primary_is_white(tmp_path: Path):
    try:
        import postprocess_ppt_visuals as visuals
    except SystemExit:
        pytest.skip("postprocess_ppt_visuals not importable")

    profile = tmp_path / "template_profile.json"
    profile.write_text(
        """
{
  "visual_style": {
    "colors": {
      "brand_primary": "#FFFFFF",
      "accent_red": "#AA3322",
      "grid_gray": "#CCCCCC",
      "text_gray": "#555555"
    },
    "typography": {
      "body": "Arial",
      "table_header": "Arial",
      "table_body": "Arial"
    }
  }
}
""".strip()
        + "\n",
        encoding="utf-8",
    )
    warnings: list[str] = []

    visuals._apply_template_profile_style(profile, warnings)

    assert tuple(visuals.BRAND_BLUE) == tuple(visuals.DEFAULT_BRAND_BLUE)
    assert any("too light" in item for item in warnings)


def test_visible_metric_claims_render_as_style_guided_kpi_cards():
    try:
        from pptx import Presentation
        from postprocess_ppt_visuals import render_metric_cards
    except (ImportError, SystemExit):
        pytest.skip("python-pptx or postprocess_ppt_visuals not importable")

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    result = render_metric_cards(
        slide,
        {
            "visible_metric_claims": [
                {"metric_id": "MET-001", "display_text": "2024 market indicator: 110 index"},
                {"metric_id": "MET-002", "claim": "Online channel share continues to expand"},
            ]
        },
        (500_000, 500_000, 6_000_000, 2_000_000),
    )

    visible_text = "\n".join(shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False))
    assert result["rendered"] is True
    assert result["items"] == 2
    assert "2024 market indicator: 110 index" in visible_text
    assert "Online channel share continues to expand" in visible_text


def test_visible_metric_claims_trigger_split_layout_with_body_copy():
    try:
        from pptx import Presentation
        from postprocess_ppt_visuals import generic_visual_layout
    except (ImportError, SystemExit):
        pytest.skip("python-pptx or postprocess_ppt_visuals not importable")

    prs = Presentation()
    layout = generic_visual_layout(
        prs,
        {
            "body_copy": {"point_1": "Evidence-backed mechanism."},
            "visible_metric_claims": [{"display_text": "2024 market indicator: 110 index"}],
        },
    )

    assert layout["body_box"][2] < layout["chart_box"][2]
    assert layout["chart_box"][0] > layout["body_box"][0]
