#!/usr/bin/env python3

from __future__ import annotations

import json
import contextlib
import io
import subprocess
import sys
import zipfile
from copy import deepcopy
from pathlib import Path

import pytest

from conftest import FIXED_PAGE_ROLES, ROLE_SCRIPT_DIRS, SCRIPT_IMPORT_PATHS, SKILL_DIR, SLIDE_NUMBERS, _write_json


def _body_copy_for(slide_no: int, page_type: str, template_registry: dict, blocks: list[dict]) -> dict[str, str]:
    from deck_blueprint_utils import active_body_fields, strict_layout_body_fields

    slide_data = {"slide_no": slide_no, "selected_page_type": page_type}
    if page_type == "compare_table_page":
        slide_data["compare_table_data"] = {"headers": ["Dimension"], "rows": [{"label": "Sample", "cells": ["Sample"]}]}
    fields = active_body_fields(strict_layout_body_fields(template_registry, slide_no, page_type), page_type, slide_data)
    result: dict[str, str] = {}
    for idx, field in enumerate(fields):
        result[field] = blocks[idx % len(blocks)]["copy"]
    return result


def _banker_page_pack(deck_blueprint_data: dict, template_registry: dict) -> dict:
    slides = []
    for slide in deck_blueprint_data["slides"]:
        slide_no = int(slide["slide_no"])
        page_evidence_ids = [
            f"EV-{((slide_no - 1) % 6) + 1:03d}",
            f"EV-{((slide_no) % 6) + 1:03d}",
        ]
        evidence_ids = sorted(
            {
                *page_evidence_ids,
                *(item for item in slide.get("evidence_ids", []) if item),
                *(
                    item
                    for block in slide.get("body_blocks", [])
                    for item in block.get("evidence_ids", [])
                    if item
                ),
            }
        ) or ["EV-001"]
        metric_id = f"MET-{((slide_no - 1) % 4) + 1:03d}"
        secondary_metric_id = f"MET-{((slide_no) % 4) + 1:03d}"
        metric_ids = sorted(
            {
                metric_id,
                secondary_metric_id,
                *(item for item in slide.get("metric_ids", []) if item),
                *(
                    item
                    for block in slide.get("body_blocks", [])
                    for item in block.get("metric_ids", [])
                    if item
                ),
            }
        )
        blocks = [
            {
                "role": f"point_{idx}",
                "copy": (
                    f"Page {slide_no} point {idx} interprets evidence for a pre-mandate client discussion, "
                    "linking industry structure, transaction framing, and market economics rather than listing generic facts."
                ),
                "evidence_ids": evidence_ids,
                "metric_ids": [metric_id] if idx in {1, 2} else [],
            }
            for idx in range(1, 5)
        ]
        chart_data = slide.get("chart_data") if isinstance(slide.get("chart_data"), dict) else {}
        compare_table_data = slide.get("compare_table_data") if isinstance(slide.get("compare_table_data"), dict) else {}
        if slide_no in {1, 2}:
            chart_data = {
                "chart_type": "bar",
                "title": f"Slide {slide_no} metric trend",
                "categories": ["2022", "2023", "2024"],
                "series": [{"name": "Indexed metric", "values": [80.0 + slide_no, 92.0 + slide_no, 108.0 + slide_no]}],
                "unit": "index",
                "source_rows": [
                    {"label": "2022", "value": 80.0 + slide_no, "metric_id": metric_id},
                    {"label": "2023", "value": 92.0 + slide_no, "metric_id": secondary_metric_id},
                    {"label": "2024", "value": 108.0 + slide_no, "metric_id": metric_id},
                ],
            }
        slides.append(
            {
                "slide_no": slide_no,
                "fixed_page_role": slide["fixed_page_role"],
                "banker_judgment": (
                    f"Page {slide_no} should communicate a banker judgment about market structure, transaction framing, "
                    "and sector economics using evidence rather than broad industry commentary, while also explaining "
                    "how the client should frame growth quality, competitive risk, and transaction logic before a mandate is signed."
                ),
                "page_argument": slide["page_argument"],
                "selected_page_type": slide["selected_page_type"],
                "deck_use": "可作标题",
                "headline": f"Page {slide_no} industry read",
                "main_message": (
                    f"Evidence links sector structure to transaction framing for page {slide_no}."
                ),
                "exhibit": {
                    **slide["exhibit"],
                    "data_or_evidence_inputs": [*evidence_ids, *metric_ids],
                    "visual_structure": "A dense exhibit using chart/table/card fields plus body blocks to fill the formal page.",
                },
                "body_blocks": blocks,
                "body_copy": _body_copy_for(slide_no, slide["selected_page_type"], template_registry, blocks),
                "evidence_ids": evidence_ids,
                "metric_ids": metric_ids,
                "visible_metric_claims": [
                    {
                        "location": "main exhibit",
                        "display_text": f"Indexed metric supports page {slide_no}",
                        "metric_ids": [metric_id],
                        "usage_type": "context_only",
                        "basis_note": "Test fixture metric binding.",
                    }
                ],
                "chart_data": chart_data,
                "compare_table_data": compare_table_data,
                "project_relevance_note": (
                    f"For a pre-mandate pitch, page {slide_no} turns the industry evidence into a transaction-framing discussion point."
                    if slide_no in {7, 8}
                    else ""
                ),
                "source_note": "Sources: " + "; ".join(evidence_ids),
                "caveats": [],
                "source_limitations": [],
            }
        )
    return {
        "schema_version": "banker_page_pack",
        "section_meta": {"target_company": "Example Target", "industry": "Example sector"},
        "deck_storyline": (
            "The section moves from industry scale and structure into competition, economics, and selective project relevance, "
            "using traceable data and banker judgment to support a pre-mandate client conversation. It should make the bank's industry view visible through charts, tables, caveats, and selective project relevance rather than through generic summary language."
        ),
        "evidence_policy": {
            "important_data": "Use MET rows from research_evidence_db with audit-grade source fields.",
            "normal_claims": "Use EV/source IDs and visible caveats.",
        },
        "deliverable_readiness": {
            "business_action": "client_ready",
            "decision_note": "Fixture has enough evidence-linked page copy, visible metrics, and exhibits to proceed as a pre-mandate client pitch section.",
        },
        "key_data_audit": [
            {
                "metric_id": f"MET-{idx:03d}",
                "indicator": f"Fixture indexed metric {idx}",
                "value": str(100 + idx),
                "unit": "index",
                "period": "2024",
                "geography": "Fixture market",
                "source": "Contract fixture source",
                "original_location": "Fixture table 1",
                "original_excerpt": "Fixture excerpt supports the indexed value.",
                "usage_in_deck": "Visible metric claim and exhibit support.",
                "remarks": "Contract-test audit row.",
            }
            for idx in range(1, 5)
        ],
        "conflict_data_notes": [],
        "slides": slides,
    }


def _run(script_name: str, args: list[str]) -> subprocess.CompletedProcess:
    script_path = SKILL_DIR / "scripts" / "pipeline.py" if script_name == "pipeline.py" else ROLE_SCRIPT_DIRS[script_name]
    return subprocess.run(
        [sys.executable, str(script_path), *args],
        text=True,
        capture_output=True,
        cwd=str(SKILL_DIR),
        env={**__import__("os").environ, "PYTHONPATH": ":".join(str(path) for path in SCRIPT_IMPORT_PATHS)},
    )


def _compile_page_pack(run_dir: Path) -> subprocess.CompletedProcess:
    from pipeline import compile_page_pack

    args = ["pipeline.compile_page_pack", "--run-dir", str(run_dir)]
    stdout = io.StringIO()
    stderr = io.StringIO()
    try:
        with contextlib.redirect_stdout(stdout), contextlib.redirect_stderr(stderr):
            compile_page_pack(run_dir, sys.executable)
    except Exception as exc:  # pragma: no cover - caller asserts stdout/stderr on failures
        stderr.write(f"{type(exc).__name__}: {exc}")
        return subprocess.CompletedProcess(args, 1, stdout.getvalue(), stderr.getvalue())
    return subprocess.CompletedProcess(args, 0, stdout.getvalue(), stderr.getvalue())


def _chart_led_pack_without_fixed_role() -> dict:
    return {
        "schema_version": "banker_page_pack",
        "section_meta": {},
        "deck_storyline": "Exhibit-led page pack for a concise client-facing section.",
        "slides": [
            {
                "slide_no": 1,
                "banker_judgment": "The market signal is carried by comparable audited metrics rather than supporting bullets.",
                "page_argument": "A chart-led page can be page-worthy when the visual itself contains the main proof.",
                "selected_page_type": "chart_page",
                "deck_use": "可作标题",
                "headline": "Chart-led evidence can carry a page when the data is explicit",
                "main_message": "The visual is the primary evidence; body copy is optional when the exhibit is sufficiently structured.",
                "exhibit": {
                    "exhibit_type": "chart",
                    "why_this_exhibit": "Show a compact market scale trend.",
                    "data_or_evidence_inputs": ["MET-001"],
                    "visual_structure": "One chart with source-bound trend values.",
                },
                "body_blocks": [],
                "chart_data": {
                    "chart_type": "bar",
                    "title": "Indexed market scale",
                    "categories": ["2022", "2023", "2024"],
                    "series": [{"name": "Index", "values": [80, 95, 110]}],
                    "unit": "index",
                },
                "evidence_ids": [],
                "metric_ids": ["MET-001"],
                "project_relevance_note": "",
                "source_note": "Sources: metric audit table.",
            }
        ],
    }


def test_banker_page_pack_flags_sparse_page_without_blocking_llm_revision(tmp_path: Path) -> None:
    pack = {
        "schema_version": "banker_page_pack",
        "section_meta": {},
        "deck_storyline": "too short",
        "slides": [
            {
                "slide_no": idx,
                "fixed_page_role": FIXED_PAGE_ROLES[idx],
                "banker_judgment": "thin",
                "page_argument": "thin",
                "selected_page_type": "summary_page",
                "deck_use": "可作标题",
                "headline": "Thin page",
                "main_message": "thin",
                "exhibit": {"exhibit_type": "driver_cards", "data_or_evidence_inputs": ["EV-001"], "visual_structure": "thin"},
                "body_blocks": [],
                "evidence_ids": ["EV-001"],
                "metric_ids": [],
                "project_relevance_note": "",
                "source_note": "Source: fixture evidence database",
            }
            for idx in SLIDE_NUMBERS
        ],
    }
    path = tmp_path / "banker_page_pack.json"
    _write_json(path, pack)

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(path)])
    assert result.returncode == 0, result.stdout + result.stderr
    assert "page lacks substantive visible content" in result.stdout
    assert "LLM editorial prompt" in result.stdout
    assert "instead of filling template slots automatically" in result.stdout


def test_banker_page_pack_allows_exhibit_led_page_without_body_blocks(tmp_path: Path) -> None:
    pack = _chart_led_pack_without_fixed_role()
    pack["slides"][0]["fixed_page_role"] = FIXED_PAGE_ROLES[1]
    path = tmp_path / "banker_page_pack.json"
    _write_json(path, pack)

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(path)])
    assert result.returncode == 0, result.stdout
    assert "no body_blocks supplied" not in result.stdout


def test_banker_page_pack_allows_llm_visual_design_payload_without_body_blocks(tmp_path: Path) -> None:
    pack = _chart_led_pack_without_fixed_role()
    slide = pack["slides"][0]
    slide["selected_page_type"] = "freeform_page"
    slide["exhibit"] = {
        "exhibit_type": "custom_visual",
        "why_this_exhibit": "LLM chose a visual sequence rather than a fixed chart or table.",
    }
    slide["body_blocks"] = []
    slide["body_copy"] = {}
    slide["chart_data"] = {}
    slide["compare_table_data"] = {}
    slide["visual_design"] = {
        "layout": "Three source-labeled cards connected by a short flow arrow",
        "card_1": "Demand signal: category evidence and channel behavior",
        "card_2": "Competition signal: peer positioning and product economics",
        "card_3": "Transaction relevance: why the category lens matters for a pre-mandate discussion",
    }
    path = tmp_path / "banker_page_pack.json"
    _write_json(path, pack)

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(path)])

    assert result.returncode == 0, result.stdout + result.stderr
    assert "page lacks substantive visible content" not in result.stdout


def test_banker_page_pack_allows_metric_claim_led_page_without_body_blocks(tmp_path: Path) -> None:
    pack = _chart_led_pack_without_fixed_role()
    slide = pack["slides"][0]
    slide.pop("chart_data", None)
    slide.pop("compare_table_data", None)
    slide.pop("visual_design", None)
    slide["body_blocks"] = []
    slide["body_copy"] = {}
    slide["exhibit"] = {
        "exhibit_type": "kpi_cards",
        "why_this_exhibit": "Visible metric cards carry the proof instead of prose blocks.",
    }
    slide["visible_metric_claims"] = [
        {
            "metric_id": "MET-001",
            "display_text": "2024 market indicator: 110 index",
            "source_note": "Source: metric audit table.",
        }
    ]
    slide["key_data_audit"] = [
        {
            "metric_id": "MET-001",
            "indicator": "Indexed market scale",
            "value": "110",
            "source": "Metric audit table",
        }
    ]
    path = tmp_path / "banker_page_pack.json"
    _write_json(path, pack)

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(path)])

    assert result.returncode == 0, result.stdout + result.stderr
    assert "page lacks substantive visible content" not in result.stdout


def test_metric_claims_infer_data_page_type_when_style_guided() -> None:
    from renderer_compile_utils import infer_selected_page_type

    slide = {
        "slide_no": 7,
        "headline": "KPI cards carry the evidence without fixed template fields",
        "page_argument": "Metric cards can be the visible payload when the numbers are sourced.",
        "visible_metric_claims": [{"metric_id": "MET-001", "display_text": "2024 indicator: 110 index"}],
        "metric_ids": ["MET-001"],
    }

    assert infer_selected_page_type(slide, 7, {}) == "chart_page"


def test_style_guided_page_type_inference_is_content_led_not_slide_number() -> None:
    from renderer_compile_utils import infer_selected_page_type

    table_slide = {
        "headline": "Peer economics require a comparison table",
        "page_argument": "The page compares peers across price band, channel mix, and proof points.",
        "compare_table_data": {
            "headers": ["Price band", "Channel mix", "Evidence"],
            "rows": [{"label": "Brand A", "cells": ["Mass", "Douyin-led", "Platform ranking"]}],
        },
    }
    metric_slide = {
        "headline": "Sourced KPIs carry the market argument",
        "page_argument": "The page uses visible metric cards to support the category story.",
        "visible_metric_claims": [{"metric_id": "MET-001", "display_text": "2024 index: 110"}],
        "metric_ids": ["MET-001"],
    }

    assert infer_selected_page_type(table_slide, 2, {}) == "compare_table_page"
    assert infer_selected_page_type(table_slide, 7, {}) == "compare_table_page"
    assert infer_selected_page_type(metric_slide, 1, {}) == "chart_page"
    assert infer_selected_page_type(metric_slide, 5, {}) == "chart_page"


def test_banker_page_pack_fixed_page_role_optional_in_style_guided_mode(tmp_path: Path) -> None:
    pack = _chart_led_pack_without_fixed_role()
    path = tmp_path / "banker_page_pack.json"
    _write_json(path, pack)

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(path)])
    assert result.returncode == 0, result.stdout


def test_style_guided_compile_does_not_inject_registry_page_roles() -> None:
    from deck_blueprint_utils import normalize_deck_blueprint_for_page_plan
    from renderer_compile_utils import build_banker_page_contract, build_internal_deck_blueprint, build_renderer_spec_from_deck_blueprint

    pack = {
        "schema_version": "banker_page_pack",
        "slides": [
            {
                "page_argument": "Evidence-backed category framing should drive the first page.",
                "headline": "Category framing leads the pitch, not the bundled template role",
                "deck_use": "只用于正文",
                "body_blocks": [{"copy": "The visible point is authored by the LLM, not copied from the registry."}],
                "source_note": "Source: fixture",
            }
        ],
    }
    deck = build_internal_deck_blueprint(pack)
    contract = build_banker_page_contract(deck)
    renderer_spec = build_renderer_spec_from_deck_blueprint(deck, {}, contract)
    page_plan = normalize_deck_blueprint_for_page_plan(deck)

    assert deck["rendering_policy"]["template_contract_mode"] == "style_guided"
    assert deck["slides"][0]["fixed_page_role"] == ""
    assert contract["slides"][0]["page_role"] == ""
    assert renderer_spec["slides"][0]["fixed_page_role"] == ""
    assert renderer_spec["slides"][0]["slide_role"] == ""
    assert renderer_spec["slides"][0]["page_role"] == ""
    assert page_plan["slides"][0]["fixed_page_role"] == ""


def test_style_guided_explicit_body_copy_keys_are_not_template_slots() -> None:
    from renderer_compile_utils import build_banker_page_contract, build_internal_deck_blueprint, build_renderer_spec_from_deck_blueprint

    pack = {
        "schema_version": "banker_page_pack",
        "slides": [
            {
                "page_argument": "The body composition should follow the authored market point.",
                "headline": "Body copy keys should not force template placeholder mapping",
                "deck_use": "只用于正文",
                "body_copy": {
                    "main_body": "First client-facing market point from the LLM-authored page.",
                    "right_top": "Second client-facing proof point that should not depend on a template slot.",
                },
                "source_note": "Source: fixture",
            }
        ],
    }
    deck = build_internal_deck_blueprint(pack)
    contract = build_banker_page_contract(deck)
    renderer_spec = build_renderer_spec_from_deck_blueprint(deck, {}, contract)
    body_copy = renderer_spec["slides"][0]["body_copy"]

    assert body_copy == {
        "point_1": "First client-facing market point from the LLM-authored page.",
        "point_2": "Second client-facing proof point that should not depend on a template slot.",
    }
    assert "main_body" not in body_copy
    assert "right_top" not in body_copy


def test_strict_layout_compile_keeps_registry_page_roles() -> None:
    from deck_blueprint_utils import normalize_deck_blueprint_for_page_plan
    from renderer_compile_utils import build_banker_page_contract, build_internal_deck_blueprint

    pack = {
        "schema_version": "banker_page_pack",
        "rendering_policy": {"template_contract_mode": "strict_layout"},
        "slides": [
            {
                "slide_no": 1,
                "page_argument": "Strict layout intentionally keeps the template registry lineage.",
                "headline": "Strict-layout registry role remains available",
                "deck_use": "只用于正文",
                "body_blocks": [{"copy": "Strict layout is the explicit placeholder-level mode."}],
                "source_note": "Source: fixture",
            }
        ],
    }

    deck = build_internal_deck_blueprint(pack)
    contract = build_banker_page_contract(deck)
    page_plan = normalize_deck_blueprint_for_page_plan(deck)

    assert deck["slides"][0]["fixed_page_role"] == FIXED_PAGE_ROLES[1]
    assert contract["slides"][0]["page_role"] == FIXED_PAGE_ROLES[1]
    assert page_plan["slides"][0]["fixed_page_role"] == FIXED_PAGE_ROLES[1]


def test_strict_layout_explicit_body_copy_preserves_template_slots() -> None:
    from renderer_compile_utils import build_banker_page_contract, build_internal_deck_blueprint, build_renderer_spec_from_deck_blueprint

    pack = {
        "schema_version": "banker_page_pack",
        "rendering_policy": {"template_contract_mode": "strict_layout"},
        "slides": [
            {
                "slide_no": 1,
                "selected_page_type": "summary_page",
                "page_argument": "Strict layout intentionally uses template placeholder fields.",
                "headline": "Strict layout preserves placeholder keys",
                "deck_use": "只用于正文",
                "body_copy": {"main_body": "Placeholder-level body copy remains available only in strict layout."},
                "source_note": "Source: fixture",
            }
        ],
    }
    template_registry = {
        "slides": [
            {
                "slide_no": 1,
                "variants": {
                    "summary_page": {"strict_layout_body_fields": ["main_body"]},
                },
            }
        ]
    }

    deck = build_internal_deck_blueprint(pack)
    contract = build_banker_page_contract(deck)
    renderer_spec = build_renderer_spec_from_deck_blueprint(deck, template_registry, contract)

    assert renderer_spec["slides"][0]["body_copy"] == {
        "main_body": "Placeholder-level body copy remains available only in strict layout."
    }


def test_style_guided_long_headline_is_advisory_not_shorten_blocker(tmp_path: Path) -> None:
    pack = {
        "schema_version": "banker_page_pack",
        "deliverable_readiness": {
            "business_action": "client_ready",
            "decision_note": "Fixture is intentionally long to test advisory line-fit behavior.",
        },
        "slides": [
            {
                "page_argument": "A compact first page makes the long-title page second by array order.",
                "headline": "First page keeps the section sequence clear",
                "main_message": "Style-guided mode uses page order rather than forcing template slide numbers.",
                "deck_use": "只用于正文",
                "body_blocks": [{"copy": "This page exists so the next page reaches the slide 2 text-fit hint by array order."}],
                "source_note": "Source: fixture",
            },
            {
                "selected_page_type": "chart_page",
                "page_argument": "The page uses a deliberately long headline to verify that style-guided mode preserves the argument and treats line fit as an editorial choice.",
                "headline": (
                    "A deliberately detailed market headline should trigger a line-fit warning without forcing the LLM "
                    "to delete the banker point before seeing the rendered page"
                ),
                "main_message": "The visual carries the point and the title can be rewritten, split, or handled through layout if the rendered page looks crowded.",
                "deck_use": "可作标题",
                "chart_data": {
                    "chart_type": "bar",
                    "title": "Indexed market signal",
                    "categories": ["2022", "2023", "2024"],
                    "series": [{"name": "Index", "values": [80, 95, 112]}],
                    "unit": "index",
                },
                "source_note": "Source: fixture",
            }
        ],
    }
    path = tmp_path / "banker_page_pack.json"
    _write_json(path, pack)
    _write_json(tmp_path / "template_registry.json", {"schema_version": "template_registry_v1", "slides": []})

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(path)])

    assert result.returncode == 0, result.stdout
    assert "line-fit advisory" in result.stdout
    assert "style-guided mode, preserve the page argument" in result.stdout
    assert "shorten before render" not in result.stdout


def test_banker_page_pack_fixed_page_role_required_in_strict_layout_mode(tmp_path: Path) -> None:
    (tmp_path / "artifacts").mkdir(parents=True)
    _write_json(
        tmp_path / "artifacts" / "rendering_policy.json",
        {"schema_version": "rendering_policy_v1", "template_contract_mode": "strict_layout"},
    )
    pack = _chart_led_pack_without_fixed_role()
    path = tmp_path / "banker_page_pack.json"
    _write_json(path, pack)

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(path)])
    assert result.returncode != 0
    assert "fixed_page_role is required" in result.stdout


def test_banker_page_pack_leaves_target_drift_to_llm_qc(tmp_path: Path) -> None:
    slides = []
    for idx in SLIDE_NUMBERS:
        slides.append(
            {
                "slide_no": idx,
                "fixed_page_role": FIXED_PAGE_ROLES[idx],
                "banker_judgment": "Industry judgment with source-backed market mechanism.",
                "page_argument": "Industry page argument with evidence and transaction relevance.",
                "selected_page_type": "summary_page",
                "deck_use": "可作标题",
                "headline": "标的交易故事强",
                "main_message": "Industry message.",
                "exhibit": {
                    "exhibit_type": "driver_cards",
                    "why_this_exhibit": "Shows market mechanisms.",
                    "data_or_evidence_inputs": ["EV-001"],
                    "visual_structure": "Four industry cards.",
                    "density_target": "Dense.",
                },
                "body_blocks": [
                    {
                        "role": f"point_{n}",
                        "copy": "Industry mechanism with evidence and transaction relevance.",
                        "evidence_ids": ["EV-001"],
                        "metric_ids": [],
                    }
                    for n in range(4)
                ],
                "evidence_ids": ["EV-001"],
                "metric_ids": [],
                "project_relevance_note": "",
                "source_note": "Source: fixture evidence database",
            }
        )
    pack = {
        "schema_version": "banker_page_pack",
        "section_meta": {},
        "deck_storyline": "Industry-led storyline.",
        "deliverable_readiness": {
            "business_action": "client_ready",
            "decision_note": "Fixture.",
        },
        "slides": slides,
    }
    path = tmp_path / "banker_page_pack.json"
    _write_json(path, pack)

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(path)])
    assert result.returncode == 0, result.stdout + result.stderr
    assert "headline contains target/project terms" not in result.stdout
    qc_text = (SKILL_DIR / "references/qc.md").read_text(encoding="utf-8")
    assert "project-context drift" in qc_text


def test_banker_page_pack_warns_on_internal_working_paper_language_without_blocking(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][0]["body_blocks"][0]["copy"] = "工作市场：样例行业；第一步是证明赛道边界清楚。"
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert result.returncode == 0, result.stdout + result.stderr
    assert "client-visible copy may contain internal workpaper language" in result.stdout
    assert "body_blocks[1].copy" in result.stdout
    assert "market-boundary slot label" in result.stdout
    assert "工作市场" not in result.stdout
    validator_text = (SKILL_DIR / "scripts/qc/validate_artifact.py").read_text(encoding="utf-8")
    assert "CLIENT_VISIBLE_EDITORIAL_HINT_TERMS" in validator_text
    generation_text = (SKILL_DIR / "references/generation.md").read_text(encoding="utf-8")
    assert "Do not write visible slide copy that talks about how to build the pitchbook" in generation_text
    assert "Avoid workpaper labels in visible Chinese or English copy" in generation_text


def test_banker_page_pack_allows_named_source_limitations_when_client_facing(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][0]["body_blocks"][0]["copy"] = (
        "Source limitation: category-level GMV is directionally useful, but platform data should not be read as full-market share."
    )
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert result.returncode == 0, result.stdout + result.stderr
    assert "source limitation" not in result.stdout.lower()


def test_banker_page_pack_warns_on_internal_point_copy_without_blocking(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    first_block = pack["slides"][0]["body_blocks"][0]
    first_block.pop("copy", None)
    first_block["point"] = "后续验证点：这里仍是内部工作底稿语言。"
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert result.returncode == 0, result.stdout + result.stderr
    assert "client-visible copy may contain internal workpaper language" in result.stdout
    assert "body_blocks[1].point" in result.stdout
    assert "research workflow label" in result.stdout
    assert "后续验证点" not in result.stdout


def test_banker_page_pack_warns_on_internal_metric_claim_copy_without_blocking(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][0]["visible_metric_claims"] = [
        {"claim": "客户关注点：该指标仍是内部问题语言。", "metric_id": "MET-001"}
    ]
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert result.returncode == 0, result.stdout + result.stderr
    assert "client-visible copy may contain internal workpaper language" in result.stdout
    assert "visible_metric_claims[1].claim" in result.stdout
    assert "question-bucket label" in result.stdout
    assert "客户关注点" not in result.stdout


def test_banker_page_pack_warns_on_visible_readiness_or_research_loop_language_without_blocking(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][0]["body_blocks"][0]["copy"] = (
        "NOT_CLIENT_READY: this page needs a targeted research request before client-ready delivery."
    )
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert result.returncode == 0, result.stdout + result.stderr
    assert "client-visible copy may contain internal workpaper language" in result.stdout
    assert "body_blocks[1].copy" in result.stdout
    assert "NOT_CLIENT_READY" not in result.stdout
    assert "research request" not in result.stdout
    assert "delivery-status label" in result.stdout
    assert "research workflow label" in result.stdout


def test_filled_ppt_warns_on_internal_working_paper_language_without_blocking(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    import validate_artifact as validator

    pptx_path = tmp_path / "industry_section_filled_clean.pptx"
    with zipfile.ZipFile(pptx_path, "w") as archive:
        archive.writestr("[Content_Types].xml", "<Types></Types>")
    monkeypatch.setattr(
        validator,
        "_ppt_slide_texts",
        lambda _path: ["后续验证点：口径越窄，越接近目标项目的真实赛道。"],
    )

    errors, warnings = validator.validate_artifact("filled_ppt", tmp_path)

    assert errors == []
    assert any("rendered PPT text may contain internal workpaper language" in warning for warning in warnings)
    assert all("后续验证点" not in warning for warning in warnings)


def test_filled_ppt_accepts_latest_final_ppt_marker_for_direct_composition(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    import validate_artifact as validator

    pptx_path = tmp_path / "client_style_direct_composition.pptx"
    with zipfile.ZipFile(pptx_path, "w") as archive:
        archive.writestr("[Content_Types].xml", "<Types></Types>")
    (tmp_path / "LATEST_FINAL_PPT.txt").write_text(f"{pptx_path.name}\n", encoding="utf-8")
    monkeypatch.setattr(validator, "_ppt_slide_texts", lambda _path: ["面部底妆品类判断与交易相关性"])

    errors, warnings = validator.validate_artifact("filled_ppt", tmp_path)

    assert errors == []
    assert not any("missing filled PPT output" in item for item in errors + warnings)


def test_banker_page_pack_leaves_subject_mix_to_llm_qc(tmp_path: Path) -> None:
    slides = []
    for idx in SLIDE_NUMBERS:
        slides.append(
            {
                "slide_no": idx,
                "fixed_page_role": FIXED_PAGE_ROLES[idx],
                "banker_judgment": "Industry judgment with source-backed market mechanism.",
                "page_argument": "Industry page argument with evidence and transaction relevance.",
                "selected_page_type": "summary_page",
                "deck_use": "可作标题",
                "headline": "Industry structure leads",
                "main_message": "Industry message.",
                "exhibit": {
                    "exhibit_type": "driver_cards",
                    "why_this_exhibit": "Shows market mechanisms.",
                    "data_or_evidence_inputs": ["EV-001"],
                    "visual_structure": "Four industry cards.",
                    "density_target": "Dense.",
                },
                "body_blocks": [
                    {
                        "role": f"point_{n}",
                        "copy": "Industry mechanism with evidence and transaction relevance.",
                        "evidence_ids": ["EV-001"],
                        "metric_ids": [],
                    }
                    for n in range(4)
                ],
                "evidence_ids": ["EV-001"],
                "metric_ids": [],
                "project_relevance_note": "Short project bridge.",
                "source_note": "Source: fixture evidence database",
            }
        )
    pack = {
        "schema_version": "banker_page_pack",
        "section_meta": {},
        "deck_storyline": "Industry-led storyline.",
        "deliverable_readiness": {
            "business_action": "client_ready",
            "decision_note": "Fixture.",
        },
        "slides": slides,
    }
    path = tmp_path / "banker_page_pack.json"
    _write_json(path, pack)

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(path)])
    assert result.returncode == 0, result.stdout + result.stderr
    assert "fewer industry-primary pages than the advisory target" not in result.stdout
    assert "more project_relevance_note pages than the advisory target" not in result.stdout
    assert "references/content-quality.md" in (SKILL_DIR / "references/qc.md").read_text(encoding="utf-8")


def test_banker_page_pack_leaves_mixed_axis_units_to_llm_qc(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][0]["chart_data"]["source_rows"] = [
        {"label": "market size", "value": 100, "unit": "RMB bn", "metric_id": "MET-001"},
        {"label": "target sales", "value": 470, "unit": "万件", "metric_id": "MET-002"},
    ]
    path = tmp_path / "banker_page_pack.json"
    _write_json(path, pack)

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(path)])

    assert result.returncode == 0, result.stdout + result.stderr
    assert "chart_data mixes units on one chart axis" not in result.stdout
    assert "mixed units" in (SKILL_DIR / "references/qc.md").read_text(encoding="utf-8")


def test_banker_page_pack_allows_omitted_deck_use_with_conservative_default(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][0].pop("deck_use", None)
    pack["slides"][0].pop("allowed_deck_usage", None)
    path = tmp_path / "banker_page_pack.json"
    _write_json(path, pack)

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(path)])

    assert result.returncode == 0, result.stdout + result.stderr
    assert "deck_use omitted" not in result.stdout

    _write_json(tmp_path / "banker_page_pack.json", pack)
    _write_json(tmp_path / "template_registry.json", template_registry)
    compile_result = _compile_page_pack(tmp_path)

    assert compile_result.returncode == 0, compile_result.stdout + compile_result.stderr
    contract = json.loads((tmp_path / "page_evidence_contract.json").read_text(encoding="utf-8"))
    first_slide = contract["slides"][0]
    assert first_slide["allowed_deck_usage"] == "supporting_context"
    assert first_slide["headline_allowed"] is False
    assert first_slide["body_evidence_ids"]


def test_banker_page_pack_accepts_natural_deck_use_language(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][0].pop("allowed_deck_usage", None)
    pack["slides"][0]["deck_use"] = "可作标题"
    pack["slides"][1].pop("allowed_deck_usage", None)
    pack["slides"][1]["deck_use"] = "只用于正文"
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert validation.returncode == 0, validation.stdout + validation.stderr
    assert "non-standard" not in validation.stdout

    compile_result = _compile_page_pack(tmp_path)

    assert compile_result.returncode == 0, compile_result.stdout + compile_result.stderr
    contract = json.loads((tmp_path / "page_evidence_contract.json").read_text(encoding="utf-8"))
    assert contract["slides"][0]["allowed_deck_usage"] == "headline_allowed"
    assert contract["slides"][0]["headline_allowed"] is True
    assert contract["slides"][1]["allowed_deck_usage"] == "body_only"
    assert contract["slides"][1]["headline_allowed"] is False
    assert contract["slides"][1]["downstream_permission"]["body_copy_allowed"] is True


def test_banker_page_pack_unrecognized_deck_use_stays_natural_and_defaults_conservative(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][0].pop("allowed_deck_usage", None)
    pack["slides"][0]["deck_use"] = "这页适合作为市场背景铺垫，由合伙人决定是否放在正文"
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert validation.returncode == 0, validation.stdout + validation.stderr
    assert "non-standard" not in validation.stdout
    assert "supporting_context" not in validation.stdout

    compile_result = _compile_page_pack(tmp_path)

    assert compile_result.returncode == 0, compile_result.stdout + compile_result.stderr
    contract = json.loads((tmp_path / "page_evidence_contract.json").read_text(encoding="utf-8"))
    assert contract["slides"][0]["allowed_deck_usage"] == "supporting_context"
    assert contract["slides"][0]["headline_allowed"] is False


def test_style_guided_compile_skips_pages_marked_not_for_deck(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    skipped_headline = pack["slides"][0]["headline"]
    kept_headline = pack["slides"][1]["headline"]
    pack["slides"][0].pop("allowed_deck_usage", None)
    pack["slides"][0]["deck_use"] = "不可用于页面"
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert validation.returncode == 0, validation.stdout + validation.stderr
    assert "style-guided structured render will skip these pages" in validation.stdout

    compile_result = _compile_page_pack(tmp_path)

    assert compile_result.returncode == 0, compile_result.stdout + compile_result.stderr
    blueprint = json.loads((tmp_path / "deck_blueprint.json").read_text(encoding="utf-8"))
    contract = json.loads((tmp_path / "page_evidence_contract.json").read_text(encoding="utf-8"))
    renderer = json.loads((tmp_path / "renderer_spec.json").read_text(encoding="utf-8"))
    headlines = [slide["headline"] for slide in blueprint["slides"]]
    assert skipped_headline not in headlines
    assert headlines[0] == kept_headline
    assert blueprint["slides"][0]["slide_no"] == 1
    assert contract["slides"][0]["slide_no"] == 1
    assert renderer["slides"][0]["slide_no"] == 1


def test_main_message_is_optional_and_compiler_preserves_omission(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][0].pop("main_message", None)
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert validation.returncode == 0, validation.stdout + validation.stderr
    assert "main_message is required" not in validation.stdout

    compile_result = _compile_page_pack(tmp_path)

    assert compile_result.returncode == 0, compile_result.stdout + compile_result.stderr
    compiled_deck = json.loads((tmp_path / "deck_blueprint.json").read_text(encoding="utf-8"))
    renderer = json.loads((tmp_path / "renderer_spec.json").read_text(encoding="utf-8"))
    assert compiled_deck["slides"][0]["main_message"] == ""
    assert renderer["slides"][0]["main_message"] == ""


def test_internal_banker_page_id_is_optional_and_normalized(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][0]["banker_page_id"] = "first-page"
    pack["slides"][1].pop("banker_page_id", None)
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert validation.returncode == 0, validation.stdout + validation.stderr
    assert "will be normalized to BP-001" in validation.stdout

    compile_result = _compile_page_pack(tmp_path)

    assert compile_result.returncode == 0, compile_result.stdout + compile_result.stderr
    deck = json.loads((tmp_path / "deck_blueprint.json").read_text(encoding="utf-8"))
    contract = json.loads((tmp_path / "page_evidence_contract.json").read_text(encoding="utf-8"))
    assert deck["slides"][0]["banker_page_id"] == "BP-001"
    assert deck["slides"][1]["banker_page_id"] == "BP-002"
    assert contract["slides"][0]["banker_page_id"] == "BP-001"
    assert contract["slides"][1]["banker_page_id"] == "BP-002"


def test_legacy_transaction_readthrough_is_rejected_instead_of_mapped(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][0]["project_relevance_note"] = ""
    pack["slides"][0]["transaction_readthrough"] = "Pre-mandate transaction relevance stays as a short client-facing bridge."
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert validation.returncode != 0
    assert "transaction_readthrough is a legacy internal field" in validation.stdout
    assert "project_relevance_note" in validation.stdout


def test_project_relevance_note_remains_current_project_bridge_field(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][0]["project_relevance_note"] = "Pre-mandate transaction relevance stays as a short client-facing bridge."
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert validation.returncode == 0, validation.stdout + validation.stderr

    compile_result = _compile_page_pack(tmp_path)

    assert compile_result.returncode == 0, compile_result.stdout + compile_result.stderr
    deck = json.loads((tmp_path / "deck_blueprint.json").read_text(encoding="utf-8"))
    assert deck["slides"][0]["project_relevance_note"] == "Pre-mandate transaction relevance stays as a short client-facing bridge."


def test_visible_metric_claim_usage_type_is_optional(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][0]["visible_metric_claims"][0].pop("usage_type", None)
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert validation.returncode == 0, validation.stdout + validation.stderr

    compile_result = _compile_page_pack(tmp_path)

    assert compile_result.returncode == 0, compile_result.stdout + compile_result.stderr
    renderer = json.loads((tmp_path / "renderer_spec.json").read_text(encoding="utf-8"))
    assert renderer["slides"][0]["visible_metric_claims"][0]["usage_type"] == "visible_metric"


def test_visible_metric_claim_accepts_natural_metric_aliases(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][0]["visible_metric_claims"] = [
        {
            "where": "headline proof",
            "claim": "Market indicator reaches 110 index points in 2024.",
            "metric_id": "MET-001",
        }
    ]
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])
    assert validation.returncode == 0, validation.stdout + validation.stderr

    compile_result = _compile_page_pack(tmp_path)

    assert compile_result.returncode == 0, compile_result.stdout + compile_result.stderr
    renderer = json.loads((tmp_path / "renderer_spec.json").read_text(encoding="utf-8"))
    claim = renderer["slides"][0]["visible_metric_claims"][0]
    assert claim["location"] == "headline proof"
    assert claim["display_text"] == "Market indicator reaches 110 index points in 2024."
    assert claim["metric_ids"] == ["MET-001"]
    assert claim["usage_type"] == "visible_metric"


def test_compiler_uses_allowed_deck_usage_for_contract_permission(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][0]["allowed_deck_usage"] = "body_only"
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    result = _compile_page_pack(tmp_path)

    assert result.returncode == 0, result.stdout + result.stderr
    contract = json.loads((tmp_path / "page_evidence_contract.json").read_text(encoding="utf-8"))
    blueprint = json.loads((tmp_path / "deck_blueprint.json").read_text(encoding="utf-8"))
    first_slide = contract["slides"][0]
    assert first_slide["allowed_deck_usage"] == "body_only"
    assert "evidence_status" not in first_slide
    assert first_slide["headline_allowed"] is False
    assert first_slide["downstream_permission"]["body_copy_allowed"] is True
    assert "body_only" in json.dumps(blueprint, ensure_ascii=False)


def test_banker_page_pack_validates_and_compiles(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])
    assert validation.returncode == 0, validation.stdout + validation.stderr

    deck_blueprint = tmp_path / "deck_blueprint.json"
    page_contract = tmp_path / "page_evidence_contract.json"
    renderer_spec = tmp_path / "renderer_spec.json"
    result = _compile_page_pack(tmp_path)
    assert result.returncode == 0, result.stdout + result.stderr
    compiled_deck = json.loads(deck_blueprint.read_text(encoding="utf-8"))
    assert "authoring_status" not in compiled_deck
    assert "page_" + "argument_ids" not in json.dumps(compiled_deck)
    assert "issue_" + "analysis_ids" not in json.dumps(compiled_deck)
    contract = json.loads(page_contract.read_text(encoding="utf-8"))
    assert all(slide.get("banker_page_id", "").startswith("BP-") for slide in contract["slides"])
    renderer = json.loads(renderer_spec.read_text(encoding="utf-8"))
    assert len(renderer["slides"]) == 8
    assert all(slide.get("source_note") for slide in renderer["slides"])
    assert all("pitch_relevance" not in slide for slide in renderer["slides"])


def test_style_guided_page_pack_can_omit_selected_page_type(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    for slide in pack["slides"]:
        slide.pop("selected_page_type", None)
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert validation.returncode == 0, validation.stdout + validation.stderr
    assert "selected_page_type omitted" not in validation.stdout

    result = _compile_page_pack(tmp_path)

    assert result.returncode == 0, result.stdout + result.stderr
    renderer = json.loads((tmp_path / "renderer_spec.json").read_text(encoding="utf-8"))
    assert renderer["slides"][0]["selected_page_type"] == "chart_page"
    assert renderer["slides"][5]["selected_page_type"] == "compare_table_page"


def test_page_question_is_legacy_input_context_only(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    expected_headline = pack["slides"][0]["headline"]
    for slide in pack["slides"]:
        slide.pop("page_question", None)
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])
    assert validation.returncode == 0, validation.stdout + validation.stderr

    result = _compile_page_pack(tmp_path)

    assert result.returncode == 0, result.stdout + result.stderr
    blueprint = json.loads((tmp_path / "deck_blueprint.json").read_text(encoding="utf-8"))
    contract = json.loads((tmp_path / "page_evidence_contract.json").read_text(encoding="utf-8"))
    assert "page_question" not in blueprint["slides"][0]
    assert "page_question" not in contract["slides"][0]
    assert contract["slides"][0]["headline_claim"] == expected_headline


def test_source_note_is_optional_when_evidence_ids_are_bound(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    for slide in pack["slides"]:
        slide.pop("source_note", None)
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert validation.returncode == 0, validation.stdout + validation.stderr
    assert "source_note omitted" not in validation.stdout

    result = _compile_page_pack(tmp_path)

    assert result.returncode == 0, result.stdout + result.stderr
    renderer = json.loads((tmp_path / "renderer_spec.json").read_text(encoding="utf-8"))
    assert renderer["slides"][0]["source_note"] == ""


def test_source_note_is_optional_when_metric_ids_are_bound(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    first_slide = pack["slides"][0]
    first_slide.pop("source_note", None)
    first_slide["evidence_ids"] = []
    first_slide["metric_ids"] = ["MET-001"]
    first_slide["visible_metric_claims"] = [
        {"display_text": "2024 sample metric: 110", "metric_id": "MET-001"}
    ]
    for block in first_slide.get("body_blocks", []):
        block["evidence_ids"] = []
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert validation.returncode == 0, validation.stdout + validation.stderr
    assert "source_note omitted" not in validation.stdout


def test_deliverable_readiness_uses_substantive_judgment_fields(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    assert "decision_status" not in pack["deliverable_readiness"]
    assert "decision_owner" not in pack["deliverable_readiness"]
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])
    assert validation.returncode == 0, validation.stdout + validation.stderr

    result = _compile_page_pack(tmp_path)

    assert result.returncode == 0, result.stdout + result.stderr
    compiled = json.loads((tmp_path / "deck_blueprint.json").read_text(encoding="utf-8"))
    assert len(compiled["slides"]) == len(pack["slides"])


def test_section_meta_is_optional_run_metadata(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack.pop("section_meta", None)
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])
    assert validation.returncode == 0, validation.stdout + validation.stderr

    result = _compile_page_pack(tmp_path)

    assert result.returncode == 0, result.stdout + result.stderr
    renderer = json.loads((tmp_path / "renderer_spec.json").read_text(encoding="utf-8"))
    assert renderer["section_meta"] == {}


def test_deck_storyline_is_optional_editorial_overview(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack.pop("deck_storyline", None)
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])
    assert validation.returncode == 0, validation.stdout + validation.stderr

    result = _compile_page_pack(tmp_path)

    assert result.returncode == 0, result.stdout + result.stderr
    compiled = json.loads((tmp_path / "deck_blueprint.json").read_text(encoding="utf-8"))
    assert compiled["deck_storyline"] == ""
    assert len(compiled["slides"]) == len(pack["slides"])


def test_banker_judgment_is_optional_when_page_argument_is_clear(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    expected_argument = pack["slides"][0]["page_argument"]
    for slide in pack["slides"]:
        slide.pop("banker_judgment", None)
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])
    assert validation.returncode == 0, validation.stdout + validation.stderr

    result = _compile_page_pack(tmp_path)

    assert result.returncode == 0, result.stdout + result.stderr
    compiled = json.loads((tmp_path / "deck_blueprint.json").read_text(encoding="utf-8"))
    assert compiled["slides"][0]["page_thesis"] == expected_argument


def test_banker_page_pack_accepts_natural_thesis_and_title_fields(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    first_slide = pack["slides"][0]
    expected_thesis = "行业需求正在从泛流量增长转向有证据支撑的细分场景增长。"
    expected_title = "细分场景增长正在重塑底妆品牌的交易叙事"
    first_slide.pop("page_argument", None)
    first_slide.pop("headline", None)
    first_slide.pop("main_message", None)
    first_slide["page_thesis"] = expected_thesis
    first_slide["title"] = expected_title
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])
    assert validation.returncode == 0, validation.stdout + validation.stderr
    assert "page_argument is required" not in validation.stdout
    assert "headline is required" not in validation.stdout

    result = _compile_page_pack(tmp_path)

    assert result.returncode == 0, result.stdout + result.stderr
    compiled = json.loads((tmp_path / "deck_blueprint.json").read_text(encoding="utf-8"))
    renderer = json.loads((tmp_path / "renderer_spec.json").read_text(encoding="utf-8"))
    assert compiled["slides"][0]["page_argument"] == expected_thesis
    assert compiled["slides"][0]["headline"] == expected_title
    assert compiled["slides"][0]["main_message"] == ""
    assert renderer["slides"][0]["headline"] == expected_title
    assert renderer["slides"][0]["main_message"] == ""


def test_page_argument_derives_visible_headline_when_headline_omitted(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    first_slide = pack["slides"][0]
    expected_argument = "底妆品类增长正在从泛流量转向肤质适配、复购和内容转化的组合能力。"
    first_slide["page_argument"] = expected_argument
    first_slide.pop("headline", None)
    first_slide.pop("title", None)
    first_slide.pop("slide_title", None)
    first_slide.pop("page_title", None)
    first_slide.pop("main_message", None)
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert validation.returncode == 0, validation.stdout + validation.stderr
    assert "no headline/title was found" not in validation.stdout

    result = _compile_page_pack(tmp_path)

    assert result.returncode == 0, result.stdout + result.stderr
    compiled = json.loads((tmp_path / "deck_blueprint.json").read_text(encoding="utf-8"))
    renderer = json.loads((tmp_path / "renderer_spec.json").read_text(encoding="utf-8"))
    assert compiled["slides"][0]["headline"] == expected_argument
    assert compiled["slides"][0]["main_message"] == ""
    assert renderer["slides"][0]["headline"] == expected_argument
    assert renderer["slides"][0]["main_message"] == ""


def test_page_primary_subject_is_legacy_input_context_only(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    for slide in pack["slides"]:
        slide["page_primary_subject"] = "target_context"
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])
    assert validation.returncode == 0, validation.stdout + validation.stderr

    result = _compile_page_pack(tmp_path)

    assert result.returncode == 0, result.stdout + result.stderr
    compiled = json.loads((tmp_path / "deck_blueprint.json").read_text(encoding="utf-8"))
    assert all("page_primary_subject" not in slide for slide in compiled["slides"])


def test_exhibit_metadata_is_optional_when_visible_payload_exists(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    for slide in pack["slides"]:
        slide.pop("exhibit", None)
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])
    assert validation.returncode == 0, validation.stdout + validation.stderr

    result = _compile_page_pack(tmp_path)

    assert result.returncode == 0, result.stdout + result.stderr
    compiled = json.loads((tmp_path / "deck_blueprint.json").read_text(encoding="utf-8"))
    assert compiled["slides"][0]["visual_intent"] == "Slide 1 metric trend"
    assert compiled["slides"][0]["exhibit"] == {}


def test_page_level_evidence_metric_rollups_and_body_roles_are_optional(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    for slide in pack["slides"]:
        slide.pop("evidence_ids", None)
        slide.pop("metric_ids", None)
        slide.pop("body_copy", None)
        for block in slide.get("body_blocks", []):
            block.pop("role", None)
            block.pop("evidence_ids", None)
            block.pop("metric_ids", None)
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])
    assert validation.returncode == 0, validation.stdout + validation.stderr

    result = _compile_page_pack(tmp_path)

    assert result.returncode == 0, result.stdout + result.stderr
    renderer = json.loads((tmp_path / "renderer_spec.json").read_text(encoding="utf-8"))
    assert renderer["slides"][0]["body_copy"]["point_1"]
    assert renderer["slides"][0]["visible_metric_claims"]


def test_body_blocks_accept_point_and_text_synonyms(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    slide = pack["slides"][0]
    slide.pop("body_copy", None)
    slide["body_blocks"] = [
        {
            "role": "mechanism",
            "point": "Category demand is shifting from broad cosmetic discovery to purpose-led base-makeup use cases.",
            "evidence_ids": ["EV-001"],
        },
        {
            "role": "proof",
            "text": "Channel evidence and peer behavior should be translated into a market-facing proof point.",
            "evidence_ids": ["EV-002"],
        },
    ]
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])
    assert validation.returncode == 0, validation.stdout + validation.stderr

    result = _compile_page_pack(tmp_path)

    assert result.returncode == 0, result.stdout + result.stderr
    renderer = json.loads((tmp_path / "renderer_spec.json").read_text(encoding="utf-8"))
    first_slide = renderer["slides"][0]
    assert "purpose-led base-makeup use cases" in first_slide["body_copy"]["point_1"]
    assert "market-facing proof point" in first_slide["body_copy"]["point_2"]


def test_body_blocks_accept_singular_evidence_metric_aliases(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    slide = pack["slides"][0]
    slide.pop("body_copy", None)
    slide["body_blocks"] = [
        {
            "role": "market_mechanism",
            "point": "Base-makeup demand should be framed through usage occasion and channel conversion.",
            "evidence_id": "EV-005",
            "metric_id": "MET-004",
        },
        {
            "role": "supporting_evidence",
            "text": "Peer and platform behavior provide the supporting evidence for the page claim.",
            "evidence_id": "EV-006",
        },
    ]
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])
    assert validation.returncode == 0, validation.stdout + validation.stderr

    compile_result = _compile_page_pack(tmp_path)
    assert compile_result.returncode == 0, compile_result.stdout + compile_result.stderr

    contract = json.loads((tmp_path / "page_evidence_contract.json").read_text(encoding="utf-8"))
    first_slide = contract["slides"][0]
    assert "EV-005" in first_slide["body_evidence_ids"]
    assert "EV-006" in first_slide["body_evidence_ids"]
    assert "MET-004" in first_slide["body_metric_ids"]


def test_style_guided_page_pack_can_omit_slide_no_and_use_array_order(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    for slide in pack["slides"]:
        slide.pop("slide_no", None)
        slide.pop("selected_page_type", None)
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert validation.returncode == 0, validation.stdout + validation.stderr
    assert "omits slide_no" in validation.stdout

    result = _compile_page_pack(tmp_path)

    assert result.returncode == 0, result.stdout + result.stderr
    renderer = json.loads((tmp_path / "renderer_spec.json").read_text(encoding="utf-8"))
    assert [slide["slide_no"] for slide in renderer["slides"]] == [1, 2, 3, 4, 5, 6, 7, 8]


def test_style_guided_duplicate_slide_no_hints_normalize_to_array_order(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    for slide in pack["slides"]:
        slide["slide_no"] = 1
        slide["banker_page_id"] = "BP-001"
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert validation.returncode == 0, validation.stdout + validation.stderr
    assert "duplicate slide_no hints" in validation.stdout
    assert "style-guided structured render will use array order" in validation.stdout

    compile_result = _compile_page_pack(tmp_path)

    assert compile_result.returncode == 0, compile_result.stdout + compile_result.stderr
    compiled = json.loads((tmp_path / "deck_blueprint.json").read_text(encoding="utf-8"))
    contract = json.loads((tmp_path / "page_evidence_contract.json").read_text(encoding="utf-8"))
    renderer = json.loads((tmp_path / "renderer_spec.json").read_text(encoding="utf-8"))
    assert [slide["slide_no"] for slide in compiled["slides"]] == [1, 2, 3, 4, 5, 6, 7, 8]
    assert [slide["banker_page_id"] for slide in compiled["slides"]] == [f"BP-{idx:03d}" for idx in range(1, 9)]
    assert [slide["slide_no"] for slide in contract["slides"]] == [1, 2, 3, 4, 5, 6, 7, 8]
    assert [slide["slide_no"] for slide in renderer["slides"]] == [1, 2, 3, 4, 5, 6, 7, 8]


def test_strict_layout_still_requires_selected_page_type(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][0].pop("selected_page_type", None)
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    _write_json(tmp_path / "artifacts/rendering_policy.json", {"template_contract_mode": "strict_layout"})
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert result.returncode != 0
    assert "selected_page_type is required" in result.stdout


def test_strict_layout_still_requires_slide_no(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][0].pop("slide_no", None)
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    _write_json(tmp_path / "artifacts/rendering_policy.json", {"template_contract_mode": "strict_layout"})
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert result.returncode != 0
    assert "missing positive slide_no" in result.stdout


def test_style_guided_compile_uses_llm_selected_page_count(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"] = pack["slides"][:5]
    for idx, slide in enumerate(pack["slides"], start=1):
        slide["slide_no"] = idx
        slide["banker_page_id"] = f"BP-{idx:03d}"
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])
    assert validation.returncode == 0, validation.stdout + validation.stderr

    result = _compile_page_pack(tmp_path)

    assert result.returncode == 0, result.stdout + result.stderr
    renderer = json.loads((tmp_path / "renderer_spec.json").read_text(encoding="utf-8"))
    assert [slide["slide_no"] for slide in renderer["slides"]] == [1, 2, 3, 4, 5]


def test_pre_ppt_short_page_count_is_advisory_not_fixed_blocker(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"] = pack["slides"][:3]
    for idx, slide in enumerate(pack["slides"], start=1):
        slide["slide_no"] = idx
        slide["banker_page_id"] = f"BP-{idx:03d}"
    _write_json(tmp_path / "banker_page_pack.json", pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    compile_result = _compile_page_pack(tmp_path)
    assert compile_result.returncode == 0, compile_result.stdout + compile_result.stderr

    from validate_artifact import validate_artifact

    errors, warnings = validate_artifact("pre_ppt", tmp_path)
    output = "\n".join(errors + warnings)
    assert errors == []
    assert "only 3 renderable page(s)" not in output
    assert "fixed blocker" not in output


def test_final_delivery_allows_direct_ppt_composition_without_renderer_artifacts(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"] = pack["slides"][:2]
    for idx, slide in enumerate(pack["slides"], start=1):
        slide["slide_no"] = idx
        slide["banker_page_id"] = f"BP-{idx:03d}"
    _write_json(tmp_path / "banker_page_pack.json", pack)
    with zipfile.ZipFile(tmp_path / "industry_section_filled_clean.pptx", "w") as archive:
        archive.writestr("[Content_Types].xml", "<Types></Types>")

    from validate_artifact import validate_artifact

    errors, warnings = validate_artifact("final_delivery", tmp_path)
    output = "\n".join(errors + warnings)

    assert errors == []
    assert "direct PPT composition path" in output
    assert "structured-render helper artifacts are absent" in output
    assert "renderer_spec.json" not in output
    assert not (tmp_path / "renderer_spec.json").exists()


def test_pre_ppt_allows_direct_ppt_composition_without_renderer_artifacts(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"] = pack["slides"][:2]
    for idx, slide in enumerate(pack["slides"], start=1):
        slide["slide_no"] = idx
        slide["banker_page_id"] = f"BP-{idx:03d}"
    _write_json(tmp_path / "banker_page_pack.json", pack)

    from validate_artifact import validate_artifact

    errors, warnings = validate_artifact("pre_ppt", tmp_path)
    output = "\n".join(errors + warnings)

    assert errors == []
    assert "direct PPT composition path" in output
    assert "structured-render helper artifacts are absent" in output
    assert "renderer_spec.json" not in output
    assert not (tmp_path / "renderer_spec.json").exists()


def test_final_delivery_accepts_direct_composition_ppt_named_by_marker(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"] = pack["slides"][:2]
    for idx, slide in enumerate(pack["slides"], start=1):
        slide["slide_no"] = idx
        slide["banker_page_id"] = f"BP-{idx:03d}"
    _write_json(tmp_path / "banker_page_pack.json", pack)
    with zipfile.ZipFile(tmp_path / "client_style_direct_composition.pptx", "w") as archive:
        archive.writestr("[Content_Types].xml", "<Types></Types>")
    (tmp_path / "LATEST_FINAL_PPT.txt").write_text("client_style_direct_composition.pptx\n", encoding="utf-8")

    from validate_artifact import validate_artifact

    errors, warnings = validate_artifact("final_delivery", tmp_path)
    output = "\n".join(errors + warnings)

    assert errors == []
    assert "direct PPT composition path" in output
    assert "missing filled PPT output" not in output


def test_style_guided_allows_llm_selected_longer_page_count(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    base_slide = deepcopy(pack["slides"][-1])
    for slide_no in range(9, 14):
        slide = deepcopy(base_slide)
        slide["slide_no"] = slide_no
        slide["banker_page_id"] = f"BP-{slide_no:03d}"
        slide.pop("fixed_page_role", None)
        slide.pop("selected_page_type", None)
        slide["headline"] = f"Additional industry page {slide_no}"
        slide["main_message"] = f"Additional page {slide_no} remains an LLM-selected section page."
        slide["page_argument"] = f"Additional page {slide_no} adds a distinct industry point selected by the LLM."
        slide["body_blocks"] = [
            {
                "copy": f"Additional page {slide_no} point {idx} carries source-linked industry interpretation.",
                "evidence_ids": ["EV-001"],
                "metric_ids": [],
            }
            for idx in range(1, 4)
        ]
        slide["body_copy"] = {f"point_{idx}": block["copy"] for idx, block in enumerate(slide["body_blocks"], start=1)}
        slide["evidence_ids"] = ["EV-001"]
        slide["metric_ids"] = []
        slide["visible_metric_claims"] = []
        slide["chart_data"] = {}
        slide["compare_table_data"] = {}
        slide["source_note"] = "Source: fixture evidence database"
        pack["slides"].append(slide)
    _write_json(tmp_path / "banker_page_pack.json", pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path)])
    assert validation.returncode == 0, validation.stdout + validation.stderr
    assert "style-guided rendering allows this" not in validation.stdout

    compile_result = _compile_page_pack(tmp_path)
    assert compile_result.returncode == 0, compile_result.stdout + compile_result.stderr
    renderer = json.loads((tmp_path / "renderer_spec.json").read_text(encoding="utf-8"))
    assert [slide["slide_no"] for slide in renderer["slides"]] == list(range(1, 14))

    from validate_artifact import validate_artifact

    errors, warnings = validate_artifact("pre_ppt", tmp_path)
    output = "\n".join(errors + warnings)
    assert errors == []
    assert "style-guided rendering does not block on page count alone" not in output


def test_evidence_limited_page_pack_routes_to_bounded_research_before_default_render(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"] = pack["slides"][:3]
    pack["deliverable_readiness"] = {
        "business_action": "targeted_research",
        "targeted_research_rationale": "Fixture intentionally lacks enough evidence-backed pages for PPT render.",
    }
    _write_json(tmp_path / "banker_page_pack.json", pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    compile_result = _compile_page_pack(tmp_path)
    assert compile_result.returncode == 0, compile_result.stdout + compile_result.stderr

    from validate_artifact import validate_artifact

    errors, warnings = validate_artifact("pre_ppt", tmp_path)
    output = "\n".join(errors + warnings)
    assert errors == []
    assert "LLM readiness prompt" in output
    assert "asks for bounded targeted research" in output
    assert "does not clearly mark the page pack client-ready" not in output

    render = subprocess.run(
        [sys.executable, str(SKILL_DIR / "scripts" / "pipeline.py"), "render", "--run-dir", str(tmp_path)],
        text=True,
        capture_output=True,
        cwd=str(SKILL_DIR),
        env={**__import__("os").environ, "PYTHONPATH": ":".join(str(path) for path in SCRIPT_IMPORT_PATHS)},
    )
    assert render.returncode != 0
    assert "next business action is targeted research before formal render" in render.stderr
    assert "max 2 targeted cycle" in render.stderr
    assert "max 5 active request" in render.stderr
    assert "max 3 actual search" in render.stderr
    assert not (tmp_path / "industry_section_filled_clean.pptx").exists()
    assert not (tmp_path / "RESEARCH_LIMITED_REVIEW_industry_section_filled_clean.pptx").exists()

    review_render = subprocess.run(
        [
            sys.executable,
            str(SKILL_DIR / "scripts" / "pipeline.py"),
            "render",
            "--run-dir",
            str(tmp_path),
            "--allow-research-limited-review-render",
        ],
        text=True,
        capture_output=True,
        cwd=str(SKILL_DIR),
        env={**__import__("os").environ, "PYTHONPATH": ":".join(str(path) for path in SCRIPT_IMPORT_PATHS)},
    )
    assert review_render.returncode == 0, review_render.stdout + review_render.stderr
    assert not (tmp_path / "industry_section_filled_clean.pptx").exists()
    assert (tmp_path / "RESEARCH_LIMITED_REVIEW_industry_section_filled_clean.pptx").exists()
    marker = (tmp_path / "RESEARCH_LIMITED_REVIEW_OUTPUT.txt").read_text(encoding="utf-8")
    assert "research-limited review copy" in marker
    assert "Default owner action is targeted research" in marker
    assert "page pack asks for bounded targeted research" in marker


def test_default_render_checks_llm_readiness_before_compiling_internal_artifacts(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["deliverable_readiness"] = {
        "business_action": "targeted_research",
        "targeted_research_rationale": "Needs one opened source that could change chart readiness.",
    }
    _write_json(tmp_path / "banker_page_pack.json", pack)

    render = subprocess.run(
        [sys.executable, str(SKILL_DIR / "scripts" / "pipeline.py"), "render", "--run-dir", str(tmp_path)],
        text=True,
        capture_output=True,
        cwd=str(SKILL_DIR),
        env={**__import__("os").environ, "PYTHONPATH": ":".join(str(path) for path in SCRIPT_IMPORT_PATHS)},
    )

    assert render.returncode != 0
    assert "next business action is targeted research before formal render" in render.stderr
    assert not (tmp_path / "template_registry.json").exists()
    assert not (tmp_path / "artifacts/template_selection.json").exists()
    assert not (tmp_path / "artifacts/runtime_dependencies.json").exists()
    assert not (tmp_path / "deck_blueprint.json").exists()
    assert not (tmp_path / "page_evidence_contract.json").exists()
    assert not (tmp_path / "renderer_spec.json").exists()


def test_render_without_page_pack_stops_before_mechanical_side_effects(tmp_path: Path) -> None:
    render = subprocess.run(
        [sys.executable, str(SKILL_DIR / "scripts" / "pipeline.py"), "render", "--run-dir", str(tmp_path)],
        text=True,
        capture_output=True,
        cwd=str(SKILL_DIR),
        env={**__import__("os").environ, "PYTHONPATH": ":".join(str(path) for path in SCRIPT_IMPORT_PATHS)},
    )

    assert render.returncode != 0
    assert "banker_page_pack.json is missing, so render cannot start" in render.stderr
    assert "Generation LLM should author the banker page pack first" in render.stderr
    assert not (tmp_path / "artifacts/template_selection.json").exists()
    assert not (tmp_path / "artifacts/runtime_dependencies.json").exists()
    assert not (tmp_path / "template_registry.json").exists()
    assert not (tmp_path / "deck_blueprint.json").exists()
    assert not (tmp_path / "renderer_spec.json").exists()


def test_strict_runtime_readiness_does_not_obscure_page_pack_readiness(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["deliverable_readiness"] = {
        "business_action": "targeted_research",
        "targeted_research_rationale": "Needs one opened source that could change chart readiness.",
    }
    _write_json(tmp_path / "banker_page_pack.json", pack)

    render = subprocess.run(
        [
            sys.executable,
            str(SKILL_DIR / "scripts" / "pipeline.py"),
            "render",
            "--run-dir",
            str(tmp_path),
            "--strict-runtime-readiness",
        ],
        text=True,
        capture_output=True,
        cwd=str(SKILL_DIR),
        env={**__import__("os").environ, "PYTHONPATH": ":".join(str(path) for path in SCRIPT_IMPORT_PATHS)},
    )

    assert render.returncode != 0
    assert "next business action is targeted research before formal render" in render.stderr
    assert "runtime readiness diagnostics" not in render.stderr
    assert not (tmp_path / "artifacts/runtime_dependencies.json").exists()
    assert not (tmp_path / "artifacts/template_selection.json").exists()


def test_review_render_flag_does_not_bypass_page_pack_repair_route(tmp_path: Path) -> None:
    _write_json(
        tmp_path / "banker_page_pack.json",
        {
            "schema_version": "banker_page_pack",
            "deliverable_readiness": {
                "business_action": "repair_page_pack",
                "readiness_note": (
                    "Not client-ready: visible wording still sounds like internal workpaper language "
                    "and pages need denser client-facing exhibits."
                )
            },
            "slides": [
                {
                    "page_argument": "The section needs client-facing page repair.",
                    "headline": "Client-facing market view",
                    "body_blocks": [{"copy": "The page should be rewritten before render."}],
                }
            ],
        },
    )

    render = subprocess.run(
        [
            sys.executable,
            str(SKILL_DIR / "scripts" / "pipeline.py"),
            "render",
            "--run-dir",
            str(tmp_path),
            "--allow-research-limited-review-render",
        ],
        text=True,
        capture_output=True,
        cwd=str(SKILL_DIR),
        env={**__import__("os").environ, "PYTHONPATH": ":".join(str(path) for path in SCRIPT_IMPORT_PATHS)},
    )

    assert render.returncode != 0
    assert "Repair banker_page_pack first" in render.stderr
    assert "Route to targeted research only after naming the specific evidence question" in render.stderr
    assert not (tmp_path / "template_registry.json").exists()
    assert not (tmp_path / "renderer_spec.json").exists()
    assert not (tmp_path / "RESEARCH_LIMITED_REVIEW_OUTPUT.txt").exists()


def test_source_limit_readiness_routes_to_qc_decision_not_more_research(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["deliverable_readiness"] = {
        "business_action": "qc_user_decision",
        "source_unavailable": True,
        "readiness_note": (
            "Evidence-limited after bounded targeted research loop exhausted: "
            "public sources unavailable for audit-grade market size."
        )
    }
    _write_json(tmp_path / "banker_page_pack.json", pack)

    render = subprocess.run(
        [sys.executable, str(SKILL_DIR / "scripts" / "pipeline.py"), "render", "--run-dir", str(tmp_path)],
        text=True,
        capture_output=True,
        cwd=str(SKILL_DIR),
        env={**__import__("os").environ, "PYTHONPATH": ":".join(str(path) for path in SCRIPT_IMPORT_PATHS)},
    )

    assert render.returncode != 0
    assert "bounded targeted research loop is exhausted" in render.stderr
    assert "Do not start another search loop by default" in render.stderr
    assert "create only a non-final research-limited review copy" in render.stderr
    assert "accept an evidence-limited review state" not in render.stderr
    assert "next business action is targeted research before formal render" not in render.stderr
    assert not (tmp_path / "template_registry.json").exists()
    assert not (tmp_path / "renderer_spec.json").exists()

    review_render = subprocess.run(
        [
            sys.executable,
            str(SKILL_DIR / "scripts" / "pipeline.py"),
            "render",
            "--run-dir",
            str(tmp_path),
            "--allow-research-limited-review-render",
        ],
        text=True,
        capture_output=True,
        cwd=str(SKILL_DIR),
        env={**__import__("os").environ, "PYTHONPATH": ":".join(str(path) for path in SCRIPT_IMPORT_PATHS)},
    )

    assert review_render.returncode == 0, review_render.stdout + review_render.stderr
    marker = (tmp_path / "RESEARCH_LIMITED_REVIEW_OUTPUT.txt").read_text(encoding="utf-8")
    assert "research-limited review copy" in marker
    assert "Default owner action is QC/user decision" in marker
    assert "create only a non-final research-limited review copy" in marker
    assert "accept an evidence-limited review state" not in marker
    assert "Default next step is targeted research" not in marker


def test_compare_table_body_blocks_do_not_expose_template_fields_in_style_guided_mode(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    slide_6 = pack["slides"][5]
    slide_6["body_blocks"][0]["target_field"] = "table_row_1"
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert result.returncode == 0, result.stdout + result.stderr
    assert "takes table content from compare_table_data" not in result.stdout
    assert "active body fields" not in result.stdout
    review_payload = json.loads(result.stdout)
    assert "template_diagnostics" not in review_payload
    assert review_payload["template_diagnostics_path"].endswith("banker_page_pack_template_diagnostics.json")
    diagnostics = json.loads((tmp_path / "artifacts/banker_page_pack_template_diagnostics.json").read_text(encoding="utf-8"))
    slide_6_diagnostics = next(item for item in diagnostics["slides"] if item["slide_no"] == 6)
    assert slide_6_diagnostics["template_fit_mode"] == "style_guided"
    assert "template_field_contract_mode" not in slide_6_diagnostics
    assert "active_template_body_fields" not in slide_6_diagnostics
    assert "inactive_template_fields_when_compare_table_data_present" not in slide_6_diagnostics
    assert "template_body_field_hints" not in slide_6_diagnostics
    assert "body_field_unit_limits" not in slide_6_diagnostics
    assert "default_body_field_unit_limit" not in slide_6_diagnostics
    assert "headline_main_message_line_rules" not in slide_6_diagnostics
    assert slide_6_diagnostics["strict_layout_placeholders"] == []
    assert slide_6_diagnostics["strict_layout_active_placeholders"] == []
    assert slide_6_diagnostics["strict_layout_capacity_hints"] == {}
    assert "style cues only" in slide_6_diagnostics["style_guided_template_note"]

    result = _compile_page_pack(tmp_path)
    assert result.returncode == 0, result.stdout + result.stderr
    renderer = json.loads((tmp_path / "renderer_spec.json").read_text(encoding="utf-8"))
    slide_6 = next(slide for slide in renderer["slides"] if slide["slide_no"] == 6)
    assert "table_row_1" not in slide_6["body_copy"]


def test_compare_table_body_blocks_fail_in_strict_layout_mode(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    slide_6 = pack["slides"][5]
    slide_6["body_blocks"][0]["target_field"] = "table_row_1"
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    _write_json(tmp_path / "artifacts/rendering_policy.json", {"template_contract_mode": "strict_layout"})
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert result.returncode != 0
    assert "takes table content from compare_table_data" in result.stdout


def test_compare_table_columns_alias_compiles_with_canonical_warning(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][5]["compare_table_data"] = {
        "columns": ["Dimension", "Evidence read", "Pitch use"],
        "rows": [
            ["Demand", "Evidence-backed demand signal", "Support market framing"],
            ["Competition", "Peer variation visible", "Frame positioning"],
        ],
    }
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])
    assert validation.returncode == 0, validation.stdout + validation.stderr
    assert "uses columns; structured-render helper accepts it" not in validation.stdout
    assert "row 1 is a list; structured-render helper accepts it" not in validation.stdout

    result = _compile_page_pack(tmp_path)
    assert result.returncode == 0, result.stdout + result.stderr
    renderer = json.loads((tmp_path / "renderer_spec.json").read_text(encoding="utf-8"))
    slide_6 = next(slide for slide in renderer["slides"] if slide["slide_no"] == 6)
    assert slide_6["compare_table_data"]["headers"] == ["Dimension", "Evidence read", "Pitch use"]
    assert slide_6["compare_table_data"]["rows"][0]["label"] == "Demand"


def test_compare_table_ppt_like_rows_compile_without_field_memory(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][5]["compare_table_data"] = {
        "table_header": "Dimension｜Evidence read｜Pitch use",
        "table_row_1": "Demand｜Evidence-backed demand signal｜Support market framing",
        "table_row_2": "Competition｜Peer variation visible｜Frame positioning",
    }
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])
    assert validation.returncode == 0, validation.stdout + validation.stderr
    assert "uses table_header; structured-render helper accepts it" not in validation.stdout
    assert "uses table_row_* fields; structured-render helper accepts them" not in validation.stdout

    result = _compile_page_pack(tmp_path)
    assert result.returncode == 0, result.stdout + result.stderr
    renderer = json.loads((tmp_path / "renderer_spec.json").read_text(encoding="utf-8"))
    slide_6 = next(slide for slide in renderer["slides"] if slide["slide_no"] == 6)
    assert slide_6["compare_table_data"]["headers"] == ["Dimension", "Evidence read", "Pitch use"]
    assert slide_6["compare_table_data"]["rows"][1]["label"] == "Competition"


def test_compare_table_scalar_cells_are_normalized_in_style_guided_mode(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][5]["compare_table_data"] = {
        "headers": ["Dimension", "Evidence read", "Pitch use"],
        "rows": [
            {"label": "Demand", "cells": "Evidence-backed demand signal｜Support market framing"},
            {"label": "Competition", "cells": "Peer variation visible|Frame positioning"},
        ],
    }
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    validation = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])
    assert validation.returncode == 0, validation.stdout + validation.stderr
    assert "row 1.cells is scalar" in validation.stdout
    assert "style-guided render will normalize it into cells" in validation.stdout

    result = _compile_page_pack(tmp_path)
    assert result.returncode == 0, result.stdout + result.stderr
    renderer = json.loads((tmp_path / "renderer_spec.json").read_text(encoding="utf-8"))
    slide_6 = next(slide for slide in renderer["slides"] if slide["slide_no"] == 6)
    assert slide_6["compare_table_data"]["rows"][0]["cells"] == [
        "Evidence-backed demand signal",
        "Support market framing",
    ]
    assert slide_6["compare_table_data"]["rows"][1]["cells"] == [
        "Peer variation visible",
        "Frame positioning",
    ]


def test_compare_table_scalar_cells_fail_in_strict_layout(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][5]["compare_table_data"] = {
        "headers": ["Dimension", "Evidence read", "Pitch use"],
        "rows": [
            {"label": "Demand", "cells": "Evidence-backed demand signal｜Support market framing"},
            {"label": "Competition", "cells": ["Peer variation visible", "Frame positioning"]},
        ],
    }
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    _write_json(tmp_path / "artifacts/rendering_policy.json", {"template_contract_mode": "strict_layout"})
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert result.returncode != 0
    assert "row 1.cells must be a list" in result.stdout


def test_compare_table_incomplete_payload_is_advisory_in_style_guided_mode(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][5]["compare_table_data"] = {
        "rows": [
            {"label": "Demand", "cells": ["Channel signal remains useful"]},
            {"label": "Competition", "cells": ["Peer read remains useful"]},
        ],
    }
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert result.returncode == 0, result.stdout + result.stderr
    assert "compare_table_data requires non-empty headers list" in result.stdout
    assert "style-guided render treats this table payload as advisory" in result.stdout


def test_compare_table_incomplete_payload_fails_in_strict_layout(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][5]["compare_table_data"] = {
        "rows": [
            {"label": "Demand", "cells": ["Channel signal remains useful"]},
            {"label": "Competition", "cells": ["Peer read remains useful"]},
        ],
    }
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    _write_json(tmp_path / "artifacts/rendering_policy.json", {"template_contract_mode": "strict_layout"})
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert result.returncode != 0
    assert "compare_table_data requires non-empty headers list" in result.stdout


def test_compare_table_width_mismatch_is_advisory_in_style_guided_mode(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][5]["compare_table_data"] = {
        "headers": ["Dimension", "Evidence read", "Pitch use"],
        "rows": [
            {"label": "Demand", "cells": ["Evidence-backed demand signal", "Support market framing", "Extra repeated column"]},
            {"label": "Competition", "cells": ["Peer variation visible", "Frame positioning"]},
            {"label": "Economics", "cells": ["Profit-pool evidence", "Frame margin quality"]},
        ],
    }
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert result.returncode == 0, result.stdout + result.stderr
    assert "3 headers require 2 cells after label" in result.stdout
    assert "style-guided render will merge extra cells or pad missing cells" in result.stdout

    compile_result = _compile_page_pack(tmp_path)
    assert compile_result.returncode == 0, compile_result.stdout + compile_result.stderr


def test_compare_table_width_mismatch_fails_in_strict_layout(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    pack["slides"][5]["compare_table_data"] = {
        "headers": ["Dimension", "Evidence read", "Pitch use"],
        "rows": [
            {"label": "Demand", "cells": ["Evidence-backed demand signal", "Support market framing", "Extra repeated column"]},
            {"label": "Competition", "cells": ["Peer variation visible", "Frame positioning"]},
        ],
    }
    pack_path = tmp_path / "banker_page_pack.json"
    _write_json(pack_path, pack)
    _write_json(tmp_path / "artifacts/rendering_policy.json", {"template_contract_mode": "strict_layout"})
    (tmp_path / "template_registry.json").write_text(template_registry_path.read_text(encoding="utf-8"), encoding="utf-8")

    result = _run("pipeline.py", ["review", "--artifact", "banker_page_pack", "--run-dir", str(tmp_path), "--path", str(pack_path)])

    assert result.returncode != 0
    assert "3 headers require 2 cells after label" in result.stdout


def test_render_auto_refreshes_template_and_compiled_artifacts(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    _write_json(tmp_path / "banker_page_pack.json", pack)

    result = subprocess.run(
        [sys.executable, str(SKILL_DIR / "scripts" / "pipeline.py"), "render", "--run-dir", str(tmp_path)],
        text=True,
        capture_output=True,
        cwd=str(SKILL_DIR),
        env={**__import__("os").environ, "PYTHONPATH": ":".join(str(path) for path in SCRIPT_IMPORT_PATHS)},
    )

    assert result.returncode == 0, result.stdout + result.stderr
    assert (tmp_path / "template_registry.json").exists()
    assert (tmp_path / "deck_blueprint.json").exists()
    assert (tmp_path / "page_evidence_contract.json").exists()
    assert (tmp_path / "renderer_spec.json").exists()
    assert (tmp_path / "industry_section_filled_clean.pptx").exists()
    assert (tmp_path / "artifacts/banker_page_pack_template_diagnostics.json").exists()


def test_style_guided_render_accepts_simple_non_token_template(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    pptx = pytest.importorskip("pptx")
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(914400, 914400, 7315200, 700000)
    title.text = "Simple style reference"
    note = slide.shapes.add_textbox(914400, 5800000, 7315200, 500000)
    note.text = "Source footer style"
    simple_template = tmp_path / "simple_style_template.pptx"
    prs.save(simple_template)

    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    for slide_payload in pack["slides"]:
        slide_payload["selected_page_type"] = "freeform_page"
        slide_payload["body_copy"] = {
            "mechanism": f"Page {slide_payload['slide_no']} mechanism remains authored by the LLM.",
            "data_read": "Visible data and evidence shape the page composition.",
            "banker_view": "Python should place this reliably without requiring template fields.",
        }
    pack["slides"][5]["compare_table_data"] = {
        "headers": ["Dimension", "Read"],
        "rows": [
            {"label": "Demand", "cells": ["Growth mechanism is visible"]},
            {"label": "Competition", "cells": ["Peer differences stay evidence-bound"]},
        ],
    }
    _write_json(tmp_path / "banker_page_pack.json", pack)

    result = subprocess.run(
        [
            sys.executable,
            str(SKILL_DIR / "scripts" / "pipeline.py"),
            "render",
            "--run-dir",
            str(tmp_path),
            "--template",
            str(simple_template),
        ],
        text=True,
        capture_output=True,
        cwd=str(SKILL_DIR),
        env={**__import__("os").environ, "PYTHONPATH": ":".join(str(path) for path in SCRIPT_IMPORT_PATHS)},
    )

    assert result.returncode == 0, result.stdout + result.stderr
    assert (tmp_path / "industry_section_filled_clean.pptx").exists()
    assert not (tmp_path / "replacement_dict.json").exists()
    assert not (tmp_path / "artifacts/template_token_check.json").exists()
    postprocess = json.loads((tmp_path / "artifacts/postprocess_ppt_visuals.log.json").read_text(encoding="utf-8"))
    assert postprocess["style_guided_render"] is True


def test_style_guided_render_copies_low_content_template_base_slide(
    tmp_path: Path,
    deck_blueprint_data: dict,
    template_registry_path: Path,
) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.dml.color import RGBColor

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, 320000)
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(0xE8, 0xEF, 0xF7)
    band.line.fill.background()
    simple_template = tmp_path / "blank_style_template.pptx"
    prs.save(simple_template)

    template_registry = json.loads(template_registry_path.read_text(encoding="utf-8"))
    pack = _banker_page_pack(deck_blueprint_data, template_registry)
    for slide_payload in pack["slides"]:
        slide_payload["selected_page_type"] = "freeform_page"
        slide_payload["body_copy"] = {
            "mechanism": f"Page {slide_payload['slide_no']} mechanism remains authored by the LLM.",
            "data_read": "Visible data and evidence shape the page composition.",
            "banker_view": "Python should place this reliably without requiring template fields.",
        }
    _write_json(tmp_path / "banker_page_pack.json", pack)

    result = subprocess.run(
        [
            sys.executable,
            str(SKILL_DIR / "scripts" / "pipeline.py"),
            "render",
            "--run-dir",
            str(tmp_path),
            "--template",
            str(simple_template),
        ],
        text=True,
        capture_output=True,
        cwd=str(SKILL_DIR),
        env={**__import__("os").environ, "PYTHONPATH": ":".join(str(path) for path in SCRIPT_IMPORT_PATHS)},
    )

    assert result.returncode == 0, result.stdout + result.stderr
    postprocess = json.loads((tmp_path / "artifacts/postprocess_ppt_visuals.log.json").read_text(encoding="utf-8"))
    assert postprocess["style_base_slide"] == 1
    assert postprocess["style_base_strategy"] == "copied_low_content_template_slide"
